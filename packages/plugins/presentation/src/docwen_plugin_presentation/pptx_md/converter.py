"""PPTX → Markdown converter.

Parses PowerPoint (PPTX) files using ``python-pptx`` and produces a
Markdown document with one section per slide.

The converter:
- Only writes to staging via ``WorkspaceHandle``.
- Checks cancellation before expensive operations.
- Reports progress through ``ProgressSink``.
- Returns a ``ConversionResult`` with ``ArtifactManifest`` entries.

Note: PPT (legacy binary format) → MD is handled by
``PptToMarkdownConverter`` in ``ppt_converter.py``. It uses the shared
Office bridge to preprocess PPT into PPTX, then delegates back to this
converter for Markdown extraction.
"""

from __future__ import annotations

import hashlib
import os
import re
import time
import uuid
from pathlib import PurePosixPath
from typing import TYPE_CHECKING, Any

from docwen_core.text.ocr import format_ocr_best_effort_warning
from docwen_plugin_presentation.pptx_md.request_policy import (
    PresentationMarkdownRequestPolicy,
    build_presentation_markdown_request_policy,
)

if TYPE_CHECKING:
    from docwen_core.protocols.execution_context import ConverterContext


def _report_ocr_best_effort(progress: Any, status: object, *, location: str) -> None:
    """Report one safe, request-visible warning for a fallible OCR outcome."""
    message = format_ocr_best_effort_warning(status)
    if message is None:
        return
    progress.report_diagnostic(
        "warning",
        message,
        code="OCR-BEST-EFFORT",
        location=location,
    )


class PptxToMarkdownConverter:
    """Convert a PPTX file to Markdown.

    Uses ``python-pptx`` to parse slides and extracts text content,
    tables, and images from each slide.
    """

    @staticmethod
    def _record_payload_warning(
        context: Any,
        payload_stats: dict[str, Any] | None,
        *,
        code: str,
        message: str,
        location: str,
        error: BaseException | str,
    ) -> None:
        """Record a recoverable payload loss without failing the conversion."""
        detail = str(error).strip() or type(error).__name__
        logger = getattr(context, "logger", None)
        if logger is not None:
            logger.warning(f"{message} Location: {location}. Detail: {detail}")
        if payload_stats is None:
            return
        payload_stats.setdefault("payload_warning_diagnostics", []).append(
            {
                "code": code,
                "message": message,
                "location": location,
            }
        )

    def convert(
        self,
        context: ConverterContext,
        *,
        source_path_for_naming: str | None = None,
        request_policy: PresentationMarkdownRequestPolicy | None = None,
    ) -> Any:
        """Run the PPTX → Markdown conversion.

        Args:
            context: The plugin execution context.

        Returns:
            ``ConversionResult`` with staging artifacts.
        """
        from docwen_core.models.result import (
            ConversionDiagnostic,
            ConversionMetrics,
            ConversionResult,
        )

        t_start = time.monotonic()
        task_id = context.request.request_id
        input_path = context.workspace.input_path
        naming_path = source_path_for_naming or input_path

        # 1. Check cancellation
        context.cancellation.check()

        # 2. Report start
        context.progress.report_progress(0.0, "Starting PPTX → Markdown conversion")
        context.logger.info(f"PPTX→MD: reading {input_path}")

        # 3. Parse and convert
        context.cancellation.check()
        context.progress.report_progress(10.0, "Parsing PPTX slides...")

        try:
            md_content, stats = self._parse_pptx(
                input_path,
                context,
                source_path_for_naming=naming_path,
                request_policy=request_policy,
            )
        except Exception as exc:
            context.logger.error(f"PPTX→MD conversion failed: {exc}")
            from docwen_core.models.result import ConversionErrorInfo

            return ConversionResult(
                task_id=task_id,
                success=False,
                error=ConversionErrorInfo(
                    error_type="conversion_failed",
                    message=str(exc),
                    diagnostic_code="PPTX2MD-PARSE-ERROR",
                ),
                diagnostics=[
                    ConversionDiagnostic(
                        level="error",
                        message=f"Failed to parse PPTX: {exc}",
                        code="PPTX2MD-PARSE-ERROR",
                    ),
                ],
            )

        # 4. Write markdown to staging
        context.cancellation.check()
        context.progress.report_progress(80.0, "Writing Markdown to staging...")

        output_path = context.workspace.create_artifact_path("primary", ".md")
        try:
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(md_content)
        except OSError as exc:
            context.logger.error(f"PPTX→MD write failed: {exc}")
            from docwen_core.models.result import ConversionErrorInfo

            return ConversionResult(
                task_id=task_id,
                success=False,
                error=ConversionErrorInfo(
                    error_type="conversion_failed",
                    message=f"Failed to write output file: {exc}",
                    diagnostic_code="PPTX2MD-WRITE-ERROR",
                ),
                diagnostics=[
                    ConversionDiagnostic(
                        level="error",
                        message=f"File write error at {output_path}: {exc}",
                        code="PPTX2MD-WRITE-ERROR",
                    ),
                ],
            )

        # 5. Build artifact manifest
        input_basename = os.path.basename(naming_path)
        suggested_name = input_basename.rsplit(".", 1)[0] + ".md"

        from docwen_core.models.artifact import ArtifactManifest

        artifact = ArtifactManifest(
            artifact_id=str(uuid.uuid4()),
            kind="primary",
            staging_path=output_path,
            suggested_name=suggested_name,
            media_type="text/markdown",
            metadata={
                "slide_count": stats.get("slides", 0),
                "hidden_slide_count": stats.get("hidden_slides", 0),
                "table_count": stats.get("tables", 0),
                "image_count": stats.get("images", 0),
                "smartart_text_count": stats.get("smartart_texts", 0),
                **({"chart_count": stats["charts"]} if stats.get("charts") else {}),
                **({"audio_count": stats["audio"]} if stats.get("audio") else {}),
                **({"video_count": stats["video"]} if stats.get("video") else {}),
                **(
                    {"payload_warning_count": len(stats["payload_warning_diagnostics"])}
                    if stats.get("payload_warning_diagnostics")
                    else {}
                ),
                "title": stats.get("title", ""),
            },
            is_primary=True,
        )
        context.workspace.add_artifact(artifact)
        registered_artifacts = [
            registered
            for registered in getattr(context.workspace, "registered_artifacts", [])
            if registered.artifact_id != artifact.artifact_id
        ]

        # 6. Report completion
        context.progress.report_artifact_ready(artifact.artifact_id, suggested_name)
        context.progress.report_progress(100.0, "Conversion complete")
        context.logger.info(
            f"PPTX→MD complete: {stats['slides']} slides, {stats['hidden_slides']} hidden, "
            f"{stats['tables']} tables, {stats['images']} images, {stats['smartart_texts']} SmartArt texts"
        )

        duration_ms = (time.monotonic() - t_start) * 1000.0

        diagnostics = [
            ConversionDiagnostic(
                level="info",
                message=(
                    f"Converted PPTX to Markdown: {stats['slides']} slides, "
                    f"{stats['tables']} tables, {stats['images']} images, "
                    f"{stats['charts']} charts, {stats['audio']} audio, {stats['video']} video"
                ),
                code="PPTX2MD-OK",
            )
        ]
        diagnostics.extend(
            ConversionDiagnostic(
                level="warning",
                message=(
                    "Chart semantics were preserved, but no deterministic supported "
                    "chart snapshot renderer is available."
                ),
                code="PPTX-CHART-SNAPSHOT-UNAVAILABLE",
                location=location,
            )
            for location in stats["chart_snapshot_unavailable_locations"]
        )
        diagnostics.extend(
            ConversionDiagnostic(
                level="warning",
                message="Speaker notes were requested but could not be read.",
                code="PPTX-NOTES-UNAVAILABLE",
                location=location,
            )
            for location in stats["notes_unavailable_locations"]
        )
        diagnostics.extend(
            ConversionDiagnostic(
                level="warning",
                message=str(warning["message"]),
                code=str(warning["code"]),
                location=str(warning["location"]),
            )
            for warning in stats.get("payload_warning_diagnostics", [])
        )

        return ConversionResult(
            task_id=task_id,
            success=True,
            artifacts=[artifact, *registered_artifacts],
            diagnostics=diagnostics,
            error=None,
            metrics=ConversionMetrics(
                duration_ms=duration_ms,
                input_bytes=os.path.getsize(naming_path) if os.path.isfile(naming_path) else 0,
                output_bytes=len(md_content.encode("utf-8")),
                extra=stats,
            ),
        )

    # ── Internal parsing ──────────────────────────────────────────────

    @staticmethod
    def _extract_slide_sections(pres: Any) -> dict[int, str]:
        """Extract section names for each slide ID from the PPTX XML.

        Returns a dict mapping slide_id -> section_name.
        Falls back to empty dict if no sections are found.
        """
        section_map: dict[int, str] = {}
        try:
            sections = pres.element.xpath(".//*[local-name()='sectionLst']/*[local-name()='section']")
        except Exception:
            sections = []
        for section in sections:
            name = str(section.get("name") or "").strip()
            if not name:
                continue
            try:
                slide_id_elements = section.xpath(".//*[local-name()='sldId']")
            except Exception:
                slide_id_elements = []
            for slide_id_element in slide_id_elements:
                try:
                    section_map[int(str(slide_id_element.get("id")))] = name
                except Exception:
                    continue
        return section_map

    def _parse_pptx(
        self,
        input_path: str,
        context: ConverterContext,
        *,
        source_path_for_naming: str | None = None,
        request_policy: PresentationMarkdownRequestPolicy | None = None,
    ) -> tuple[str, dict[str, Any]]:
        """Parse a PPTX file and return (markdown_text, stats_dict).

        Processes slides in order, extracting:
        - Slide titles (from title/centertitle placeholders)
        - Text frames (all text content)
        - Tables (converted to Markdown tables)
        - Images (extracted to staging)
        """
        from pptx import Presentation

        pres = Presentation(input_path)

        # Extract options from request context
        opts: dict[str, Any] = context.request.options or {}
        request_policy = request_policy or build_presentation_markdown_request_policy(context, opts)

        # Get presentation title from core properties
        title = ""
        try:
            title = str(getattr(pres.core_properties, "title", "") or "").strip()
        except Exception:
            title = ""
        if not title:
            title = os.path.splitext(os.path.basename(source_path_for_naming or input_path))[0]

        # Build section map (L9)
        section_map = self._extract_slide_sections(pres)
        resource_stem = self._resource_stem(source_path_for_naming or input_path)

        sections: list[str] = []
        total_tables = 0
        total_images = 0
        total_smartart_texts = 0
        payload_stats: dict[str, Any] = {
            "charts": 0,
            "audio": 0,
            "video": 0,
            "chart_snapshot_unavailable_locations": [],
            "notes_unavailable_locations": [],
            "payload_warning_diagnostics": [],
        }
        hidden_slide_count = 0
        slide_count = len(list(pres.slides))
        current_section: str | None = None

        for i, slide in enumerate(pres.slides, 1):
            context.cancellation.check()
            if self._is_hidden_slide(slide):
                hidden_slide_count += 1
            progress = 10.0 + 60.0 * (i / max(slide_count, 1))
            context.progress.report_progress(progress, f"Processing slide {i}/{slide_count}")

            # Emit section heading when section changes (L9)
            slide_id = getattr(slide, "slide_id", None)
            if slide_id is not None and slide_id in section_map:
                sec_name = section_map[slide_id]
                if sec_name != current_section:
                    current_section = sec_name
                    sections.append(f"# {sec_name}")
                    sections.append("")

            slide_lines, img_cnt, tbl_cnt, smartart_cnt = self._process_slide(
                slide,
                i,
                context,
                opts,
                resource_stem=resource_stem,
                request_policy=request_policy,
                payload_stats=payload_stats,
            )
            sections.extend(slide_lines)
            total_images += img_cnt
            total_tables += tbl_cnt
            total_smartart_texts += smartart_cnt

        # M10 — YAML frontmatter
        title_key = "title"
        yaml_key_labels = opts.get("yaml_key_labels")
        if isinstance(yaml_key_labels, dict):
            label_title = yaml_key_labels.get("title")
            if isinstance(label_title, str) and label_title.strip():
                title_key = label_title.strip()
        yaml_lines = ["---", "aliases:", f"  - {title}", f"{title_key}: {title}", "---"]
        yaml_block = "\n".join(yaml_lines) + "\n"
        body = yaml_block + f"# {title}\n\n" + "\n\n".join(s for s in sections if s.strip())
        body = body.strip() + "\n"

        stats = {
            "slides": slide_count,
            "hidden_slides": hidden_slide_count,
            "tables": total_tables,
            "images": total_images,
            "smartart_texts": total_smartart_texts,
            **payload_stats,
            "title": title,
        }
        return body, stats

    @staticmethod
    def _is_hidden_slide(slide: Any) -> bool:
        """Return whether the source marks a slide as hidden.

        PPTX→Markdown intentionally includes hidden slides so archival text is
        not silently discarded. The count is exposed in artifact metadata.
        """
        element = getattr(slide, "element", None)
        if element is None:
            return False
        try:
            value = str(element.get("show") or "").strip().lower()
        except Exception:
            return False
        return value in {"0", "false", "off", "no"}

    @staticmethod
    def _extract_smartart_texts(shape: Any, slide: Any) -> tuple[list[str], bool, list[str]]:
        """Return SmartArt text, whether SmartArt was found, and failed relationship IDs."""
        element = getattr(shape, "element", None)
        slide_part = getattr(slide, "part", None)
        relationships = getattr(slide_part, "rels", None)
        if element is None or relationships is None:
            return [], False, []
        try:
            relationship_ids = element.xpath(
                ".//*[local-name()='graphicData' and "
                "@uri='http://schemas.openxmlformats.org/drawingml/2006/diagram']"
                "//*[local-name()='relIds']/@*[local-name()='dm']"
            )
        except Exception:
            return [], False, []

        result: list[str] = []
        failed_relationship_ids: list[str] = []
        seen_relationship_ids: set[str] = set()
        for raw_relationship_id in relationship_ids:
            relationship_id = str(raw_relationship_id)
            if not relationship_id or relationship_id in seen_relationship_ids:
                continue
            seen_relationship_ids.add(relationship_id)
            try:
                relationship = relationships[relationship_id]
                if not str(relationship.reltype).endswith("/diagramData"):
                    failed_relationship_ids.append(relationship_id)
                    continue
                from pptx.oxml import parse_xml

                # python-pptx/lxml xpath stubs expose a scalar-or-list union;
                # diagramData queries below are element-set queries.
                diagram: Any = parse_xml(relationship.target_part.blob)
                points = list(diagram.xpath(".//*[local-name()='pt']"))
            except Exception:
                failed_relationship_ids.append(relationship_id)
                continue

            point_ids: list[str] = []
            point_text: dict[str, str] = {}
            for point in points:
                point_id = str(point.get("modelId") or "")
                if not point_id:
                    continue
                point_ids.append(point_id)
                paragraphs: list[str] = []
                for paragraph in point.xpath(".//*[local-name()='p']"):
                    runs = paragraph.xpath(".//*[local-name()='t']/text()")
                    text = " ".join("".join(str(run) for run in runs).split())
                    if text:
                        paragraphs.append(text)
                point_text[point_id] = " ".join(paragraphs)

            point_position = {point_id: index for index, point_id in enumerate(point_ids)}
            adjacency: dict[str, list[tuple[int, int, str]]] = {}
            sources: set[str] = set()
            destinations: set[str] = set()
            try:
                connections = diagram.xpath(".//*[local-name()='cxn']")
            except Exception:
                connections = []
            for connection in connections:
                if connection.get("type"):
                    continue
                source_id = str(connection.get("srcId") or "")
                destination_id = str(connection.get("destId") or "")
                if not source_id or not destination_id:
                    continue
                try:
                    source_order = int(str(connection.get("srcOrd") or "0"))
                except (TypeError, ValueError):
                    source_order = 0
                sources.add(source_id)
                destinations.add(destination_id)
                adjacency.setdefault(source_id, []).append(
                    (
                        source_order,
                        point_position.get(destination_id, len(point_ids)),
                        destination_id,
                    )
                )
            for children in adjacency.values():
                children.sort()

            roots = sorted(
                sources - destinations,
                key=lambda point_id: point_position.get(point_id, -1),
            )
            visited: set[str] = set()

            def visit(
                point_id: str,
                *,
                _visited: set[str] = visited,
                _point_text: dict[str, str] = point_text,
                _adjacency: dict[str, list[tuple[int, int, str]]] = adjacency,
            ) -> None:
                if point_id in _visited:
                    return
                _visited.add(point_id)
                text = _point_text.get(point_id, "")
                if text:
                    result.append(text)
                for _order, _position, child_id in _adjacency.get(point_id, []):
                    visit(child_id)

            for root_id in roots:
                visit(root_id)
            for point_id in point_ids:
                visit(point_id)
        return result, bool(relationship_ids), failed_relationship_ids

    @staticmethod
    def _resource_stem(source_path: str) -> str:
        """Build a bounded, link-safe source-derived label for related artifacts."""
        from docwen_core.markdown_utils import sanitize_for_wiki_link

        raw_stem = os.path.splitext(os.path.basename(source_path))[0]
        wiki_safe_stem = sanitize_for_wiki_link(raw_stem)
        safe_stem = re.sub(r"[^\w.-]+", "_", wiki_safe_stem).strip("._-")
        safe_stem = safe_stem or "presentation"
        if len(safe_stem.encode("utf-8")) > 64:
            safe_stem = safe_stem.encode("utf-8")[:64].decode("utf-8", errors="ignore")
            safe_stem = safe_stem.rstrip("._-") or "presentation"
        return safe_stem

    def _process_slide(
        self,
        slide: Any,
        slide_index: int,
        context: ConverterContext,
        options: dict[str, Any] | None = None,
        *,
        resource_stem: str = "presentation",
        request_policy: PresentationMarkdownRequestPolicy | None = None,
        payload_stats: dict[str, Any] | None = None,
    ) -> tuple[list[str], int, int, int]:
        """Process a single slide.

        Args:
            slide: python-pptx Slide object.
            slide_index: 1-based slide index.
            context: Plugin execution context.
            options: Conversion options dict.

        Returns (section_lines, image_count, table_count, smartart_text_count).
        """
        opts = options or {}
        export_notes: bool = opts.get("export_notes", False)
        keep_images: bool = opts.get("to_md_keep_images", True)
        enable_ocr: bool = opts.get("to_md_enable_ocr", False)
        ocr_language = str(opts.get("ocr_language") or "auto")
        current_locale = str(opts.get("locale") or "zh_CN")
        policy = request_policy or build_presentation_markdown_request_policy(context, opts)
        image_mode = policy.export.image_extraction_mode.strip().lower()
        if image_mode not in {"file", "base64", "embed", "omit"}:
            image_mode = "file"
        ocr_placement = policy.export.ocr_placement_mode.strip().lower()
        if ocr_placement not in {"image_md", "main_md"}:
            ocr_placement = "main_md"
        image_link_style = policy.export.image_link_style
        md_file_link_style = policy.export.md_file_link_style

        shapes = list(getattr(slide, "shapes", []))
        # Sort shapes by position (top, then left) for logical order
        shapes.sort(key=lambda s: (_safe_int(getattr(s, "top", 0)), _safe_int(getattr(s, "left", 0))))

        lines: list[str] = []
        table_count = 0
        image_count = 0
        smartart_text_count = 0
        table_shape_index = 0
        chart_shape_index = 0
        media_shape_indices = {"audio": 0, "video": 0}
        image_shape_index = 0

        # Determine slide title
        slide_title = self._extract_slide_title(slide)
        if slide_title:
            lines.append(f"## {slide_title}")
        else:
            lines.append(f"## Slide {slide_index}")
        lines.append("")

        for shape in shapes:
            # Process tables
            if getattr(shape, "has_table", False):
                table_shape_index += 1
                try:
                    table_md = self._table_to_markdown(shape.table)
                except Exception as exc:
                    self._record_payload_warning(
                        context,
                        payload_stats,
                        code="PPTX-TABLE-UNAVAILABLE",
                        message="A table was detected but could not be converted to Markdown.",
                        location=f"slide {slide_index}: table {table_shape_index}",
                        error=exc,
                    )
                    continue
                if table_md:
                    lines.append(table_md)
                    lines.append("")
                    table_count += 1
                continue

            if getattr(shape, "has_chart", False):
                chart_shape_index += 1
                try:
                    chart_lines = self._preserve_chart_payload(
                        shape,
                        slide_index,
                        context,
                        resource_stem=resource_stem,
                        chart_index=chart_shape_index,
                        payload_stats=payload_stats,
                    )
                except Exception as exc:
                    self._record_payload_warning(
                        context,
                        payload_stats,
                        code="PPTX-CHART-UNAVAILABLE",
                        message="A chart was detected but its semantic data could not be preserved.",
                        location=f"slide {slide_index}: chart {chart_shape_index}",
                        error=exc,
                    )
                    continue
                if chart_lines:
                    lines.extend(chart_lines)
                    lines.append("")
                    if payload_stats is not None:
                        payload_stats["charts"] = int(payload_stats["charts"]) + 1
                        payload_stats["chart_snapshot_unavailable_locations"].append(
                            f"slide {slide_index}: chart {chart_shape_index}"
                        )
                continue

            media_kind = self._shape_media_kind(shape)
            if media_kind is not None:
                media_shape_indices[media_kind] += 1
                media_shape_index = media_shape_indices[media_kind]
                try:
                    media_lines = self._preserve_media_payload(
                        shape,
                        slide,
                        slide_index,
                        context,
                        resource_stem=resource_stem,
                        media_kind=media_kind,
                        media_index=media_shape_index,
                        payload_stats=payload_stats,
                    )
                except Exception as exc:
                    self._record_payload_warning(
                        context,
                        payload_stats,
                        code="PPTX-MEDIA-UNAVAILABLE",
                        message=f"A {media_kind} payload was detected but could not be exported.",
                        location=f"slide {slide_index}: {media_kind} {media_shape_index}",
                        error=exc,
                    )
                    continue
                if media_lines:
                    lines.extend(media_lines)
                    lines.append("")
                    image_count += 1
                    if payload_stats is not None:
                        payload_stats[media_kind] = int(payload_stats[media_kind]) + 1
                continue

            # Process images (H5)
            if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                image_shape_index += 1
                image_count += 1
                img = shape.image
                ext = (getattr(img, "ext", "") or "png").lstrip(".").lower()
                image_blob = bytes(img.blob)
                image_digest = hashlib.sha256(image_blob).hexdigest()[:12]
                display_filename = f"slide{slide_index}_img{image_count}.{ext}"
                filename = f"{resource_stem}_{PurePosixPath(display_filename).stem}_{image_digest}.{ext}"

                img_path: str = ""
                extract_image_blob = keep_images or enable_ocr
                if extract_image_blob:
                    try:
                        img_path = context.workspace.create_artifact_path("auxiliary", f".{ext}")
                        with open(img_path, "wb") as f:
                            f.write(image_blob)
                    except OSError as exc:
                        self._record_payload_warning(
                            context,
                            payload_stats,
                            code="PPTX-IMAGE-UNAVAILABLE",
                            message="An image was detected but could not be exported.",
                            location=f"slide {slide_index}: image {image_shape_index}",
                            error=exc,
                        )
                        continue

                    if keep_images and image_mode in {"file", "embed"}:
                        from docwen_core.formats.categories import get_media_type
                        from docwen_core.models.artifact import ArtifactManifest

                        img_artifact = ArtifactManifest(
                            artifact_id=str(uuid.uuid4()),
                            kind="image",
                            staging_path=img_path,
                            suggested_name=filename,
                            media_type=get_media_type(ext),
                            is_primary=False,
                        )
                        context.workspace.add_artifact(img_artifact)

                    from docwen_core.text.image_markdown import generate_image_markdown

                    image_path = img_path if image_mode == "base64" else filename
                    image_markdown = ""
                    if keep_images:
                        image_markdown = generate_image_markdown(
                            image_path=image_path,
                            image_mode=image_mode,
                            image_link_style=image_link_style,
                            alt_text=display_filename,
                            export_semantics=policy.export,
                        )

                    # OCR support
                    if enable_ocr:
                        try:
                            from docwen_core.detection import detect_content_format
                            from docwen_core.text.ocr import run_ocr_outcome

                            outcome = run_ocr_outcome(
                                img_path,
                                source_format=detect_content_format(img_path).format,
                                ocr_language=ocr_language,
                                current_locale=current_locale,
                            )
                            _report_ocr_best_effort(
                                context.progress,
                                outcome.status,
                                location=f"slide {slide_index}: {display_filename}",
                            )
                            ocr_text = outcome.recognized_text
                            if ocr_text:
                                if ocr_placement == "image_md":
                                    from docwen_core.text.image_markdown import build_image_ocr_sidecar

                                    sidecar_base_stem = f"{PurePosixPath(filename).stem}_ocr"
                                    ocr_blockquote_title = policy.ocr_blockquote_title
                                    provisional_text, _replacement_link = build_image_ocr_sidecar(
                                        sidecar_stem=sidecar_base_stem,
                                        source_format="pptx",
                                        image_markdown=image_markdown,
                                        ocr_text=ocr_text,
                                        md_link_style=md_file_link_style,
                                        ocr_blockquote_title=ocr_blockquote_title,
                                        yaml_key_labels=opts.get("yaml_key_labels"),
                                    )
                                    sidecar_digest = hashlib.sha256(provisional_text.encode("utf-8")).hexdigest()[:12]
                                    sidecar_stem = f"{sidecar_base_stem}_{sidecar_digest}"
                                    sidecar_text, replacement_link = build_image_ocr_sidecar(
                                        sidecar_stem=sidecar_stem,
                                        source_format="pptx",
                                        image_markdown=image_markdown,
                                        ocr_text=ocr_text,
                                        md_link_style=md_file_link_style,
                                        ocr_blockquote_title=ocr_blockquote_title,
                                        yaml_key_labels=opts.get("yaml_key_labels"),
                                    )
                                    sidecar_path = context.workspace.create_artifact_path("auxiliary", ".md")
                                    with open(sidecar_path, "w", encoding="utf-8") as handle:
                                        handle.write(sidecar_text)
                                    sidecar_artifact = ArtifactManifest(
                                        artifact_id=str(uuid.uuid4()),
                                        kind="auxiliary",
                                        staging_path=sidecar_path,
                                        suggested_name=f"{sidecar_stem}.md",
                                        media_type="text/markdown",
                                        metadata={"source_format": "pptx", "ocr": True, "image": filename},
                                        is_primary=False,
                                    )
                                    context.workspace.add_artifact(sidecar_artifact)
                                    image_markdown = replacement_link
                                else:
                                    ocr_lines = []
                                    ocr_title = policy.ocr_blockquote_title
                                    if ocr_title:
                                        ocr_lines.append(f"> **{ocr_title}**")
                                        ocr_lines.append(">")
                                    for ocr_line in ocr_text.split("\n"):
                                        stripped = ocr_line.strip()
                                        if stripped:
                                            ocr_lines.append(f"> {stripped}")
                                        else:
                                            ocr_lines.append(">")
                                    if ocr_lines:
                                        image_markdown = f"{image_markdown}\n\n" + "\n".join(ocr_lines)
                        except ImportError:
                            pass

                    lines.append(image_markdown)
                    lines.append("")
                continue

            smartart_texts, smartart_detected, failed_smartart_relationships = self._extract_smartart_texts(
                shape, slide
            )
            for relationship_id in failed_smartart_relationships:
                self._record_payload_warning(
                    context,
                    payload_stats,
                    code="PPTX-SMARTART-UNAVAILABLE",
                    message="SmartArt was detected but one diagram payload could not be read.",
                    location=f"slide {slide_index}: SmartArt relationship {relationship_id}",
                    error="diagram relationship or XML is unavailable",
                )
            if smartart_texts:
                lines.extend(f"- {text}" for text in smartart_texts)
                lines.append("")
                smartart_text_count += len(smartart_texts)
                continue
            if smartart_detected:
                continue

            # Process text frames
            if getattr(shape, "has_text_frame", False):
                text = (shape.text_frame.text or "").strip()
                if text:
                    lines.append(text)
                    lines.append("")

        # H4 — Speaker notes
        if export_notes:
            try:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    lines.append("\n> Notes:")
                    for note_line in notes_text.split("\n"):
                        lines.append(f"> {note_line}")
                    lines.append("")
            except Exception as exc:
                context.logger.warning(f"Unable to read notes for slide {slide_index}: {exc}")
                if payload_stats is not None:
                    payload_stats.setdefault("notes_unavailable_locations", []).append(
                        f"slide {slide_index}: {type(exc).__name__}: {exc}"
                    )

        return lines, image_count, table_count, smartart_text_count

    @staticmethod
    def _shape_media_kind(shape: Any) -> str | None:
        element = getattr(shape, "element", None)
        if element is None:
            return None
        try:
            if element.xpath(".//*[local-name()='audioFile']"):
                return "audio"
            if element.xpath(".//*[local-name()='videoFile']"):
                return "video"
        except Exception:
            return None
        return None

    @staticmethod
    def _related_payload(owner_part: Any, relationship_id: str) -> tuple[bytes, str]:
        relationship = owner_part.rels[relationship_id]
        target_part = relationship.target_part
        blob = bytes(target_part.blob)
        partname = str(getattr(target_part, "partname", ""))
        extension = PurePosixPath(partname).suffix.lstrip(".").lower()
        return blob, extension

    @staticmethod
    def _register_payload_artifact(
        context: Any,
        *,
        payload: bytes,
        extension: str,
        suggested_name: str,
        media_type: str,
        kind: str,
        metadata: dict[str, Any],
    ) -> None:
        from docwen_core.models.artifact import ArtifactManifest

        staging_path = context.workspace.create_artifact_path("auxiliary", f".{extension}")
        with open(staging_path, "wb") as handle:
            handle.write(payload)
        context.workspace.add_artifact(
            ArtifactManifest(
                artifact_id=str(uuid.uuid4()),
                kind=kind,
                staging_path=staging_path,
                suggested_name=suggested_name,
                media_type=media_type,
                metadata=metadata,
                is_primary=False,
            )
        )

    def _preserve_chart_payload(
        self,
        shape: Any,
        slide_index: int,
        context: Any,
        *,
        resource_stem: str,
        chart_index: int,
        payload_stats: dict[str, Any] | None,
    ) -> list[str]:
        chart = shape.chart
        plots = list(chart.plots)
        series = list(chart.series)
        if not plots:
            raise ValueError("chart contains no readable plots")

        title = ""
        try:
            if chart.has_title:
                title = " ".join(
                    paragraph.text.strip()
                    for paragraph in chart.chart_title.text_frame.paragraphs
                    if paragraph.text.strip()
                )
        except Exception:
            title = ""
        series_names = [str(getattr(item, "name", "") or "").strip() for item in series]
        display_title = title or next((name for name in series_names if name), f"Chart {chart_index}")
        categories = [str(category.label) for category in plots[0].categories]
        values = [list(getattr(item, "values", ())) for item in series]

        lines = [f"### Chart: {display_title}"]
        headers = ["Category", *[name or f"Series {index}" for index, name in enumerate(series_names, 1)]]
        table_lines = [
            "| " + " | ".join(headers) + " |",
            "| " + " | ".join("---" for _ in headers) + " |",
        ]
        for category_index, category in enumerate(categories):
            row = [category]
            for series_values in values:
                value = series_values[category_index] if category_index < len(series_values) else ""
                row.append("" if value is None else str(value))
            table_lines.append("| " + " | ".join(row) + " |")
        lines.append("\n".join(table_lines))

        try:
            package_relationships = [
                relationship
                for relationship in chart.part.rels.values()
                if str(relationship.reltype).endswith("/package")
            ]
        except Exception as exc:
            self._record_payload_warning(
                context,
                payload_stats,
                code="PPTX-CHART-WORKBOOK-UNAVAILABLE",
                message="Chart semantics were preserved, but its embedded workbook relationships could not be read.",
                location=f"slide {slide_index}: chart {chart_index} workbook",
                error=exc,
            )
            return lines
        if package_relationships:
            try:
                workbook = bytes(package_relationships[0].target_part.blob)
                workbook_digest = hashlib.sha256(workbook).hexdigest()[:12]
                workbook_name = f"{resource_stem}_slide{slide_index}_chart{chart_index}_workbook_{workbook_digest}.xlsx"
                self._register_payload_artifact(
                    context,
                    payload=workbook,
                    extension="xlsx",
                    suggested_name=workbook_name,
                    media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    kind="auxiliary",
                    metadata={
                        "source_format": "pptx",
                        "payload": "chart_workbook",
                        "slide": slide_index,
                        "chart": chart_index,
                    },
                )
                lines.append(f"[Embedded chart workbook]({workbook_name})")
            except Exception as exc:
                self._record_payload_warning(
                    context,
                    payload_stats,
                    code="PPTX-CHART-WORKBOOK-UNAVAILABLE",
                    message="Chart semantics were preserved, but its embedded workbook could not be exported.",
                    location=f"slide {slide_index}: chart {chart_index} workbook",
                    error=exc,
                )
        return lines

    def _preserve_media_payload(
        self,
        shape: Any,
        slide: Any,
        slide_index: int,
        context: Any,
        *,
        resource_stem: str,
        media_kind: str,
        media_index: int,
        payload_stats: dict[str, Any] | None,
    ) -> list[str]:
        element = getattr(shape, "element", None)
        slide_part = getattr(slide, "part", None)
        if element is None or slide_part is None:
            raise ValueError("media shape has no relationship owner")
        media_ids = element.xpath(f".//*[local-name()='{media_kind}File']/@*[local-name()='link']")
        poster_ids = element.xpath(".//*[local-name()='blip']/@*[local-name()='embed']")
        if not media_ids:
            raise ValueError(f"{media_kind} relationship ID is missing")

        media_blob, media_extension = self._related_payload(slide_part, str(media_ids[0]))
        expected_extension = "mp3" if media_kind == "audio" else "mp4"
        media_extension = media_extension or expected_extension
        media_digest = hashlib.sha256(media_blob).hexdigest()[:12]
        media_name = f"{resource_stem}_slide{slide_index}_{media_kind}{media_index}_{media_digest}.{media_extension}"
        self._register_payload_artifact(
            context,
            payload=media_blob,
            extension=media_extension,
            suggested_name=media_name,
            media_type="audio/mpeg" if media_kind == "audio" else "video/mp4",
            kind="media",
            metadata={
                "source_format": "pptx",
                "payload": media_kind,
                "slide": slide_index,
                media_kind: media_index,
            },
        )

        label = "Audio" if media_kind == "audio" else "Video"
        lines = [f"[{label} payload]({media_name})"]
        if poster_ids:
            try:
                poster_blob, poster_extension = self._related_payload(slide_part, str(poster_ids[0]))
                poster_extension = poster_extension or "png"
                poster_digest = hashlib.sha256(poster_blob).hexdigest()[:12]
                poster_name = (
                    f"{resource_stem}_slide{slide_index}_{media_kind}{media_index}_poster_"
                    f"{poster_digest}.{poster_extension}"
                )
                poster_media_type = "image/jpeg" if poster_extension in {"jpg", "jpeg"} else "image/png"
                self._register_payload_artifact(
                    context,
                    payload=poster_blob,
                    extension=poster_extension,
                    suggested_name=poster_name,
                    media_type=poster_media_type,
                    kind="image",
                    metadata={
                        "source_format": "pptx",
                        "payload": f"{media_kind}_poster",
                        "slide": slide_index,
                        media_kind: media_index,
                    },
                )
                lines.append(f"![{label} poster]({poster_name})")
            except Exception as exc:
                self._record_payload_warning(
                    context,
                    payload_stats,
                    code="PPTX-MEDIA-POSTER-UNAVAILABLE",
                    message=f"The {media_kind} payload was preserved, but its poster image could not be exported.",
                    location=f"slide {slide_index}: {media_kind} {media_index} poster",
                    error=exc,
                )
        return lines

    @staticmethod
    def _extract_slide_title(slide: Any) -> str:
        """Extract the slide title from title/centertitle placeholder shapes."""
        for shape in getattr(slide, "shapes", []):
            if not getattr(shape, "has_text_frame", False):
                continue
            if not getattr(shape, "is_placeholder", False):
                continue
            try:
                ph = shape.placeholder_format
                ph_type = getattr(ph, "type", None)
                # 1 = TITLE, 3 = CENTER_TITLE
                if ph_type in {1, 3}:
                    text = (shape.text_frame.text or "").strip()
                    if text:
                        return text
            except Exception:
                continue
        return ""

    @staticmethod
    def _table_to_markdown(table: Any) -> str:
        """Convert a python-pptx table to a Markdown table string."""
        rows = table.rows
        cols = table.columns
        if len(rows) == 0 or len(cols) == 0:
            return ""

        grid: list[list[str]] = []
        for row in rows:
            row_cells: list[str] = []
            for cell in row.cells:
                row_cells.append((cell.text or "").replace("\n", " ").strip())
            grid.append(row_cells)

        header = grid[0]

        def esc(value: str) -> str:
            return value.replace("|", "\\|")

        out_lines: list[str] = []
        out_lines.append("| " + " | ".join(esc(value) for value in header) + " |")
        out_lines.append("| " + " | ".join("---" for _ in header) + " |")

        for row in grid[1:]:
            if len(row) < len(header):
                row = row + [""] * (len(header) - len(row))
            out_lines.append("| " + " | ".join(esc(value) for value in row[: len(header)]) + " |")

        return "\n".join(out_lines)


def _safe_int(v: Any) -> int:
    """Safely convert a value to int, returning 0 on failure."""
    try:
        return int(v)
    except Exception:
        return 0
