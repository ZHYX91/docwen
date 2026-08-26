"""Golden / semantic parity tests for PPTX/PPT → MD input routes.

Coverage:
  ROUTE-PPTX-001  (pptx → md)
  ROUTE-PPT-001   (ppt → md, via external office bridge)

All tests use the full runtime pipeline with PresentationPlugin.
"""

from __future__ import annotations

import base64
import json
import os
import re
import tempfile
from pathlib import Path
from typing import Any
from zipfile import ZIP_DEFLATED, ZipFile

import pytest

from docwen_core.models.file_ref import FileRef
from docwen_core.models.request import ConversionRequest, OutputPolicy

pytestmark = [pytest.mark.golden, pytest.mark.contract]

_PROJECT_ROOT = Path(__file__).resolve().parents[4]

_PPTX_OLD_SYSTEM_FIXTURE = (
    _PROJECT_ROOT / "tests" / "fixtures" / "golden" / "old_system_pptx_to_markdown_semantics.json"
)

_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

_P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"


@pytest.fixture
def pipeline():
    """Build the full runtime pipeline with PresentationPlugin."""
    from docwen_plugin_presentation import PresentationPlugin
    from docwen_runtime.engine.route_resolver import RouteResolver
    from docwen_runtime.engine.task_manager import TaskManager
    from docwen_runtime.output.finalizer import OutputFinalizer
    from docwen_runtime.plugin_registry.registry import PluginRegistry
    from docwen_runtime.workspace.manager import WorkspaceManager

    plugin = PresentationPlugin()
    registry = PluginRegistry()
    registry.register(plugin)

    resolver = RouteResolver(registry)
    ws_root = tempfile.mkdtemp(prefix="docwen_pres_")
    ws_mgr = WorkspaceManager(root_dir=ws_root)
    finalizer = OutputFinalizer()
    task_mgr = TaskManager(registry, resolver, ws_mgr, finalizer)

    yield plugin, task_mgr, ws_mgr
    ws_mgr.cleanup_all()
    import shutil

    shutil.rmtree(ws_root, ignore_errors=True)


def _run_request(
    task_mgr,
    input_path,
    source_format,
    output_dir,
    *,
    config_snapshot: dict[str, Any] | None = None,
    on_event: Any = None,
    **options,
) -> Any:
    """Run a single conversion request through the task manager."""
    opts: dict = {"to_md_keep_images": True}
    opts.update(options)

    request = ConversionRequest(
        request_id="pres-route-test",
        input_refs=[
            FileRef(
                path=str(input_path),
                format=source_format,
                category="presentation",
                size_bytes=Path(input_path).stat().st_size,
            )
        ],
        target_format="md",
        output_policy=OutputPolicy(output_dir=str(output_dir)),
        options=opts,
        config_snapshot=config_snapshot or {},
    )
    return task_mgr.execute_single(request, on_event=on_event)


def _deliverable_artifacts(result: Any) -> list[Any]:
    manifests = [
        artifact for artifact in result.artifacts if artifact.media_type == "application/vnd.docwen.document-node+json"
    ]
    assert len(manifests) == 1
    return [artifact for artifact in result.artifacts if artifact not in manifests]


def _document_node_root(path: Path, output_dir: Path) -> Path:
    relative = path.relative_to(output_dir)
    assert len(relative.parts) >= 2
    root = output_dir / relative.parts[0]
    assert root.is_dir()
    return root


def _io_path(path: str | Path) -> Path:
    from docwen_runtime.path_io import filesystem_path

    return filesystem_path(os.path.normpath(str(path)))


def _linked_path(markdown_path: Path, target: str) -> Path:
    return _io_path(os.path.join(markdown_path.parent, target))


def _legacy_markdown_projection(content: str, artifacts: list[Any]) -> str:
    for artifact in artifacts:
        if artifact.media_type != "text/markdown" or artifact.is_primary or artifact.logical_path is None:
            continue
        legacy_name = artifact.metadata.get("source_suggested_name")
        if isinstance(legacy_name, str):
            content = content.replace(artifact.logical_path.split("/", 1)[1], legacy_name)
    return content.replace("![[../", "![[").replace("](../", "](")


def _load_pptx_old_system_fixture() -> dict[str, Any]:
    return json.loads(_PPTX_OLD_SYSTEM_FIXTURE.read_text(encoding="utf-8"))


def _build_single_image_pptx(
    tmp_path: Path,
    source_dir_name: str,
    color: tuple[int, int, int],
    *,
    filename_stem: str | None = None,
) -> tuple[Path, bytes]:
    """Build one image-only PPTX and return its exact embedded PNG payload."""
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    source_dir = tmp_path / source_dir_name
    source_dir.mkdir()
    image_path = source_dir / "payload.png"
    Image.new("RGB", (4, 4), color).save(image_path)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), width=Inches(1), height=Inches(1))
    pptx_path = source_dir / f"{filename_stem or source_dir_name}.pptx"
    presentation.save(str(pptx_path))
    return pptx_path, image_path.read_bytes()


def _build_policy03_chart_pptx(tmp_path: Path) -> tuple[Path, bytes]:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    chart_data: Any = CategoryChartData()
    chart_data.categories = ["1st Qtr", "2nd Qtr", "3rd Qtr", "4th Qtr"]
    chart_data.add_series("Sales", (8.2, 3.2, 1.4, 1.2))
    chart_frame: Any = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(4),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Quarterly Sales"
    pptx_path = tmp_path / "policy03-chart.pptx"
    presentation.save(str(pptx_path))
    with ZipFile(pptx_path) as package:
        workbook_name = next(name for name in package.namelist() if name.startswith("ppt/embeddings/"))
        workbook_bytes = package.read(workbook_name)
    return pptx_path, workbook_bytes


def _build_policy03_media_pptx(tmp_path: Path, kind: str) -> tuple[Path, bytes, bytes]:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    poster_path = tmp_path / f"{kind}-poster.png"
    Image.new("RGB", (12, 8), (21, 42, 84)).save(poster_path)
    poster_bytes = poster_path.read_bytes()
    extension = "mp3" if kind == "audio" else "mp4"
    media_bytes = f"policy03-{kind}-payload".encode()
    media_path = tmp_path / f"payload.{extension}"
    media_path.write_bytes(media_bytes)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_movie(
        str(media_path),
        Inches(1),
        Inches(1),
        Inches(3),
        Inches(2),
        poster_frame_image=str(poster_path),
        mime_type="audio/mpeg" if kind == "audio" else "video/mp4",
    )
    pptx_path = tmp_path / f"policy03-{kind}.pptx"
    presentation.save(str(pptx_path))

    if kind == "audio":
        replacement = tmp_path / "policy03-audio-rewritten.pptx"
        with ZipFile(pptx_path) as source, ZipFile(replacement, "w", ZIP_DEFLATED) as target:
            for info in source.infolist():
                name = info.filename.replace("ppt/media/media1.mp4", "ppt/media/media1.mp3")
                payload = source.read(info.filename)
                if info.filename in {
                    "[Content_Types].xml",
                    "ppt/slides/slide1.xml",
                    "ppt/slides/_rels/slide1.xml.rels",
                }:
                    payload = (
                        payload.replace(b"videoFile", b"audioFile")
                        .replace(b"/relationships/video", b"/relationships/audio")
                        .replace(b"media1.mp4", b"media1.mp3")
                        .replace(
                            b'Extension="mp4" ContentType="video/mp4"', b'Extension="mp3" ContentType="audio/mpeg"'
                        )
                    )
                target.writestr(name, payload)
        pptx_path = replacement
    return pptx_path, media_bytes, poster_bytes


def _markdown_targets(markdown: str) -> list[str]:
    """Return regular Markdown and wiki-link targets used by artifact tests."""
    targets = re.findall(r"!?\[[^\]]*\]\(([^)]+)\)", markdown)
    for raw_target in re.findall(r"!?\[\[([^\]]+)\]\]", markdown):
        targets.append(raw_target.rsplit("|", 1)[-1])
    return targets


def _data_uri_payload(markdown: str) -> tuple[str, bytes]:
    match = re.search(r"data:([^;]+);base64,([A-Za-z0-9+/=]+)", markdown)
    assert match is not None
    return match.group(1), base64.b64decode(match.group(2))


@pytest.fixture
def sample_pptx_file(tmp_path: Path) -> Path:
    """Create a sample PPTX file with slides containing text and a table."""
    from pptx import Presentation
    from pptx.shapes.placeholder import SlidePlaceholder
    from pptx.util import Inches

    pres = Presentation()
    pres.core_properties.title = "Test Presentation"

    # Slide 1: Title slide
    slide_layout = pres.slide_layouts[0]
    slide1 = pres.slides.add_slide(slide_layout)
    title_shape = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    assert isinstance(title_shape, SlidePlaceholder)
    assert isinstance(subtitle, SlidePlaceholder)
    title_shape.text = "Presentation Title"
    subtitle.text = "A subtitle for the presentation"

    # Slide 2: Content with bullets
    slide_layout2 = pres.slide_layouts[1]
    slide2 = pres.slides.add_slide(slide_layout2)
    title_shape2 = slide2.shapes.title
    body = slide2.shapes.placeholders[1]
    assert isinstance(title_shape2, SlidePlaceholder)
    assert isinstance(body, SlidePlaceholder)
    title_shape2.text = "Agenda"
    tf = body.text_frame
    tf.text = "First item"
    p = tf.add_paragraph()
    p.text = "Second item"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Third item"

    # Slide 3: Content with table
    slide_layout3 = pres.slide_layouts[1]
    slide3 = pres.slides.add_slide(slide_layout3)
    title_shape3 = slide3.shapes.title
    assert isinstance(title_shape3, SlidePlaceholder)
    title_shape3.text = "Data Summary"

    rows, cols = 3, 3
    table_shape = slide3.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(3))
    table = table_shape.table
    headers = ["Category", "Q1", "Q2"]
    data = [["Alpha", "100", "200"], ["Beta", "150", "250"]]
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            table.cell(i + 1, j).text = val

    path = tmp_path / "test.pptx"
    pres.save(str(path))
    return path


@pytest.fixture
def sample_pptx_with_image(tmp_path: Path) -> Path:
    """Create a PPTX with one embedded image for link-style assertions."""
    from pptx import Presentation
    from pptx.util import Inches

    image_path = tmp_path / "tiny.png"
    try:
        from PIL import Image

        Image.new("RGB", (2, 2), color=(255, 0, 0)).save(image_path)
    except Exception:
        image_path.write_bytes(
            base64.b64decode(
                "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAIAAACQd1PeAAAADUlEQVR42mP8z8BQDwAFgwJ/lI9QOgAAAABJRU5ErkJggg=="
            )
        )

    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), width=Inches(1), height=Inches(1))

    path = tmp_path / "image-link-style.pptx"
    pres.save(str(path))
    return path


@pytest.fixture
def sample_pptx_with_sections(tmp_path: Path) -> Path:
    """Create a PPTX with PowerPoint section metadata like the old-system tests."""
    from pptx import Presentation
    from pptx.shapes.placeholder import SlidePlaceholder

    pres = Presentation()
    pres.core_properties.title = "Sectioned Presentation"

    slide1 = pres.slides.add_slide(pres.slide_layouts[0])
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    assert isinstance(title1, SlidePlaceholder)
    assert isinstance(subtitle1, SlidePlaceholder)
    title1.text = "Slide One"
    subtitle1.text = "First section body"

    slide2 = pres.slides.add_slide(pres.slide_layouts[0])
    title2 = slide2.shapes.title
    subtitle2 = slide2.placeholders[1]
    assert isinstance(title2, SlidePlaceholder)
    assert isinstance(subtitle2, SlidePlaceholder)
    title2.text = "Slide Two"
    subtitle2.text = "Second section body"

    path = tmp_path / "with-sections.pptx"
    pres.save(str(path))
    _inject_pptx_sections(path, {"First Section": 0, "Second Section": 1})
    return path


@pytest.fixture
def sample_pptx_with_notes(tmp_path: Path) -> Path:
    """Create a PPTX with speaker notes for export_notes parity."""
    from pptx import Presentation
    from pptx.shapes.placeholder import SlidePlaceholder

    pres = Presentation()
    pres.core_properties.title = "Notes Presentation"
    slide = pres.slides.add_slide(pres.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    notes_text_frame = slide.notes_slide.notes_text_frame
    assert isinstance(title, SlidePlaceholder)
    assert isinstance(subtitle, SlidePlaceholder)
    assert notes_text_frame is not None
    title.text = "Notes Slide"
    subtitle.text = "Visible body"
    notes_text_frame.text = "Speaker note line"

    path = tmp_path / "with-notes.pptx"
    pres.save(str(path))
    return path


def _inject_pptx_sections(path: Path, section_to_slide_index: dict[str, int]) -> None:
    from lxml import etree

    with ZipFile(path, "r") as zin:
        file_map = {info.filename: zin.read(info.filename) for info in zin.infolist()}

    root = etree.fromstring(file_map["ppt/presentation.xml"])
    slide_ids_raw = root.xpath(".//p:sldId/@id", namespaces={"p": _P_NS})
    assert isinstance(slide_ids_raw, list)
    slide_ids = [str(slide_id) for slide_id in slide_ids_raw]
    assert len(slide_ids) >= len(section_to_slide_index)

    ext_lst = root.find(f"{{{_P_NS}}}extLst")
    if ext_lst is None:
        ext_lst = etree.SubElement(root, f"{{{_P_NS}}}extLst")

    ext = etree.SubElement(ext_lst, f"{{{_P_NS}}}ext", uri="{521415D9-36F7-43E2-AB2F-B90AF26B5E84}")
    section_lst = etree.SubElement(ext, f"{{{_P14_NS}}}sectionLst", nsmap={"p14": _P14_NS})
    for idx, (section_name, slide_index) in enumerate(section_to_slide_index.items(), start=1):
        section = etree.SubElement(
            section_lst,
            f"{{{_P14_NS}}}section",
            name=section_name,
            id=f"{{00000000-0000-0000-0000-00000000000{idx}}}",
        )
        slide_id_list = etree.SubElement(section, f"{{{_P14_NS}}}sldIdLst")
        etree.SubElement(slide_id_list, f"{{{_P14_NS}}}sldId", id=str(slide_ids[slide_index]))

    file_map["ppt/presentation.xml"] = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

    with ZipFile(path, "w", compression=ZIP_DEFLATED) as zout:
        for filename, data in file_map.items():
            zout.writestr(filename, data)


__all__ = (
    "Any",
    "ConversionRequest",
    "FileRef",
    "OutputPolicy",
    "Path",
    "_build_policy03_chart_pptx",
    "_build_policy03_media_pptx",
    "_build_single_image_pptx",
    "_data_uri_payload",
    "_deliverable_artifacts",
    "_document_node_root",
    "_io_path",
    "_legacy_markdown_projection",
    "_linked_path",
    "_load_pptx_old_system_fixture",
    "_markdown_targets",
    "_run_request",
    "pipeline",
    "pytest",
    "pytestmark",
    "re",
    "sample_pptx_file",
    "sample_pptx_with_image",
    "sample_pptx_with_notes",
    "sample_pptx_with_sections",
)
