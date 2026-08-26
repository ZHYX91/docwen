"""Focused contracts for SmartArt text and hidden-slide handling."""

from __future__ import annotations

from types import SimpleNamespace
from typing import Any

import pytest
from lxml import etree
from tests.support.config import FakeConfigView

from docwen_plugin_presentation.pptx_md.converter import PptxToMarkdownConverter

pytestmark = pytest.mark.unit


def _context() -> Any:
    return SimpleNamespace(
        cancellation=SimpleNamespace(check=lambda: None),
        progress=SimpleNamespace(report_progress=lambda *_args, **_kwargs: None),
        request=SimpleNamespace(options={}),
        config=FakeConfigView(),
    )


def test_process_slide_extracts_smartart_nodes_in_relationship_order() -> None:
    shape_xml = etree.fromstring(
        b"""
        <p:graphicFrame
          xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
          xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"
          xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
          <a:graphic>
            <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/diagram">
              <dgm:relIds r:dm="rId3"/>
            </a:graphicData>
          </a:graphic>
        </p:graphicFrame>
        """
    )
    diagram_xml = b"""
        <dgm:dataModel
          xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <dgm:ptLst>
            <dgm:pt modelId="root"/>
            <dgm:pt modelId="parent"><dgm:t><a:p><a:r><a:t>Parent</a:t></a:r></a:p></dgm:t></dgm:pt>
            <dgm:pt modelId="sibling"><dgm:t><a:p><a:r><a:t>Sibling</a:t></a:r></a:p></dgm:t></dgm:pt>
            <dgm:pt modelId="child"><dgm:t><a:p><a:r><a:t>Child</a:t></a:r></a:p></dgm:t></dgm:pt>
          </dgm:ptLst>
          <dgm:cxnLst>
            <dgm:cxn srcId="root" destId="sibling" srcOrd="1"/>
            <dgm:cxn srcId="parent" destId="child" srcOrd="0"/>
            <dgm:cxn srcId="root" destId="parent" srcOrd="0"/>
          </dgm:cxnLst>
        </dgm:dataModel>
    """
    relationship = SimpleNamespace(
        reltype="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData",
        target_part=SimpleNamespace(blob=diagram_xml),
    )
    shape = SimpleNamespace(
        element=shape_xml,
        top=0,
        left=0,
        has_table=False,
        shape_type=None,
        has_text_frame=False,
    )
    slide = SimpleNamespace(shapes=[shape], part=SimpleNamespace(rels={"rId3": relationship}))

    lines, image_count, table_count, smartart_text_count = PptxToMarkdownConverter()._process_slide(
        slide,
        1,
        _context(),
    )

    assert lines == ["## Slide 1", "", "- Parent", "- Child", "- Sibling", ""]
    assert image_count == 0
    assert table_count == 0
    assert smartart_text_count == 3


def test_parse_pptx_includes_hidden_slide_content_and_counts_policy(tmp_path) -> None:
    from pptx import Presentation

    source = tmp_path / "hidden-policy.pptx"
    presentation = Presentation()
    visible = presentation.slides.add_slide(presentation.slide_layouts[5])
    visible_title = visible.shapes.title
    assert visible_title is not None
    visible_title.text = "Visible slide"
    hidden = presentation.slides.add_slide(presentation.slide_layouts[5])
    hidden_title = hidden.shapes.title
    assert hidden_title is not None
    hidden_title.text = "Hidden source content"
    hidden.element.set("show", "0")
    presentation.save(str(source))

    content, stats = PptxToMarkdownConverter()._parse_pptx(str(source), _context())

    assert "## Visible slide" in content
    assert "## Hidden source content" in content
    assert stats["slides"] == 2
    assert stats["hidden_slides"] == 1
    assert stats["smartart_texts"] == 0
