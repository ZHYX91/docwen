"""Focused tests split from test_presentation_to_md.py."""

from __future__ import annotations

import pytest

from ._presentation_to_md_support import (
    Any,
    Path,
    _build_single_image_pptx,
    _data_uri_payload,
    _deliverable_artifacts,
    _document_node_root,
    _io_path,
    _legacy_markdown_projection,
    _linked_path,
    _load_pptx_old_system_fixture,
    _markdown_targets,
    _run_request,
    re,
)
from ._presentation_to_md_support import (
    pipeline as pipeline,
)

pytestmark = [pytest.mark.golden, pytest.mark.contract]
from ._presentation_to_md_support import (
    sample_pptx_file as sample_pptx_file,
)
from ._presentation_to_md_support import (
    sample_pptx_with_image as sample_pptx_with_image,
)
from ._presentation_to_md_support import (
    sample_pptx_with_notes as sample_pptx_with_notes,
)
from ._presentation_to_md_support import (
    sample_pptx_with_sections as sample_pptx_with_sections,
)


class TestPptxToMd:
    """Golden parity tests for ROUTE-PPTX-001: pptx → md."""

    def test_pptx_shared_output_directory_keeps_each_primary_linked_to_its_own_image(
        self,
        pipeline,
        tmp_path: Path,
    ) -> None:
        """Two normal conversions must not silently cross-link generic image names."""
        _plugin, task_mgr, _ws_mgr = pipeline
        output_dir = tmp_path / "shared-image-output"
        output_dir.mkdir()
        sources = [
            _build_single_image_pptx(
                tmp_path,
                "red-source",
                (255, 0, 0),
                filename_stem="共享 deck (v1)",
            ),
            _build_single_image_pptx(
                tmp_path,
                "blue-source",
                (0, 0, 255),
                filename_stem="共享 deck (v1)",
            ),
        ]
        finalized_image_results: list[tuple[Path, bytes]] = []

        for input_path, expected_image_bytes in sources:
            result = _run_request(
                task_mgr,
                input_path,
                "pptx",
                output_dir,
                image_link_style="markdown_embed",
            )

            assert result.success
            primary = next(artifact for artifact in result.artifacts if artifact.is_primary)
            finalized_image_results.append((Path(primary.staging_path), expected_image_bytes))

        for primary_path, expected_image_bytes in finalized_image_results:
            markdown = primary_path.read_text(encoding="utf-8")
            image_targets = [target for target in _markdown_targets(markdown) if target.endswith(".png")]
            assert len(image_targets) == 1
            assert not {" ", "(", ")"}.intersection(image_targets[0])
            linked_image = primary_path.parent / image_targets[0]
            assert linked_image.read_bytes() == expected_image_bytes

        from docwen_plugin_presentation.pptx_md.converter import PptxToMarkdownConverter

        long_cjk_stem = PptxToMarkdownConverter._resource_stem(f"{'演示' * 40} (v1).pptx")
        assert len(long_cjk_stem.encode("utf-8")) <= 64
        longest_sidecar_name = f"{long_cjk_stem}_slide9999_img9999_{'a' * 12}_ocr_{'b' * 12}.md"
        assert len(longest_sidecar_name.encode("utf-8")) <= 255

    def test_pptx_shared_output_directory_keeps_ocr_sidecar_and_image_links_together(
        self,
        pipeline,
        tmp_path: Path,
        monkeypatch,
    ) -> None:
        """Identical images with different OCR text must keep distinct sidecars."""
        from docwen_core.text.ocr import OcrOutcome, OcrStatus

        _plugin, task_mgr, _ws_mgr = pipeline
        output_dir = tmp_path / "shared-ocr-output"
        output_dir.mkdir()
        sources = [
            (
                *_build_single_image_pptx(
                    tmp_path,
                    "first-ocr-source",
                    (32, 64, 128),
                    filename_stem="shared-ocr-deck",
                ),
                "FIRST DECK OCR",
            ),
            (
                *_build_single_image_pptx(
                    tmp_path,
                    "second-ocr-source",
                    (32, 64, 128),
                    filename_stem="shared-ocr-deck",
                ),
                "SECOND DECK OCR",
            ),
        ]
        assert sources[0][1] == sources[1][1]
        ocr_results = iter(ocr_text for _input_path, _image_bytes, ocr_text in sources)
        finalized_ocr_results: list[tuple[Path, bytes, str]] = []

        def _run_payload_ocr(path: str, **_kwargs: Any) -> OcrOutcome:
            assert Path(path).read_bytes() == sources[0][1]
            return OcrOutcome(OcrStatus.SUCCESS, text=next(ocr_results))

        monkeypatch.setattr("docwen_core.text.ocr.run_ocr_outcome", _run_payload_ocr)

        for input_path, expected_image_bytes, expected_ocr_text in sources:
            result = _run_request(
                task_mgr,
                input_path,
                "pptx",
                output_dir,
                to_md_enable_ocr=True,
                ocr_placement="image_md",
                image_link_style="markdown_embed",
            )

            assert result.success
            primary = next(artifact for artifact in result.artifacts if artifact.is_primary)
            finalized_ocr_results.append((Path(primary.staging_path), expected_image_bytes, expected_ocr_text))

        for primary_path, expected_image_bytes, expected_ocr_text in finalized_ocr_results:
            markdown = primary_path.read_text(encoding="utf-8")
            sidecar_targets = [target for target in _markdown_targets(markdown) if target.endswith(".md")]
            assert len(sidecar_targets) == 1
            sidecar_path = _linked_path(primary_path, sidecar_targets[0])
            sidecar = sidecar_path.read_text(encoding="utf-8")
            assert expected_ocr_text in sidecar

            image_targets = [target for target in _markdown_targets(sidecar) if target.endswith(".png")]
            assert len(image_targets) == 1
            linked_image = _linked_path(sidecar_path, image_targets[0])
            assert linked_image.read_bytes() == expected_image_bytes

    def test_pptx_jpeg_image_uses_standard_media_type(self, pipeline, tmp_path: Path) -> None:
        """JPEG artifacts use the IANA ``image/jpeg`` type, not ``image/jpg``."""
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        jpeg_path = tmp_path / "photo.jpg"
        Image.new("RGB", (8, 8), (40, 80, 120)).save(jpeg_path, format="JPEG")
        pres = Presentation()
        slide = pres.slides.add_slide(pres.slide_layouts[6])
        slide.shapes.add_picture(str(jpeg_path), Inches(1), Inches(1))
        input_path = tmp_path / "jpeg-image.pptx"
        pres.save(str(input_path))

        _plugin, task_mgr, _ws_mgr = pipeline
        output_dir = tmp_path / "output_pptx_jpeg"
        output_dir.mkdir()
        result = _run_request(task_mgr, input_path, "pptx", output_dir)

        assert result.success
        image_artifact = next(artifact for artifact in result.artifacts if artifact.kind == "image")
        assert re.fullmatch(
            r"jpeg-image_slide1_img1_[0-9a-f]{12}\.jpg",
            image_artifact.suggested_name,
        )
        assert image_artifact.media_type == "image/jpeg"
        assert Path(image_artifact.staging_path).read_bytes().startswith(b"\xff\xd8\xff")

    def test_pptx_jpeg_base64_uses_standard_media_type(self, pipeline, tmp_path: Path) -> None:
        """Uncompressed JPEG data URIs use ``image/jpeg``, never ``image/jpg``."""
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        jpeg_path = tmp_path / "photo.jpg"
        Image.new("RGB", (8, 8), (40, 80, 120)).save(jpeg_path, format="JPEG")
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(str(jpeg_path), Inches(1), Inches(1))
        input_path = tmp_path / "jpeg-base64.pptx"
        presentation.save(str(input_path))

        _plugin, task_mgr, _ws_mgr = pipeline
        output_dir = tmp_path / "output_pptx_jpeg_base64"
        output_dir.mkdir()
        result = _run_request(
            task_mgr,
            input_path,
            "pptx",
            output_dir,
            config_snapshot={
                "conversion": {
                    "export": {
                        "base64_compress_enabled": False,
                        "base64_compress_threshold_kb": 1,
                    }
                }
            },
            image_mode="base64",
            image_link_style="markdown_embed",
        )

        assert result.success
        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")
        mime, payload = _data_uri_payload(content)
        assert mime == "image/jpeg"
        assert payload.startswith(b"\xff\xd8\xff")
        assert "data:image/jpg;base64," not in content

    def test_pptx_multi_image_order_matches_old_system_projection(
        self,
        pipeline,
        tmp_path: Path,
    ) -> None:
        """Focused PPTX multi-image projection matches old systems and finalizes images."""
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        fixture = _load_pptx_old_system_fixture()
        scope = fixture["image_artifact_scope"]["multi_image_probe"]

        red_path = tmp_path / "red.png"
        blue_path = tmp_path / "blue.png"
        Image.new("RGB", tuple(scope["input"]["image_size_px"]), (255, 0, 0)).save(red_path)
        Image.new("RGB", tuple(scope["input"]["image_size_px"]), (0, 0, 255)).save(blue_path)

        pres = Presentation()
        pres.core_properties.title = scope["input"]["title"]
        slide = pres.slides.add_slide(pres.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(0.4))
        text_box.text_frame.text = scope["input"]["slide_text"]
        slide.shapes.add_picture(str(red_path), Inches(1), Inches(1), width=Inches(1), height=Inches(1))
        slide.shapes.add_picture(str(blue_path), Inches(3), Inches(1), width=Inches(1), height=Inches(1))

        input_path = tmp_path / scope["input"]["filename"]
        pres.save(str(input_path))

        _plugin, task_mgr, ws_mgr = pipeline
        output_dir = tmp_path / "output_pptx_multi_image"
        output_dir.mkdir()
        result = _run_request(
            task_mgr,
            input_path,
            "pptx",
            output_dir,
            image_link_style=scope["input"]["options"]["image_link_style"],
        )

        expected = scope["current_projection"]
        assert result.success is expected["success"]
        assert [diagnostic.code for diagnostic in result.diagnostics] == expected["diagnostic_codes"]
        assert len(_deliverable_artifacts(result)) == expected["artifact_count"]
        primary_artifacts = [artifact for artifact in result.artifacts if artifact.is_primary]
        image_artifacts = [artifact for artifact in result.artifacts if artifact.kind == "image"]
        assert len(primary_artifacts) == expected["primary_artifact_count"]
        assert len(image_artifacts) == expected["image_artifact_count"]
        image_suggested_names = [artifact.suggested_name for artifact in image_artifacts]
        assert all(
            re.fullmatch(pattern, suggested_name)
            for pattern, suggested_name in zip(
                expected["image_suggested_name_patterns"],
                image_suggested_names,
                strict=True,
            )
        )

        markdown_artifact = primary_artifacts[0]
        for key, value in expected["primary_metadata"].items():
            assert markdown_artifact.metadata[key] == value
        node_root = _document_node_root(Path(markdown_artifact.staging_path), output_dir)
        assert Path(markdown_artifact.staging_path).is_file()
        for artifact in image_artifacts:
            assert _document_node_root(Path(artifact.staging_path), output_dir) == node_root
            assert Path(artifact.staging_path).is_file()
            assert Path(artifact.staging_path).read_bytes().startswith(b"\x89PNG\r\n\x1a\n")

        content = Path(markdown_artifact.staging_path).read_text(encoding="utf-8")
        assert str(output_dir) not in content
        assert str(Path(ws_mgr.root_dir)) not in content
        assert (scope["input"]["title"] in content) is expected["markdown_contains_title"]
        assert (scope["input"]["slide_text"] in content) is expected["markdown_contains_slide_text"]
        assert all(name in content for name in image_suggested_names) is expected["markdown_contains_two_image_links"]
        first_index = content.index(image_suggested_names[0])
        second_index = content.index(image_suggested_names[1])
        assert (first_index < second_index) is expected["image_links_in_order"]
        for token in expected["required_markdown_tokens"]:
            assert token in content

    def test_pptx_ocr_image_md_creates_sidecar(
        self,
        pipeline,
        sample_pptx_with_image,
        tmp_path,
        monkeypatch,
    ) -> None:
        """PPTX image_md OCR should finalize primary Markdown, image, and sidecar."""
        from docwen_core.text.ocr import OcrOutcome, OcrStatus

        monkeypatch.setattr(
            "docwen_core.text.ocr.run_ocr_outcome",
            lambda _path, **_kwargs: OcrOutcome(
                OcrStatus.SUCCESS,
                text="OCR text from PPTX image",
            ),
        )
        _plugin, task_mgr, ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_ocr_sidecar"
        output_dir.mkdir()
        result = _run_request(
            task_mgr,
            sample_pptx_with_image,
            "pptx",
            output_dir,
            to_md_enable_ocr=True,
            ocr_placement="image_md",
            image_link_style="markdown_embed",
        )

        assert result.success
        assert any(d.code == "FINALIZER_DONE" for d in result.diagnostics)
        markdown_artifact = next(artifact for artifact in result.artifacts if artifact.is_primary)
        image_artifacts = [artifact for artifact in result.artifacts if artifact.kind == "image"]
        sidecar_artifacts = [artifact for artifact in result.artifacts if artifact.kind == "auxiliary"]
        assert len(image_artifacts) == 1
        assert len(sidecar_artifacts) == 1

        markdown_path = Path(markdown_artifact.staging_path)
        image_path = Path(image_artifacts[0].staging_path)
        sidecar_path = Path(sidecar_artifacts[0].staging_path)
        node_root = _document_node_root(markdown_path, output_dir)
        assert _document_node_root(image_path, output_dir) == node_root
        assert _document_node_root(sidecar_path, output_dir) == node_root
        assert markdown_path.exists()
        assert image_path.exists()
        assert _io_path(sidecar_path).exists()

        content = markdown_path.read_text(encoding="utf-8")
        sidecar_content = _io_path(sidecar_path).read_text(encoding="utf-8")
        assert sidecar_artifacts[0].logical_path is not None
        assert sidecar_artifacts[0].logical_path.split("/", 1)[1] in content
        assert image_artifacts[0].suggested_name not in content
        assert image_artifacts[0].suggested_name in sidecar_content
        assert "OCR text from PPTX image" in sidecar_content
        assert str(Path(ws_mgr.root_dir)) not in content
        assert str(Path(ws_mgr.root_dir)) not in sidecar_content

    def test_pptx_multi_ocr_sidecar_matches_old_system_projection(
        self,
        pipeline,
        tmp_path: Path,
        monkeypatch,
    ) -> None:
        """Focused PPTX multi-image OCR sidecar projection matches old systems."""
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        fixture = _load_pptx_old_system_fixture()
        scope = fixture["image_artifact_scope"]["multi_ocr_sidecar_probe"]

        red_path = tmp_path / "red.png"
        blue_path = tmp_path / "blue.png"
        Image.new("RGB", tuple(scope["input"]["image_size_px"]), (255, 0, 0)).save(red_path)
        Image.new("RGB", tuple(scope["input"]["image_size_px"]), (0, 0, 255)).save(blue_path)

        pres = Presentation()
        pres.core_properties.title = scope["input"]["title"]
        slide = pres.slides.add_slide(pres.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(0.4))
        text_box.text_frame.text = scope["input"]["slide_text"]
        slide.shapes.add_picture(str(red_path), Inches(1), Inches(1), width=Inches(1), height=Inches(1))
        slide.shapes.add_picture(str(blue_path), Inches(3), Inches(1), width=Inches(1), height=Inches(1))

        input_path = tmp_path / scope["input"]["filename"]
        pres.save(str(input_path))

        from docwen_core.text.ocr import OcrOutcome, OcrStatus

        monkeypatch.setattr(
            "docwen_core.text.ocr.run_ocr_outcome",
            lambda _path, **_kwargs: OcrOutcome(
                OcrStatus.SUCCESS,
                text=scope["input"]["stubbed_ocr_text"],
            ),
        )

        _plugin, task_mgr, ws_mgr = pipeline
        output_dir = tmp_path / "output_pptx_multi_ocr"
        output_dir.mkdir()
        result = _run_request(
            task_mgr,
            input_path,
            "pptx",
            output_dir,
            to_md_enable_ocr=True,
            ocr_placement=scope["input"]["options"]["ocr_placement"],
            image_link_style=scope["input"]["options"]["image_link_style"],
        )

        expected = scope["current_projection"]
        assert result.success is expected["success"]
        assert [diagnostic.code for diagnostic in result.diagnostics] == expected["diagnostic_codes"]
        assert len(_deliverable_artifacts(result)) == expected["artifact_count"]
        primary_artifacts = [artifact for artifact in result.artifacts if artifact.is_primary]
        image_artifacts = [artifact for artifact in result.artifacts if artifact.kind == "image"]
        sidecar_artifacts = [artifact for artifact in result.artifacts if artifact.kind == "auxiliary"]
        assert len(primary_artifacts) == expected["primary_artifact_count"]
        assert len(image_artifacts) == expected["image_artifact_count"]
        assert len(sidecar_artifacts) == expected["auxiliary_artifact_count"]
        image_suggested_names = [artifact.suggested_name for artifact in image_artifacts]
        sidecar_suggested_names = [artifact.metadata["source_suggested_name"] for artifact in sidecar_artifacts]
        assert all(
            re.fullmatch(pattern, suggested_name)
            for pattern, suggested_name in zip(
                expected["image_suggested_name_patterns"],
                image_suggested_names,
                strict=True,
            )
        )
        assert all(
            re.fullmatch(pattern, suggested_name)
            for pattern, suggested_name in zip(
                expected["sidecar_suggested_name_patterns"],
                sidecar_suggested_names,
                strict=True,
            )
        )

        markdown_artifact = primary_artifacts[0]
        for key, value in expected["primary_metadata"].items():
            assert markdown_artifact.metadata[key] == value
        node_root = _document_node_root(Path(markdown_artifact.staging_path), output_dir)
        assert Path(markdown_artifact.staging_path).is_file()
        for artifact in image_artifacts + sidecar_artifacts:
            assert _document_node_root(Path(artifact.staging_path), output_dir) == node_root
            assert _io_path(artifact.staging_path).is_file()

        content = _legacy_markdown_projection(
            Path(markdown_artifact.staging_path).read_text(encoding="utf-8"), sidecar_artifacts
        )
        sidecar_contents = [
            _legacy_markdown_projection(_io_path(artifact.staging_path).read_text(encoding="utf-8"), sidecar_artifacts)
            for artifact in sidecar_artifacts
        ]
        assert str(output_dir) not in content
        assert str(Path(ws_mgr.root_dir)) not in content
        assert all(str(output_dir) not in sidecar_content for sidecar_content in sidecar_contents)
        assert all(str(Path(ws_mgr.root_dir)) not in sidecar_content for sidecar_content in sidecar_contents)
        assert (scope["input"]["stubbed_ocr_text"] in content) is expected["primary_markdown_contains_ocr_text"]
        assert (
            all(name in content for name in sidecar_suggested_names)
            is expected["primary_markdown_contains_sidecar_links"]
        )
        first_index = content.index(sidecar_suggested_names[0])
        second_index = content.index(sidecar_suggested_names[1])
        assert (first_index < second_index) is expected["sidecar_links_in_order"]
        assert (
            all(scope["input"]["stubbed_ocr_text"] in text for text in sidecar_contents)
            is expected["sidecars_contain_ocr_text"]
        )
        assert (
            all("![slide1_img" in text and ".png)" in text for text in sidecar_contents)
            is expected["sidecars_contain_image_markdown"]
        )
        for token in expected["required_primary_markdown_tokens"]:
            assert token in content
        for sidecar_content in sidecar_contents:
            for token in expected["required_sidecar_markdown_tokens"]:
                assert token in sidecar_content
