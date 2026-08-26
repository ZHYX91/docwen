"""Focused tests split from test_presentation_to_md.py."""

from __future__ import annotations

import pytest

from ._presentation_to_md_support import (
    Path,
    _data_uri_payload,
    _document_node_root,
    _load_pptx_old_system_fixture,
    _run_request,
)
from ._presentation_to_md_support import (
    pipeline as pipeline,
)

pytestmark = [pytest.mark.golden, pytest.mark.contract]
from ._presentation_to_md_support import (
    sample_pptx_file as sample_pptx_file,
)
from ._presentation_to_md_support import (
    sample_pptx_with_image as sample_pptx_with_image,
)
from ._presentation_to_md_support import (
    sample_pptx_with_notes as sample_pptx_with_notes,
)
from ._presentation_to_md_support import (
    sample_pptx_with_sections as sample_pptx_with_sections,
)


class TestPptxToMd:
    """Golden parity tests for ROUTE-PPTX-001: pptx → md."""

    def test_pptx_to_md_matches_old_system_semantic_fixture(self, pipeline, tmp_path) -> None:
        """Current PPTX→MD should preserve old-system core slide semantics."""
        from pptx import Presentation
        from pptx.shapes.placeholder import SlidePlaceholder
        from pptx.util import Inches

        fixture = _load_pptx_old_system_fixture()
        pptx_spec = fixture["input_pptx"]
        pres = Presentation()
        pres.core_properties.title = pptx_spec["title"]

        title_slide = pptx_spec["slides"][0]
        slide1 = pres.slides.add_slide(pres.slide_layouts[0])
        title1 = slide1.shapes.title
        subtitle1 = slide1.placeholders[1]
        assert isinstance(title1, SlidePlaceholder)
        assert isinstance(subtitle1, SlidePlaceholder)
        title1.text = title_slide["title"]
        subtitle1.text = title_slide["subtitle"]

        agenda_slide = pptx_spec["slides"][1]
        slide2 = pres.slides.add_slide(pres.slide_layouts[1])
        title2 = slide2.shapes.title
        body = slide2.shapes.placeholders[1]
        assert isinstance(title2, SlidePlaceholder)
        assert isinstance(body, SlidePlaceholder)
        title2.text = agenda_slide["title"]
        tf = body.text_frame
        tf.text = agenda_slide["paragraphs"][0]
        for paragraph in agenda_slide["paragraphs"][1:]:
            p = tf.add_paragraph()
            p.text = paragraph

        table_slide = pptx_spec["slides"][2]
        slide3 = pres.slides.add_slide(pres.slide_layouts[1])
        title3 = slide3.shapes.title
        assert isinstance(title3, SlidePlaceholder)
        title3.text = table_slide["title"]
        table_data = table_slide["table"]
        rows = len(table_data["rows"]) + 1
        cols = len(table_data["headers"])
        table_shape = slide3.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(3))
        table = table_shape.table
        for col_idx, header in enumerate(table_data["headers"]):
            table.cell(0, col_idx).text = header
        for row_idx, row_data in enumerate(table_data["rows"], start=1):
            for col_idx, value in enumerate(row_data):
                table.cell(row_idx, col_idx).text = value

        pptx_path = tmp_path / pptx_spec["filename"]
        pres.save(str(pptx_path))
        output_dir = tmp_path / "output_pptx_old_system_fixture"
        output_dir.mkdir()
        _plugin, task_mgr, ws_mgr = pipeline

        result = _run_request(task_mgr, pptx_path, "pptx", output_dir)

        assert result.success, f"PPTX conversion failed: {result.error.message if result.error else 'unknown'}"
        artifact = result.artifacts[0]
        expected_metadata = fixture["current_expected_semantics"]["artifact_metadata"]
        assert artifact.media_type == expected_metadata["media_type"]
        assert artifact.metadata["title"] == expected_metadata["title"]
        assert artifact.metadata["slide_count"] == expected_metadata["slide_count"]
        assert artifact.metadata["table_count"] == expected_metadata["table_count"]
        assert artifact.metadata["image_count"] == expected_metadata["image_count"]
        assert any(d.code == "FINALIZER_DONE" for d in result.diagnostics)

        artifact_path = Path(artifact.staging_path)
        node_root = _document_node_root(artifact_path, output_dir)
        assert artifact_path.name == f"{node_root.name}.md"
        assert artifact_path.exists()

        content = artifact_path.read_text(encoding="utf-8")
        assert str(Path(ws_mgr.root_dir)) not in content
        for token in fixture["current_expected_semantics"]["required_markdown_tokens"]:
            assert token in content

        non_empty_lines = [line.strip() for line in content.splitlines() if line.strip()]
        assert non_empty_lines[0] == fixture["probe_outputs"]["current_first_non_empty_line"]

    def test_pptx_conversion_succeeds(self, pipeline, sample_pptx_file, tmp_path) -> None:
        """PPTX file must convert to Markdown successfully."""
        _plugin, task_mgr, _ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx"
        output_dir.mkdir()
        result = _run_request(task_mgr, sample_pptx_file, "pptx", output_dir)

        assert result.success, f"PPTX conversion failed: {result.error.message if result.error else 'unknown'}"
        assert len(result.artifacts) >= 1
        assert result.artifacts[0].is_primary is True

        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")
        assert len(content.strip()) > 0

    def test_pptx_title_extraction(self, pipeline, sample_pptx_file, tmp_path) -> None:
        """Presentation title must appear as heading."""
        _plugin, task_mgr, _ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_title"
        output_dir.mkdir()
        result = _run_request(task_mgr, sample_pptx_file, "pptx", output_dir)

        assert result.success
        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")
        assert "Test Presentation" in content
        assert "# " in content

    def test_pptx_yaml_frontmatter_consumes_locale_title_label(self, pipeline, sample_pptx_file, tmp_path) -> None:
        """PPTX→Markdown consumes pre-resolved YAML labels from the app edge."""
        _plugin, task_mgr, ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_locale_yaml"
        output_dir.mkdir()
        result = _run_request(
            task_mgr,
            sample_pptx_file,
            "pptx",
            output_dir,
            yaml_key_labels={"title": "Titel"},
        )

        assert result.success
        assert any(d.code == "FINALIZER_DONE" for d in result.diagnostics)
        artifact_path = Path(result.artifacts[0].staging_path)
        node_root = _document_node_root(artifact_path, output_dir)
        assert artifact_path.name == f"{node_root.name}.md"
        assert artifact_path.exists()
        content = artifact_path.read_text(encoding="utf-8")
        assert str(Path(ws_mgr.root_dir)) not in content
        yaml_block = content.split("---", 2)[1]
        assert "Titel: Test Presentation" in yaml_block
        assert "标题: Test Presentation" not in yaml_block
        assert "title: Test Presentation" not in yaml_block

    def test_pptx_slide_content(self, pipeline, sample_pptx_file, tmp_path) -> None:
        """Slide headings and text must be preserved."""
        _plugin, task_mgr, _ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_slides"
        output_dir.mkdir()
        result = _run_request(task_mgr, sample_pptx_file, "pptx", output_dir)

        assert result.success
        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")
        assert "## " in content, f"No slide headings (##) found.\nContent:\n{content[:300]}"
        assert "Agenda" in content
        assert "Data Summary" in content
        assert "First item" in content

    def test_pptx_sections_match_old_system_heading_semantics(
        self,
        pipeline,
        sample_pptx_with_sections,
        tmp_path,
    ) -> None:
        """PowerPoint sections must emit section H1 headings before their slides."""
        _plugin, task_mgr, _ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_sections"
        output_dir.mkdir()
        result = _run_request(task_mgr, sample_pptx_with_sections, "pptx", output_dir)

        assert result.success
        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")
        assert "# First Section" in content
        assert "# Second Section" in content
        assert "## Slide One" in content
        assert "## Slide Two" in content
        assert content.index("# First Section") < content.index("## Slide One")
        assert content.index("# Second Section") < content.index("## Slide Two")

    def test_pptx_export_notes_matches_old_system_blockquote_semantics(
        self,
        pipeline,
        sample_pptx_with_notes,
        tmp_path,
    ) -> None:
        """export_notes=True should include speaker notes as blockquotes."""
        _plugin, task_mgr, _ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_notes"
        output_dir.mkdir()
        result = _run_request(task_mgr, sample_pptx_with_notes, "pptx", output_dir, export_notes=True)

        assert result.success
        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")
        assert "> Notes:" in content
        assert "> Speaker note line" in content

    def test_pptx_table_preserved(self, pipeline, sample_pptx_file, tmp_path) -> None:
        """Tables in PPTX must be converted to Markdown tables."""
        _plugin, task_mgr, _ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_tables"
        output_dir.mkdir()
        result = _run_request(task_mgr, sample_pptx_file, "pptx", output_dir)

        assert result.success
        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")

        assert "Category" in content
        assert "Q1" in content
        assert "Q2" in content
        assert "Alpha" in content
        assert "Beta" in content
        assert "100" in content
        assert "200" in content
        assert "|" in content

    def test_pptx_artifact_metadata(self, pipeline, sample_pptx_file, tmp_path) -> None:
        """Artifact must carry correct metadata."""
        _plugin, task_mgr, _ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_meta"
        output_dir.mkdir()
        result = _run_request(task_mgr, sample_pptx_file, "pptx", output_dir)

        assert result.success
        artifact = result.artifacts[0]
        assert artifact.kind == "primary"
        assert artifact.suggested_name.endswith(".md")
        assert artifact.media_type == "text/markdown"
        assert artifact.is_primary is True

        meta = artifact.metadata
        assert meta.get("slide_count", 0) >= 2
        assert meta.get("table_count", 0) >= 1

    def test_pptx_image_link_style_option(self, pipeline, sample_pptx_with_image, tmp_path) -> None:
        """PPTX image links must honor the GUI/request link style option."""
        _plugin, task_mgr, _ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_image_link"
        output_dir.mkdir()
        result = _run_request(
            task_mgr,
            sample_pptx_with_image,
            "pptx",
            output_dir,
            image_link_style="markdown_link",
        )

        assert result.success
        image_artifact = next(artifact for artifact in result.artifacts if artifact.kind == "image")
        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")
        assert f"[slide1_img1.png]({image_artifact.suggested_name})" in content
        assert f"![[{image_artifact.suggested_name}]]" not in content

    def test_pptx_image_mode_base64_inlines_images_without_image_artifacts(
        self,
        pipeline,
        sample_pptx_with_image,
        tmp_path,
    ) -> None:
        """PPTX image_mode=base64 must preserve old-system inline image semantics."""
        _plugin, task_mgr, _ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_base64_image_mode"
        output_dir.mkdir()
        result = _run_request(
            task_mgr,
            sample_pptx_with_image,
            "pptx",
            output_dir,
            image_mode="base64",
            image_link_style="markdown_embed",
        )

        assert result.success
        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")
        assert "![slide1_img1.png](data:image/png;base64," in content
        assert "slide1_img1.png](slide1_img1.png)" not in content
        assert [artifact for artifact in result.artifacts if artifact.kind == "image"] == []

    def test_pptx_base64_consumes_compression_semantics(
        self,
        pipeline,
        tmp_path: Path,
    ) -> None:
        """Large embedded images consume the shared Base64 compression policy."""
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        image_path = tmp_path / "noise.png"
        Image.effect_noise((512, 512), 100).convert("RGB").save(image_path, format="PNG")
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
        input_path = tmp_path / "base64-compression.pptx"
        presentation.save(str(input_path))

        _plugin, task_mgr, _ws_mgr = pipeline
        output_dir = tmp_path / "output_pptx_base64_compression"
        output_dir.mkdir()
        result = _run_request(
            task_mgr,
            input_path,
            "pptx",
            output_dir,
            config_snapshot={
                "conversion": {
                    "export": {
                        "base64_compress_enabled": True,
                        "base64_compress_threshold_kb": 100,
                    }
                }
            },
            image_mode="base64",
            image_link_style="markdown_embed",
        )

        assert result.success
        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")
        mime, payload = _data_uri_payload(content)
        assert mime == "image/jpeg"
        assert payload.startswith(b"\xff\xd8\xff")
        assert len(payload) <= 100 * 1024
        assert len(payload) < image_path.stat().st_size
        assert [artifact for artifact in result.artifacts if artifact.kind == "image"] == []

    def test_pptx_base64_snapshot_controls_compression_semantics(
        self,
        pipeline,
        tmp_path: Path,
    ) -> None:
        """An admitted request can disable compression explicitly."""
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        image_path = tmp_path / "snapshot-noise.png"
        Image.effect_noise((512, 512), 100).convert("RGB").save(image_path, format="PNG")
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
        input_path = tmp_path / "snapshot-base64-compression.pptx"
        presentation.save(str(input_path))

        output_dir = tmp_path / "output_pptx_snapshot_base64"
        output_dir.mkdir()
        result = _run_request(
            pipeline[1],
            input_path,
            "pptx",
            output_dir,
            config_snapshot={
                "conversion": {
                    "export": {
                        "base64_compress_enabled": False,
                        "base64_compress_threshold_kb": 1,
                    }
                },
                "export": {
                    "to_md_image_extraction_mode": "base64",
                    "to_md_ocr_placement_mode": "main_md",
                },
            },
            image_mode="base64",
            image_link_style="markdown_embed",
        )

        assert result.success
        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")
        mime, payload = _data_uri_payload(content)
        assert mime == "image/png"
        assert payload == image_path.read_bytes()
        assert [artifact for artifact in result.artifacts if artifact.kind == "image"] == []

    def test_pptx_base64_ocr_request_forces_main_md_without_sidecar(
        self,
        pipeline,
        sample_pptx_with_image,
        tmp_path,
        monkeypatch,
    ) -> None:
        """base64 image mode follows old-system semantics and writes OCR inline."""
        from docwen_core.text.ocr import OcrOutcome, OcrStatus

        monkeypatch.setattr(
            "docwen_core.text.ocr.run_ocr_outcome",
            lambda _path, **_kwargs: OcrOutcome(OcrStatus.SUCCESS, text="PPTX OCR inline"),
        )
        _plugin, task_mgr, _ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_base64_ocr_main"
        output_dir.mkdir()
        result = _run_request(
            task_mgr,
            sample_pptx_with_image,
            "pptx",
            output_dir,
            image_mode="base64",
            to_md_enable_ocr=True,
            ocr_placement="image_md",
            image_link_style="markdown_embed",
        )

        assert result.success
        markdown_artifact = next(artifact for artifact in result.artifacts if artifact.is_primary)
        content = Path(markdown_artifact.staging_path).read_text(encoding="utf-8")
        assert "![slide1_img1.png](data:image/png;base64," in content
        assert "> PPTX OCR inline" in content
        assert [artifact for artifact in result.artifacts if artifact.kind == "image"] == []
        assert [artifact for artifact in result.artifacts if artifact.kind == "auxiliary"] == []

    def test_pptx_images_are_finalized_as_relative_markdown_artifacts(
        self,
        pipeline,
        sample_pptx_with_image,
        tmp_path,
    ) -> None:
        """PPTX images must survive runtime finalization, not only staging."""
        _plugin, task_mgr, ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_images"
        output_dir.mkdir()
        result = _run_request(
            task_mgr,
            sample_pptx_with_image,
            "pptx",
            output_dir,
            image_link_style="markdown_embed",
        )

        assert result.success
        assert any(d.code == "FINALIZER_DONE" for d in result.diagnostics)
        markdown_artifact = next(artifact for artifact in result.artifacts if artifact.is_primary)
        markdown_path = Path(markdown_artifact.staging_path)
        node_root = _document_node_root(markdown_path, output_dir)
        assert markdown_path.is_file()

        image_artifacts = [artifact for artifact in result.artifacts if artifact.kind == "image"]
        assert markdown_artifact.metadata["image_count"] == 1
        assert len(image_artifacts) == 1
        image_path = Path(image_artifacts[0].staging_path)
        assert _document_node_root(image_path, output_dir) == node_root
        assert image_path.is_file()
        assert image_path.read_bytes().startswith(b"\x89PNG\r\n\x1a\n")

        content = markdown_path.read_text(encoding="utf-8")
        assert image_artifacts[0].suggested_name in content
        assert str(output_dir) not in content
        assert str(Path(ws_mgr.root_dir)) not in content
