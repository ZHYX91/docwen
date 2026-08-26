"""Focused tests split from test_presentation_to_md.py."""

from __future__ import annotations

from ._presentation_to_md_support import (
    Any,
    ConversionRequest,
    FileRef,
    OutputPolicy,
    Path,
    _run_request,
    pytest,
)
from ._presentation_to_md_support import (
    pipeline as pipeline,
)

pytestmark = [pytest.mark.golden, pytest.mark.contract]
from ._presentation_to_md_support import (
    sample_pptx_file as sample_pptx_file,
)
from ._presentation_to_md_support import (
    sample_pptx_with_image as sample_pptx_with_image,
)
from ._presentation_to_md_support import (
    sample_pptx_with_notes as sample_pptx_with_notes,
)
from ._presentation_to_md_support import (
    sample_pptx_with_sections as sample_pptx_with_sections,
)


class TestPptxToMd:
    """Golden parity tests for ROUTE-PPTX-001: pptx → md."""

    def test_pptx_all_ocr_outcomes_warn_and_continue_later_images(
        self,
        pipeline,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """Every OCR outcome warns while base slide content and later OCR survive."""
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        import docwen_core.text.ocr as ocr

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(6), Inches(0.4))
        text_box.text_frame.text = "presentation base content"
        for index in range(6):
            image_path = tmp_path / f"presentation-ocr-{index}.png"
            Image.new("RGB", (2, 2), color=(index * 20, 0, 0)).save(image_path)
            slide.shapes.add_picture(
                str(image_path),
                Inches(0.5 + index),
                Inches(1),
                width=Inches(0.5),
                height=Inches(0.5),
            )
        input_path = tmp_path / "presentation-ocr-best-effort.pptx"
        presentation.save(str(input_path))

        outcomes = iter(
            [
                ocr.OcrOutcome(ocr.OcrStatus.UNAVAILABLE, message="private unavailable detail"),
                ocr.OcrOutcome(ocr.OcrStatus.MODEL_MISSING, message="private model path"),
                ocr.OcrOutcome(ocr.OcrStatus.INITIALIZATION_FAILED, message="private init detail"),
                ocr.OcrOutcome(ocr.OcrStatus.RECOGNITION_FAILED, message="private recognition detail"),
                ocr.OcrOutcome(ocr.OcrStatus.NO_TEXT),
                ocr.OcrOutcome(ocr.OcrStatus.SUCCESS, text="later presentation OCR"),
            ]
        )
        calls: list[str] = []

        def _run_ocr_outcome(path: str, **_kwargs: Any) -> ocr.OcrOutcome:
            calls.append(str(path))
            return next(outcomes)

        monkeypatch.setattr(ocr, "run_ocr_outcome", _run_ocr_outcome)

        _plugin, task_mgr, _ws_mgr = pipeline
        output_dir = tmp_path / "presentation-ocr-best-effort-output"
        output_dir.mkdir()
        events: list[Any] = []

        result = _run_request(
            task_mgr,
            input_path,
            "pptx",
            output_dir,
            on_event=events.append,
            to_md_enable_ocr=True,
            ocr_placement="main_md",
            image_mode="file",
        )

        assert result.success is True
        markdown = Path(next(artifact for artifact in result.artifacts if artifact.is_primary).staging_path).read_text(
            encoding="utf-8"
        )
        assert "presentation base content" in markdown
        assert "later presentation OCR" in markdown
        assert len(calls) == 6
        warnings = [
            event.payload
            for event in events
            if event.event_type == "diagnostic" and event.payload.get("code") == "OCR-BEST-EFFORT"
        ]
        assert len(warnings) == 6
        assert [warning["level"] for warning in warnings] == ["warning"] * 6
        assert [warning["message"].split("status=", 1)[1].split(";", 1)[0] for warning in warnings] == [
            "unavailable",
            "model_missing",
            "initialization_failed",
            "recognition_failed",
            "no_text",
            "success",
        ]
        assert all(warning["location"].startswith("slide 1:") for warning in warnings)
        assert all("private" not in warning["message"] for warning in warnings)

    def test_pptx_image_link_style_uses_request_export_semantics(
        self,
        pipeline,
        sample_pptx_with_image,
        tmp_path,
    ) -> None:
        """PPTX image links honor the admitted Link setting when options omit it."""

        _plugin, task_mgr, _ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_image_link_default"
        output_dir.mkdir()
        result = _run_request(
            task_mgr,
            sample_pptx_with_image,
            "pptx",
            output_dir,
            config_snapshot={"link": {"format": {"image_link_style": "markdown_embed"}}},
        )

        assert result.success
        image_artifact = next(artifact for artifact in result.artifacts if artifact.kind == "image")
        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")
        assert f"![slide1_img1.png]({image_artifact.suggested_name})" in content
        assert f"![[{image_artifact.suggested_name}]]" not in content

    def test_pptx_nonempty_snapshot_owns_export_semantics(
        self,
        pipeline,
        sample_pptx_with_image,
        tmp_path,
    ) -> None:
        """An admitted snapshot owns the complete export policy."""
        _plugin, task_mgr, _ws_mgr = pipeline
        output_dir = tmp_path / "output_pptx_snapshot_link"
        output_dir.mkdir()
        result = _run_request(
            task_mgr,
            sample_pptx_with_image,
            "pptx",
            output_dir,
            config_snapshot={
                "link": {
                    "format": {
                        "image_link_style": "wiki_embed",
                        "md_file_link_style": "wiki_link",
                    }
                },
                "export": {
                    "to_md_image_extraction_mode": "file",
                    "to_md_ocr_placement_mode": "main_md",
                },
            },
        )

        assert result.success
        image_artifact = next(artifact for artifact in result.artifacts if artifact.kind == "image")
        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")
        assert f"![[{image_artifact.suggested_name}]]" in content
        assert f"![slide1_img1.png]({image_artifact.suggested_name})" not in content

    def test_pptx_ocr_title_is_frozen_from_nonempty_snapshot(
        self,
        pipeline,
        sample_pptx_with_image,
        tmp_path,
        monkeypatch,
    ) -> None:
        """OCR presentation must use the admitted request generation."""
        from docwen_core.text.ocr import OcrOutcome, OcrStatus

        monkeypatch.setattr(
            "docwen_core.text.ocr.run_ocr_outcome",
            lambda _path, **_kwargs: OcrOutcome(OcrStatus.SUCCESS, text="snapshot OCR body"),
        )
        _plugin, task_mgr, _ws_mgr = pipeline
        output_dir = tmp_path / "output_pptx_snapshot_ocr_title"
        output_dir.mkdir()
        result = _run_request(
            task_mgr,
            sample_pptx_with_image,
            "pptx",
            output_dir,
            config_snapshot={
                "conversion": {
                    "ocr_output": {
                        "show_blockquote_title": True,
                        "blockquote_title_override_by_locale": {"en_US": "Snapshot OCR title"},
                    }
                },
                "export": {
                    "to_md_image_extraction_mode": "file",
                    "to_md_ocr_placement_mode": "main_md",
                },
                "gui": {"language": {"locale": "en_US"}},
            },
            to_md_enable_ocr=True,
        )

        assert result.success
        content = Path(result.artifacts[0].staging_path).read_text(encoding="utf-8")
        assert "> **Snapshot OCR title**" in content
        assert "Global OCR title" not in content

    def test_pptx_builds_one_request_policy_for_every_slide(
        self,
        pipeline,
        sample_pptx_file,
        tmp_path,
        monkeypatch,
    ) -> None:
        """One immutable policy must be shared by every slide in a request."""
        from docwen_plugin_presentation.pptx_md import converter as converter_module

        built_policies: list[Any] = []
        slide_policies: list[Any] = []
        original_builder = converter_module.build_presentation_markdown_request_policy
        original_process_slide = converter_module.PptxToMarkdownConverter._process_slide

        def record_builder(context: Any, options: Any) -> Any:
            policy = original_builder(context, options)
            built_policies.append(policy)
            return policy

        def record_process_slide(
            self: Any,
            *args: Any,
            request_policy: Any = None,
            **kwargs: Any,
        ) -> Any:
            slide_policies.append(request_policy)
            return original_process_slide(
                self,
                *args,
                request_policy=request_policy,
                **kwargs,
            )

        monkeypatch.setattr(
            converter_module,
            "build_presentation_markdown_request_policy",
            record_builder,
        )
        monkeypatch.setattr(
            converter_module.PptxToMarkdownConverter,
            "_process_slide",
            record_process_slide,
        )
        output_dir = tmp_path / "output_pptx_single_policy"
        output_dir.mkdir()

        result = _run_request(
            pipeline[1],
            sample_pptx_file,
            "pptx",
            output_dir,
            config_snapshot={
                "export": {
                    "to_md_image_extraction_mode": "file",
                    "to_md_ocr_placement_mode": "main_md",
                }
            },
        )

        assert result.success
        assert len(built_policies) == 1
        assert len(slide_policies) >= 2
        assert all(policy is built_policies[0] for policy in slide_policies)

    def test_pptx_cancellation(self, pipeline, sample_pptx_file, tmp_path) -> None:
        """PPTX→MD must support cancellation."""
        _plugin, task_mgr, _ws_mgr = pipeline

        output_dir = tmp_path / "output_pptx_cancel"
        output_dir.mkdir()

        request = ConversionRequest(
            request_id="pptx-cancel-test",
            input_refs=[
                FileRef(
                    path=str(sample_pptx_file),
                    format="pptx",
                    category="presentation",
                    size_bytes=sample_pptx_file.stat().st_size,
                )
            ],
            target_format="md",
            output_policy=OutputPolicy(output_dir=str(output_dir)),
        )

        task_mgr.cancel("pptx-cancel-test")
        result = task_mgr.execute_single(request)

        assert result.success is False
        assert result.error is not None
        assert result.error.error_type == "cancelled"
