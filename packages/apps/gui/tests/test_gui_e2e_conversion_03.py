"""Focused tests split from test_gui_e2e_conversion.py."""

from __future__ import annotations

from ._gui_e2e_conversion_support import (
    _E2E_CONVERSION_TIMEOUT_MS,
    Path,
    _assert_markdown_node,
    _create_ocr_smoke_png,
    _wait_for,
    pytest,
)


@pytest.mark.gui
class TestImageGuiExecution:
    def test_png_to_markdown_real_ocr_runs_through_gui_action_route(
        self,
        main_window_with_controller,
        tmp_path: Path,
    ) -> None:
        from PySide6.QtWidgets import QApplication

        from docwen_gui.main_window import _normalize_path

        source = _create_ocr_smoke_png(tmp_path / "hello_docwen_ocr.png")

        window = main_window_with_controller
        controller = window.view_model.controller
        assert controller is not None
        assert controller.config_port is not None
        output_dir = tmp_path / "image_ocr_gui_exports"
        assert controller.config_port.set("output.directory.mode", "custom")
        assert controller.config_port.set("output.directory.custom_path", str(output_dir))
        assert controller.config_port.set("output.directory.create_date_subfolder", False)
        # The product default is ``main_md``.  This smoke intentionally exercises
        # the still-supported sidecar route, so keep that choice explicit.
        assert controller.config_port.set("export.to_md_ocr_placement_mode", "image_md")

        normalized = _normalize_path(str(source))
        window.view_model.add_files([str(source)])

        app = QApplication.instance()
        if app is not None:
            app.processEvents()

        assert _wait_for(lambda: window._action_area_vm.visible, timeout_ms=5000, interval_ms=20)
        assert window._action_area_vm.file_type == "image"
        assert window._action_area_vm.extract_ocr is True
        assert window._action_area_vm.collect_options()["ocr_placement"] == "image_md"
        window._action_area_vm.request_conversion("md")

        if app is not None:
            app.processEvents()

        def _conversion_finished() -> bool:
            entry = window._batch_list_vm.get_file_entry(normalized)
            return entry is not None and entry.status in ("completed", "failed")

        assert _wait_for(_conversion_finished, timeout_ms=_E2E_CONVERSION_TIMEOUT_MS)
        entry = window._batch_list_vm.get_file_entry(normalized)
        assert entry is not None
        assert entry.status == "completed"
        assert entry.output_path

        output_path = Path(entry.output_path)
        image_path = output_path.with_name("hello_docwen_ocr.png")
        _assert_markdown_node(
            output_path,
            source_stem="hello_docwen_ocr",
            source_tag="Png",
            output_root=output_dir,
        )
        sidecars = [path for path in output_path.parent.rglob("*.md") if path != output_path]
        assert len(sidecars) == 1
        sidecar_path = sidecars[0]
        assert sidecar_path.parent.name == sidecar_path.stem
        assert image_path.exists()

        primary_text = output_path.read_text(encoding="utf-8")
        sidecar_text = sidecar_path.read_text(encoding="utf-8")
        assert sidecar_path.name in primary_text
        assert "HELLO DOCWEN OCR" in sidecar_text
        assert "> HELLO DOCWEN OCR" in sidecar_text
        assert window._info_area_vm.has_task_summary
        assert window._info_area_vm._task_summary.state == "success"


@pytest.mark.gui
class TestPresentationGuiExecution:
    def test_pptx_to_markdown_runs_through_gui_action_only_route(
        self,
        main_window_with_controller,
        tmp_path: Path,
    ) -> None:
        from pptx import Presentation
        from pptx.shapes.placeholder import SlidePlaceholder
        from PySide6.QtWidgets import QApplication

        from docwen_gui.main_window import _normalize_path

        source = tmp_path / "presentation-smoke.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title_placeholder = slide.shapes.title
        body_placeholder = slide.placeholders[1]
        assert isinstance(title_placeholder, SlidePlaceholder)
        assert isinstance(body_placeholder, SlidePlaceholder)
        title_placeholder.text = "Presentation Smoke"
        body_placeholder.text = "GUI action-only route to Markdown"
        presentation.core_properties.title = "Presentation Smoke"
        presentation.save(str(source))

        window = main_window_with_controller
        normalized = _normalize_path(str(source))
        window.view_model.add_files([str(source)])

        app = QApplication.instance()
        if app is not None:
            app.processEvents()

        assert _wait_for(lambda: window._action_area_vm.visible, timeout_ms=5000, interval_ms=20)
        assert window._action_area_vm.file_type == "pptx"
        window._action_area_vm.request_conversion("md")

        if app is not None:
            app.processEvents()

        def _conversion_finished() -> bool:
            entry = window._batch_list_vm.get_file_entry(normalized)
            return entry is not None and entry.status in ("completed", "failed")

        assert _wait_for(_conversion_finished, timeout_ms=_E2E_CONVERSION_TIMEOUT_MS)
        entry = window._batch_list_vm.get_file_entry(normalized)
        assert entry is not None
        assert entry.status == "completed"
        assert entry.output_path

        output_path = Path(entry.output_path)
        _assert_markdown_node(output_path, source_stem="presentation-smoke", source_tag="Pptx")
        content = output_path.read_text(encoding="utf-8")
        assert "Presentation Smoke" in content
        assert "GUI action-only route to Markdown" in content
        assert window._info_area_vm.has_task_summary
        assert window._info_area_vm._task_summary.state == "success"


@pytest.mark.gui
class TestMarkupGuiExecution:
    def test_html_to_markdown_runs_through_gui_action_only_route(
        self,
        main_window_with_controller,
        tmp_path: Path,
    ) -> None:
        from PySide6.QtWidgets import QApplication

        from docwen_gui.main_window import _normalize_path

        source = tmp_path / "html-smoke.html"
        source.write_text(
            "<!doctype html><html><head><title>HTML Smoke</title></head>"
            "<body><h1>HTML Smoke</h1><p>GUI action-only route to Markdown.</p></body></html>",
            encoding="utf-8",
        )

        window = main_window_with_controller
        normalized = _normalize_path(str(source))
        window.view_model.add_files([str(source)])

        app = QApplication.instance()
        if app is not None:
            app.processEvents()

        assert _wait_for(lambda: window._action_area_vm.visible, timeout_ms=5000, interval_ms=20)
        assert window._action_area_vm.file_type == "html"
        window._action_area_vm.request_conversion("md")

        if app is not None:
            app.processEvents()

        def _conversion_finished() -> bool:
            entry = window._batch_list_vm.get_file_entry(normalized)
            return entry is not None and entry.status in ("completed", "failed")

        assert _wait_for(_conversion_finished, timeout_ms=_E2E_CONVERSION_TIMEOUT_MS)
        entry = window._batch_list_vm.get_file_entry(normalized)
        assert entry is not None
        assert entry.status == "completed"
        assert entry.output_path

        output_path = Path(entry.output_path)
        _assert_markdown_node(output_path, source_stem="html-smoke", source_tag="Html")
        content = output_path.read_text(encoding="utf-8")
        assert "HTML Smoke" in content
        assert "GUI action-only route to Markdown" in content
        assert window._info_area_vm.has_task_summary
        assert window._info_area_vm._task_summary.state == "success"

    def test_mhtml_to_markdown_runs_through_gui_action_only_route_with_embedded_resource(
        self,
        main_window_with_controller,
        tmp_path: Path,
    ) -> None:
        from email.mime.image import MIMEImage
        from email.mime.multipart import MIMEMultipart
        from email.mime.text import MIMEText

        from PySide6.QtWidgets import QApplication

        from docwen_gui.main_window import _normalize_path

        html_part = MIMEText(
            "<html><body><h1>MHTML Smoke</h1><p>GUI action-only route to Markdown.</p>"
            "<p><img src='cid:embedded-img'></p></body></html>",
            "html",
        )
        html_part["Content-Location"] = "mhtml-smoke.html"

        message = MIMEMultipart("related")
        message["Subject"] = "MHTML Smoke"
        message.attach(html_part)
        image_bytes = _create_ocr_smoke_png(tmp_path / "embedded-source.png").read_bytes()
        image_part = MIMEImage(image_bytes, _subtype="png")
        image_part["Content-ID"] = "<embedded-img>"
        image_part["Content-Location"] = "embedded.png"
        message.attach(image_part)

        source = tmp_path / "mhtml-smoke.mhtml"
        source.write_bytes(message.as_bytes())

        window = main_window_with_controller
        normalized = _normalize_path(str(source))
        window.view_model.add_files([str(source)])

        app = QApplication.instance()
        if app is not None:
            app.processEvents()

        assert _wait_for(lambda: window._action_area_vm.visible, timeout_ms=5000, interval_ms=20)
        assert window._action_area_vm.file_type == "mhtml"
        window._action_area_vm.request_conversion("md")

        if app is not None:
            app.processEvents()

        def _conversion_finished() -> bool:
            entry = window._batch_list_vm.get_file_entry(normalized)
            return entry is not None and entry.status in ("completed", "failed")

        assert _wait_for(_conversion_finished, timeout_ms=_E2E_CONVERSION_TIMEOUT_MS)
        entry = window._batch_list_vm.get_file_entry(normalized)
        assert entry is not None
        assert entry.status == "completed"
        assert entry.output_path

        output_path = Path(entry.output_path)
        embedded_path = output_path.with_name("embedded.png")
        _assert_markdown_node(output_path, source_stem="mhtml-smoke", source_tag="Mhtml")
        assert embedded_path.exists()
        assert embedded_path.read_bytes() == image_bytes
        content = output_path.read_text(encoding="utf-8")
        assert "MHTML Smoke" in content
        assert "GUI action-only route to Markdown" in content
        assert "embedded.png" in content
        assert window._info_area_vm.has_task_summary
        assert window._info_area_vm._task_summary.state == "success"

    def test_mhtml_to_markdown_ocr_runs_through_gui_action_only_route(
        self,
        main_window_with_controller,
        tmp_path: Path,
        monkeypatch,
    ) -> None:
        from email.mime.image import MIMEImage
        from email.mime.multipart import MIMEMultipart
        from email.mime.text import MIMEText

        from PySide6.QtWidgets import QApplication

        from docwen_core.text.ocr import OcrOutcome, OcrStatus
        from docwen_gui.main_window import _normalize_path

        monkeypatch.setattr(
            "docwen_plugin_markup.markdown_resources.run_ocr_outcome",
            lambda _path, **_kwargs: OcrOutcome(
                OcrStatus.SUCCESS,
                text="OCR text from GUI MHTML image",
            ),
        )

        html_part = MIMEText(
            "<html><body><h1>MHTML OCR Smoke</h1><p>GUI action-only OCR route.</p>"
            "<p><img src='cid:embedded-img'></p></body></html>",
            "html",
        )
        html_part["Content-Location"] = "mhtml-ocr-smoke.html"

        message = MIMEMultipart("related")
        message["Subject"] = "MHTML OCR Smoke"
        message.attach(html_part)
        image_bytes = _create_ocr_smoke_png(tmp_path / "mhtml-ocr-image.png").read_bytes()
        image_part = MIMEImage(image_bytes, _subtype="png")
        image_part["Content-ID"] = "<embedded-img>"
        image_part["Content-Location"] = "embedded.png"
        message.attach(image_part)

        source = tmp_path / "mhtml-ocr-smoke.mhtml"
        source.write_bytes(message.as_bytes())

        window = main_window_with_controller
        normalized = _normalize_path(str(source))
        window.view_model.add_files([str(source)])

        app = QApplication.instance()
        if app is not None:
            app.processEvents()

        assert _wait_for(lambda: window._action_area_vm.visible, timeout_ms=5000, interval_ms=20)
        assert window._action_area_vm.file_type == "mhtml"
        window._action_area_vm.set_file_to_md_option("extract_ocr", True)
        collected_options = window._action_area_vm.collect_options()
        assert collected_options["to_md_enable_ocr"] is True
        assert collected_options["ocr_placement"] == "main_md"
        window._action_area_vm.request_conversion("md")

        if app is not None:
            app.processEvents()

        def _conversion_finished() -> bool:
            entry = window._batch_list_vm.get_file_entry(normalized)
            return entry is not None and entry.status in ("completed", "failed")

        assert _wait_for(_conversion_finished, timeout_ms=_E2E_CONVERSION_TIMEOUT_MS)
        entry = window._batch_list_vm.get_file_entry(normalized)
        assert entry is not None
        assert entry.status == "completed"
        assert entry.output_path

        output_path = Path(entry.output_path)
        embedded_path = output_path.with_name("embedded.png")
        _assert_markdown_node(output_path, source_stem="mhtml-ocr-smoke", source_tag="Mhtml")
        assert embedded_path.exists()
        assert embedded_path.read_bytes() == image_bytes
        assert list(output_path.parent.rglob("*.md")) == [output_path]
        content = output_path.read_text(encoding="utf-8")
        assert "MHTML OCR Smoke" in content
        assert "GUI action-only OCR route" in content
        assert "embedded.png" in content
        assert "> OCR text from GUI MHTML image" in content
        assert window._info_area_vm.has_task_summary
        assert window._info_area_vm._task_summary.state == "success"

    def test_epub_to_markdown_runs_through_gui_action_only_route_with_resource(
        self,
        main_window_with_controller,
        tmp_path: Path,
    ) -> None:
        from ebooklib import epub
        from PySide6.QtWidgets import QApplication

        from docwen_gui.main_window import _normalize_path

        image_bytes = _create_ocr_smoke_png(tmp_path / "epub-image.png").read_bytes()
        book = epub.EpubBook()
        book.set_identifier("docwen-gui-epub-smoke")
        book.set_title("EPUB Smoke")
        book.set_language("en")

        image = epub.EpubItem(
            uid="pic",
            file_name="images/pic.png",
            media_type="image/png",
            content=image_bytes,
        )
        chapter = epub.EpubHtml(title="Chapter 1", file_name="chapter.xhtml", lang="en")
        chapter.content = (
            "<html><body><h1>EPUB Smoke</h1>"
            "<p>GUI action-only route to Markdown.</p>"
            "<p><img src='images/pic.png' /></p></body></html>"
        )
        book.add_item(chapter)
        book.add_item(image)
        book.toc = [chapter]
        book.spine = ["nav", chapter]
        book.add_item(epub.EpubNcx())
        book.add_item(epub.EpubNav())

        source = tmp_path / "epub-smoke.epub"
        epub.write_epub(str(source), book)

        window = main_window_with_controller
        normalized = _normalize_path(str(source))
        window.view_model.add_files([str(source)])

        app = QApplication.instance()
        if app is not None:
            app.processEvents()

        assert _wait_for(lambda: window._action_area_vm.visible, timeout_ms=5000, interval_ms=20)
        assert window._action_area_vm.file_type == "epub"
        window._action_area_vm.request_conversion("md")

        if app is not None:
            app.processEvents()

        def _conversion_finished() -> bool:
            entry = window._batch_list_vm.get_file_entry(normalized)
            return entry is not None and entry.status in ("completed", "failed")

        assert _wait_for(_conversion_finished, timeout_ms=_E2E_CONVERSION_TIMEOUT_MS)
        entry = window._batch_list_vm.get_file_entry(normalized)
        assert entry is not None
        assert entry.status == "completed"
        assert entry.output_path

        output_path = Path(entry.output_path)
        image_path = output_path.with_name("pic.png")
        _assert_markdown_node(output_path, source_stem="epub-smoke", source_tag="Epub")
        assert image_path.exists()
        assert image_path.read_bytes() == image_bytes
        content = output_path.read_text(encoding="utf-8")
        assert "EPUB Smoke" in content
        assert "GUI action-only route to Markdown" in content
        assert "pic.png" in content
        assert window._info_area_vm.has_task_summary
        assert window._info_area_vm._task_summary.state == "success"

    def test_enex_to_markdown_runs_through_gui_action_only_route_with_resource(
        self,
        main_window_with_controller,
        tmp_path: Path,
    ) -> None:
        import base64
        import hashlib

        from PySide6.QtWidgets import QApplication

        from docwen_gui.main_window import _normalize_path

        image_bytes = _create_ocr_smoke_png(tmp_path / "enex-image.png").read_bytes()
        image_hash = hashlib.md5(image_bytes).hexdigest()
        image_b64 = base64.b64encode(image_bytes).decode("ascii")
        source = tmp_path / "enex-smoke.enex"
        source.write_text(
            f"""<?xml version="1.0" encoding="UTF-8"?>
<!DOCTYPE en-export SYSTEM "http://xml.evernote.com/pub/evernote-export3.dtd">
<en-export export-date="20260101T000000Z" application="Evernote" version="10.x">
<note>
    <title>ENEX Smoke</title>
    <content>
        <![CDATA[<?xml version="1.0" encoding="UTF-8"?>
<!DOCTYPE en-note SYSTEM "http://xml.evernote.com/pub/enml2.dtd">
<en-note>
    <p>GUI action-only route to Markdown.</p>
    <en-media hash="{image_hash}" type="image/png" />
</en-note>
]]>
    </content>
    <created>20260101T000000Z</created>
    <updated>20260101T000000Z</updated>
    <resource>
        <data encoding="base64">{image_b64}</data>
        <mime>image/png</mime>
        <resource-attributes>
            <file-name>enex-image.png</file-name>
        </resource-attributes>
    </resource>
</note>
</en-export>""",
            encoding="utf-8",
        )

        window = main_window_with_controller
        normalized = _normalize_path(str(source))
        window.view_model.add_files([str(source)])

        app = QApplication.instance()
        if app is not None:
            app.processEvents()

        assert _wait_for(lambda: window._action_area_vm.visible, timeout_ms=5000, interval_ms=20)
        assert window._action_area_vm.file_type == "enex"
        window._action_area_vm.request_conversion("md")

        if app is not None:
            app.processEvents()

        def _conversion_finished() -> bool:
            entry = window._batch_list_vm.get_file_entry(normalized)
            return entry is not None and entry.status in ("completed", "failed")

        assert _wait_for(_conversion_finished, timeout_ms=_E2E_CONVERSION_TIMEOUT_MS)
        entry = window._batch_list_vm.get_file_entry(normalized)
        assert entry is not None
        assert entry.status == "completed"
        assert entry.output_path

        output_path = Path(entry.output_path)
        image_path = output_path.with_name("enex-image.png")
        _assert_markdown_node(output_path, source_stem="enex-smoke", source_tag="Enex")
        assert image_path.exists()
        assert image_path.read_bytes() == image_bytes
        content = output_path.read_text(encoding="utf-8")
        assert "ENEX Smoke" in content
        assert "GUI action-only route to Markdown" in content
        assert "enex-image.png" in content
        assert window._info_area_vm.has_task_summary
        assert window._info_area_vm._task_summary.state == "success"
