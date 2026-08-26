"""Contract guards for the packaged GUI Office conversion release gate."""

from __future__ import annotations

import json
import subprocess
from pathlib import Path

import pytest

pytestmark = pytest.mark.contract


def _write_pdf(path: Path, text: str) -> None:
    import fitz

    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), text)
    document.save(path)
    document.close()


def test_packaged_gui_office_smoke_drives_panel_and_action_routes(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from scripts.release import verify_packaged_gui

    binary_dir = tmp_path / "dist"
    binary_dir.mkdir()
    binary_name = "DocWen.exe"
    (binary_dir / binary_name).write_text("placeholder", encoding="utf-8")
    monkeypatch.setattr(verify_packaged_gui, "_verify_resource_layout", lambda _path: None)

    calls: list[tuple[str, str, str]] = []

    def fake_run_with_env(
        binary_path: Path,
        *args: str,
        cwd: Path,
        env: dict[str, str],
        timeout: int,
    ) -> subprocess.CompletedProcess[str]:
        del args, timeout
        source = Path(env["DOCWEN_GUI_TEST_CONVERSION_INPUT"])
        surface = env["DOCWEN_GUI_TEST_CONVERSION_SURFACE"]
        tokens = {
            ".docx": "DOCWEN PACKAGED GUI WORD 2026 semantic readback",
            ".xlsx": "DOCWEN PACKAGED GUI SHEET 2026 Alpha 95",
            ".md": "DOCWEN PACKAGED GUI MARKDOWN 2026 semantic readback",
        }
        if source.suffix == ".xlsx":
            from openpyxl import load_workbook

            workbook = load_workbook(source, read_only=False)
            sheet = workbook.active
            assert sheet is not None
            assert sheet.column_dimensions["A"].width >= 40
            assert sheet.column_dimensions["B"].width >= 10
            workbook.close()
        output_dir = Path(env["DOCWEN_GUI_TEST_CONVERSION_OUTPUT_DIR"])
        output_path = output_dir / f"{source.stem}.pdf"
        _write_pdf(output_path, tokens[source.suffix])
        report_path = Path(env["DOCWEN_GUI_TEST_CONVERSION_REPORT"])
        backends = {
            ".docx": "fixture-word",
            ".xlsx": "fixture-spreadsheet",
            ".md": "fixture-markdown-office",
        }
        report_path.write_text(
            json.dumps(
                {
                    "success": True,
                    "status": "completed",
                    "inputPath": str(source),
                    "outputPath": str(output_path),
                    "outputExists": True,
                    "outputBytes": output_path.stat().st_size,
                    "conversionMetrics": {
                        "durationMs": 12.5,
                        "inputBytes": source.stat().st_size,
                        "outputBytes": output_path.stat().st_size,
                        "engine": "office_bridge",
                        "backend": backends[source.suffix],
                    },
                    "error": None,
                }
            ),
            encoding="utf-8",
        )
        log_dir = cwd / "log_home" / "logs"
        log_dir.mkdir(parents=True, exist_ok=True)
        (log_dir / "docwen.log").write_text("ok", encoding="utf-8")
        calls.append((source.suffix, surface, env["DOCWEN_GUI_TEST_CONVERSION_TARGET"]))
        return subprocess.CompletedProcess([str(binary_path)], 0, stdout="", stderr="")

    monkeypatch.setattr(verify_packaged_gui, "_run_with_env", fake_run_with_env)

    exit_code = verify_packaged_gui.main(
        ["--binary-dir", str(binary_dir), "--binary-name", binary_name, "--office-smoke"]
    )

    assert exit_code == 0
    assert calls == [
        (".docx", "panel", "pdf"),
        (".xlsx", "panel", "pdf"),
        (".md", "action", "pdf"),
    ]


@pytest.mark.parametrize(
    "metrics",
    [
        None,
        {"durationMs": 1.0, "inputBytes": 10, "outputBytes": 20, "engine": "", "backend": "Excel"},
        {
            "durationMs": 1.0,
            "inputBytes": 10,
            "outputBytes": 20,
            "engine": "office_bridge",
            "backend": "",
        },
        {
            "durationMs": -1.0,
            "inputBytes": 10,
            "outputBytes": 20,
            "engine": "office_bridge",
            "backend": "Excel",
        },
        {
            "durationMs": 1.0,
            "inputBytes": 0,
            "outputBytes": 20,
            "engine": "office_bridge",
            "backend": "Excel",
        },
        {
            "durationMs": 1.0,
            "inputBytes": 10,
            "outputBytes": 0,
            "engine": "office_bridge",
            "backend": "Excel",
        },
        {
            "durationMs": True,
            "inputBytes": 10,
            "outputBytes": 20,
            "engine": "office_bridge",
            "backend": "Excel",
        },
        {
            "durationMs": 1.0,
            "inputBytes": True,
            "outputBytes": 20,
            "engine": "office_bridge",
            "backend": "Excel",
        },
        {
            "durationMs": 1.0,
            "inputBytes": 10,
            "outputBytes": True,
            "engine": "office_bridge",
            "backend": "Excel",
        },
        {
            "durationMs": float("nan"),
            "inputBytes": 10,
            "outputBytes": 20,
            "engine": "office_bridge",
            "backend": "Excel",
        },
        {
            "durationMs": float("inf"),
            "inputBytes": 10,
            "outputBytes": 20,
            "engine": "office_bridge",
            "backend": "Excel",
        },
        {
            "durationMs": float("-inf"),
            "inputBytes": 10,
            "outputBytes": 20,
            "engine": "office_bridge",
            "backend": "Excel",
        },
        {
            "durationMs": 1.0,
            "inputBytes": 10,
            "outputBytes": 20,
            "engine": "office_bridge",
            "backend": True,
        },
    ],
)
def test_packaged_gui_office_smoke_rejects_invalid_metrics(tmp_path: Path, metrics: object) -> None:
    from scripts.release import verify_packaged_gui

    input_path = tmp_path / "input.xlsx"
    input_path.write_bytes(b"i" * 10)
    output_path = tmp_path / "output.pdf"
    output_path.write_bytes(b"o" * 20)
    report_path = tmp_path / "report.json"
    report_path.write_text(
        json.dumps(
            {
                "inputPath": str(input_path),
                "outputPath": str(output_path),
                "outputBytes": output_path.stat().st_size,
                "conversionMetrics": metrics,
            }
        ),
        encoding="utf-8",
    )

    with pytest.raises(RuntimeError, match="packaged_gui_office_metrics_invalid: xlsx"):
        verify_packaged_gui._verify_office_conversion_metrics(
            report_path,
            case_name="xlsx",
            input_path=input_path,
            output_path=output_path,
        )


@pytest.mark.parametrize(
    "mismatch",
    [
        "input_path",
        "output_path",
        "input_metric_bytes",
        "output_metric_bytes",
        "reported_output_bytes",
    ],
)
def test_packaged_gui_office_smoke_rejects_metrics_provenance_mismatch(
    tmp_path: Path,
    mismatch: str,
) -> None:
    from scripts.release import verify_packaged_gui

    input_path = tmp_path / "input.xlsx"
    input_path.write_bytes(b"input-bytes")
    output_path = tmp_path / "output.pdf"
    output_path.write_bytes(b"%PDF-output-bytes")
    report = {
        "inputPath": str(input_path),
        "outputPath": str(output_path),
        "outputBytes": output_path.stat().st_size,
        "conversionMetrics": {
            "durationMs": 1.0,
            "inputBytes": input_path.stat().st_size,
            "outputBytes": output_path.stat().st_size,
            "engine": "office_bridge",
            "backend": "fixture-office",
        },
    }
    metrics = report["conversionMetrics"]
    assert isinstance(metrics, dict)
    if mismatch == "input_path":
        report["inputPath"] = str(tmp_path / "other.xlsx")
    elif mismatch == "output_path":
        report["outputPath"] = str(tmp_path / "other.pdf")
    elif mismatch == "input_metric_bytes":
        metrics["inputBytes"] = input_path.stat().st_size + 1
    elif mismatch == "output_metric_bytes":
        metrics["outputBytes"] = output_path.stat().st_size + 1
    else:
        report["outputBytes"] = output_path.stat().st_size + 1
    report_path = tmp_path / "report.json"
    report_path.write_text(json.dumps(report), encoding="utf-8")

    with pytest.raises(RuntimeError, match="packaged_gui_office_metrics_invalid: xlsx"):
        verify_packaged_gui._verify_office_conversion_metrics(
            report_path,
            case_name="xlsx",
            input_path=input_path,
            output_path=output_path,
        )


@pytest.mark.parametrize("duration_ms", [0.0, 1e308])
def test_packaged_gui_office_smoke_accepts_finite_duration_without_product_threshold(
    tmp_path: Path,
    duration_ms: float,
) -> None:
    from scripts.release import verify_packaged_gui

    input_path = tmp_path / "input.xlsx"
    input_path.write_bytes(b"input-bytes")
    output_path = tmp_path / "output.pdf"
    output_path.write_bytes(b"%PDF-output-bytes")
    report_path = tmp_path / "report.json"
    report_path.write_text(
        json.dumps(
            {
                "inputPath": str(input_path),
                "outputPath": str(output_path),
                "outputBytes": output_path.stat().st_size,
                "conversionMetrics": {
                    "durationMs": duration_ms,
                    "inputBytes": input_path.stat().st_size,
                    "outputBytes": output_path.stat().st_size,
                    "engine": "office_bridge",
                    "backend": "fixture-office",
                },
            }
        ),
        encoding="utf-8",
    )

    verify_packaged_gui._verify_office_conversion_metrics(
        report_path,
        case_name="xlsx",
        input_path=input_path,
        output_path=output_path,
    )


def test_generic_pdf_smoke_does_not_require_office_metrics(tmp_path: Path) -> None:
    from scripts.release import verify_packaged_gui

    output_path = tmp_path / "output.pdf"
    expected = "generic packaged GUI smoke"
    _write_pdf(output_path, expected)
    report_path = tmp_path / "report.json"
    report_path.write_text(
        json.dumps(
            {
                "success": True,
                "status": "completed",
                "outputPath": str(output_path),
                "outputExists": True,
                "error": None,
            }
        ),
        encoding="utf-8",
    )

    assert (
        verify_packaged_gui._verify_conversion_smoke_report(
            report_path,
            expected_tokens=(expected,),
        )
        == output_path
    )


def test_packaged_gui_office_smoke_rejects_pdf_without_expected_semantics(tmp_path: Path) -> None:
    from scripts.release import verify_packaged_gui

    output_path = tmp_path / "output.pdf"
    _write_pdf(output_path, "wrong output")
    report_path = tmp_path / "report.json"
    report_path.write_text(
        json.dumps(
            {
                "success": True,
                "status": "completed",
                "outputPath": str(output_path),
                "outputExists": True,
                "error": None,
            }
        ),
        encoding="utf-8",
    )

    with pytest.raises(RuntimeError, match="packaged_gui_conversion_pdf_semantics_missing"):
        verify_packaged_gui._verify_conversion_smoke_report(
            report_path,
            expected_tokens=("DOCWEN PACKAGED GUI WORD 2026",),
        )


def test_packaged_gui_presentation_fixture_uses_legacy_powerpoint_format(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from pptx import Presentation
    from scripts.release import verify_packaged_gui

    from docwen_core.office_bridge import BridgeResult

    calls: list[tuple[str, str, str, list[str], list[tuple[str, int]], str | None]] = []

    def fake_convert(
        input_path: str,
        output_path: str,
        *,
        source_format: str,
        backend_priority: list[str],
        com_candidates: dict[str, object],
        libreoffice_format: str | None,
        **_kwargs: object,
    ) -> BridgeResult:
        candidates = [(candidate.prog_id, candidate.save_format) for candidate in com_candidates.values()]  # type: ignore[attr-defined]
        calls.append((input_path, output_path, source_format, backend_priority, candidates, libreoffice_format))
        Path(output_path).write_bytes(b"legacy ppt fixture")
        return BridgeResult(True, output_path=output_path, backend="fixture PowerPoint")

    monkeypatch.setattr("docwen_core.office_bridge.convert_with_backend_priority", fake_convert)

    ppt_path, backend = verify_packaged_gui._write_presentation_smoke_input(tmp_path)

    assert backend == "fixture PowerPoint"
    assert ppt_path.name == "packaged-gui-presentation.ppt"
    assert calls[0][2] == "pptx"
    assert calls[0][3] == ["wps_presentation", "msoffice_powerpoint", "libreoffice"]
    assert calls[0][4] == [("Kwpp.Application", 1), ("PowerPoint.Application", 1)]
    assert calls[0][5] == "ppt"
    source = Presentation(calls[0][0])
    assert source.core_properties.title in (None, "")
    assert source.slides[0].shapes.title.text == "DOCWEN PACKAGED GUI PRESENTATION 2026"  # type: ignore[union-attr]
    assert "Alpha" in source.slides[1].placeholders[1].text


def test_packaged_gui_presentation_smoke_drives_action_route(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from scripts.release import verify_packaged_gui

    binary_dir = tmp_path / "dist"
    binary_dir.mkdir()
    binary_name = "DocWen.exe"
    (binary_dir / binary_name).write_text("placeholder", encoding="utf-8")
    monkeypatch.setattr(verify_packaged_gui, "_verify_resource_layout", lambda _path: None)

    def fake_fixture(work_dir: Path) -> tuple[Path, str]:
        source = work_dir / "inputs" / "packaged-gui-presentation.ppt"
        source.parent.mkdir(parents=True)
        source.write_bytes(b"legacy ppt fixture")
        return source, "fixture WPS Presentation"

    observed: dict[str, str] = {}

    def fake_run_with_env(
        binary_path: Path,
        *args: str,
        cwd: Path,
        env: dict[str, str],
        timeout: int,
    ) -> subprocess.CompletedProcess[str]:
        del args, timeout
        observed.update(env)
        output_dir = Path(env["DOCWEN_GUI_TEST_CONVERSION_OUTPUT_DIR"])
        root_name = "packaged-gui-presentation_20260824_120102_fromPpt"
        output_path = output_dir / root_name / f"{root_name}.md"
        output_path.parent.mkdir()
        output_path.write_text(
            "# DOCWEN PACKAGED GUI PRESENTATION 2026\n\n"
            "Legacy PowerPoint bridge semantic readback.\n\n"
            "## Acceptance Matrix\n\nAlpha\nBeta\n",
            encoding="utf-8",
        )
        Path(env["DOCWEN_GUI_TEST_CONVERSION_REPORT"]).write_text(
            json.dumps(
                {
                    "success": True,
                    "status": "completed",
                    "outputPath": str(output_path),
                    "outputExists": True,
                    "error": None,
                }
            ),
            encoding="utf-8",
        )
        log_dir = cwd / "log_home" / "logs"
        log_dir.mkdir(parents=True, exist_ok=True)
        (log_dir / "docwen.log").write_text("ok", encoding="utf-8")
        return subprocess.CompletedProcess([str(binary_path)], 0, stdout="", stderr="")

    monkeypatch.setattr(verify_packaged_gui, "_write_presentation_smoke_input", fake_fixture)
    monkeypatch.setattr(verify_packaged_gui, "_run_with_env", fake_run_with_env)

    exit_code = verify_packaged_gui.main(
        ["--binary-dir", str(binary_dir), "--binary-name", binary_name, "--presentation-smoke"]
    )

    assert exit_code == 0
    assert observed["DOCWEN_GUI_TEST_CONVERSION_TARGET"] == "md"
    assert observed["DOCWEN_GUI_TEST_CONVERSION_SURFACE"] == "action"
    assert observed["DOCWEN_GUI_TEST_CONVERSION_INPUT"].endswith("packaged-gui-presentation.ppt")


def test_packaged_gui_presentation_smoke_rejects_staging_name(tmp_path: Path) -> None:
    from scripts.release import verify_packaged_gui

    output_path = tmp_path / "auxiliary_1.md"
    output_path.write_text("DOCWEN PACKAGED GUI PRESENTATION 2026", encoding="utf-8")
    report_path = tmp_path / "report.json"
    report_path.write_text(
        json.dumps(
            {
                "success": True,
                "status": "completed",
                "outputPath": str(output_path),
                "outputExists": True,
                "error": None,
            }
        ),
        encoding="utf-8",
    )

    with pytest.raises(RuntimeError, match="packaged_gui_conversion_markdown_name_unexpected"):
        verify_packaged_gui._verify_markdown_conversion_smoke_report(
            report_path,
            expected_source_stem="packaged-gui-presentation",
            expected_source_format="Ppt",
            expected_tokens=("DOCWEN PACKAGED GUI PRESENTATION 2026",),
        )
