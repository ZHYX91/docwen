from __future__ import annotations

import argparse
import contextlib
import csv
import hashlib
import json
import math
import os
import re
import shutil
import signal
import stat
import subprocess
import sys
import tempfile
import time
from pathlib import Path
from typing import Any, cast

try:
    from scripts.release import packaged_resources as _packaged_resources
    from scripts.release.successful_warning_fixture import (
        SUCCESSFUL_WARNING_ACTION,
        SUCCESSFUL_WARNING_CODE,
        SUCCESSFUL_WARNING_MESSAGE,
        write_successful_warning_fixture,
    )
except ModuleNotFoundError:
    import packaged_resources as _packaged_resources  # type: ignore[no-redef]
    from successful_warning_fixture import (  # type: ignore[no-redef]
        SUCCESSFUL_WARNING_ACTION,
        SUCCESSFUL_WARNING_CODE,
        SUCCESSFUL_WARNING_MESSAGE,
        write_successful_warning_fixture,
    )

_REQUIRED_ASSET_FILES = _packaged_resources.REQUIRED_ASSET_FILES
_REQUIRED_CONFIG_FILES = _packaged_resources.REQUIRED_CONFIG_FILES
_REQUIRED_LOCALE_FILES = _packaged_resources.REQUIRED_LOCALE_FILES
_REQUIRED_MODEL_FILES = _packaged_resources.REQUIRED_MODEL_FILES
_REQUIRED_TEMPLATE_FILES = _packaged_resources.REQUIRED_TEMPLATE_FILES
missing_files = _packaged_resources.missing_files
verify_common_resource_layout = _packaged_resources.verify_common_resource_layout

_RUNTIME_FAILURE_MARKERS: tuple[str, ...] = (
    "Failed to load plugin ",
    "Traceback (most recent call last)",
    "CRITICAL",
    "Unhandled exception",
)
_RELEVANT_PROCESS_NAMES: frozenset[str] = frozenset(
    {
        "docwen.exe",
        "docwencli.exe",
        "winword.exe",
        "excel.exe",
        "powerpnt.exe",
        "wps.exe",
        "wpp.exe",
        "et.exe",
        "soffice.exe",
        "soffice.bin",
    }
)
_SUCCESSFUL_WARNING_LOCALE = "zh_CN"
_GENERATED_SUCCESSFUL_WARNING_FIXTURE = Path("__docwen_generated_successful_warning_fixture__")
_REQUIRED_SETTINGS_TAB_KEYS: tuple[str, ...] = (
    "general",
    "text",
    "proofread",
    "document",
    "spreadsheet",
    "image",
    "layout",
    "link",
    "formatting",
    "output",
    "export",
    "logging",
    "other",
)
_REQUIRED_SETTINGS_PAGE_MODULES: frozenset[str] = frozenset(
    f"docwen_gui.widgets.settings.{key}_tab" for key in _REQUIRED_SETTINGS_TAB_KEYS
)
_FILE_ATTRIBUTE_REPARSE_POINT = getattr(stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0x400)


def _default_binary_name() -> str:
    return "DocWen.exe" if os.name == "nt" else "DocWen"


def _default_cli_binary_name() -> str:
    return "DocWenCLI.exe" if os.name == "nt" else "DocWenCLI"


def _is_link_or_reparse(path_stat: os.stat_result) -> bool:
    return stat.S_ISLNK(path_stat.st_mode) or bool(
        getattr(path_stat, "st_file_attributes", 0) & _FILE_ATTRIBUTE_REPARSE_POINT
    )


def _lexical_path_text(path: Path) -> str:
    return os.path.normcase(os.fspath(path))


def _is_same_or_descendant(path: Path, root: Path) -> bool:
    try:
        path.relative_to(root)
    except ValueError:
        return False
    return True


def _assert_directory_chain_without_aliases(directory: Path) -> Path:
    if not directory.is_absolute():
        raise ValueError(f"packaged_gui_evidence_parent_not_absolute: {directory}")

    current = Path(directory.anchor)
    for part in directory.parts[1:]:
        current /= part
        try:
            path_stat = current.lstat()
        except FileNotFoundError as exc:
            raise FileNotFoundError(f"packaged_gui_evidence_parent_missing: {current}") from exc
        if _is_link_or_reparse(path_stat):
            raise RuntimeError(f"packaged_gui_evidence_parent_link_rejected: {current}")
        if not stat.S_ISDIR(path_stat.st_mode):
            raise NotADirectoryError(f"packaged_gui_evidence_parent_not_directory: {current}")

    resolved = directory.resolve(strict=True)
    if _lexical_path_text(resolved) != _lexical_path_text(directory):
        raise RuntimeError(f"packaged_gui_evidence_parent_alias_rejected: {directory} -> {resolved}")
    return resolved


def _validate_evidence_destination(
    requested: Path,
    *,
    binary_dir: Path,
    verification_dir: Path | None = None,
) -> Path:
    if not requested.is_absolute():
        raise ValueError(f"packaged_gui_evidence_dir_not_absolute: {requested}")

    normalized = Path(os.path.abspath(requested))
    if _lexical_path_text(normalized) != _lexical_path_text(requested):
        raise RuntimeError(f"packaged_gui_evidence_dir_alias_rejected: {requested} -> {normalized}")
    if os.path.lexists(normalized):
        raise FileExistsError(f"packaged_gui_evidence_dir_exists: {normalized}")

    resolved_parent = _assert_directory_chain_without_aliases(normalized.parent)
    destination = resolved_parent / normalized.name
    resolved_binary_dir = binary_dir.resolve(strict=True)
    if _is_same_or_descendant(destination, resolved_binary_dir):
        raise RuntimeError(f"packaged_gui_evidence_dir_inside_binary_dir: {destination}")
    if verification_dir is not None:
        resolved_verification_dir = verification_dir.resolve(strict=True)
        if _is_same_or_descendant(destination, resolved_verification_dir):
            raise RuntimeError(f"packaged_gui_evidence_dir_inside_verification_dir: {destination}")
    return destination


def _hash_regular_file(path: Path) -> tuple[int, str]:
    before = path.lstat()
    if _is_link_or_reparse(before) or not stat.S_ISREG(before.st_mode):
        raise RuntimeError(f"packaged_gui_evidence_unsafe_file: {path}")
    digest = hashlib.sha256()
    size = 0
    with path.open("rb") as stream:
        while chunk := stream.read(1024 * 1024):
            digest.update(chunk)
            size += len(chunk)
    after = path.lstat()
    if (
        _is_link_or_reparse(after)
        or not stat.S_ISREG(after.st_mode)
        or size != before.st_size
        or size != after.st_size
        or before.st_mtime_ns != after.st_mtime_ns
    ):
        raise RuntimeError(f"packaged_gui_evidence_file_changed: {path}")
    return size, digest.hexdigest()


def _capture_evidence_tree(root: Path) -> tuple[tuple[str, ...], dict[str, tuple[int, str]]]:
    root_stat = root.lstat()
    if _is_link_or_reparse(root_stat) or not stat.S_ISDIR(root_stat.st_mode):
        raise RuntimeError(f"packaged_gui_evidence_unsafe_root: {root}")

    directories: list[str] = []
    files: dict[str, tuple[int, str]] = {}

    def visit(directory: Path) -> None:
        for child in sorted(directory.iterdir(), key=lambda path: path.name):
            relative = child.relative_to(root).as_posix()
            child_stat = child.lstat()
            if _is_link_or_reparse(child_stat):
                raise RuntimeError(f"packaged_gui_evidence_link_rejected: {child}")
            if stat.S_ISDIR(child_stat.st_mode):
                directories.append(relative)
                visit(child)
            elif stat.S_ISREG(child_stat.st_mode):
                files[relative] = _hash_regular_file(child)
            else:
                raise RuntimeError(f"packaged_gui_evidence_unsafe_entry: {child}")

    visit(root)
    return tuple(directories), files


def _verify_file_bytes_equal(source: Path, destination: Path) -> None:
    with source.open("rb") as source_stream, destination.open("rb") as destination_stream:
        while True:
            source_chunk = source_stream.read(1024 * 1024)
            destination_chunk = destination_stream.read(1024 * 1024)
            if source_chunk != destination_chunk:
                raise RuntimeError(f"packaged_gui_evidence_bytes_mismatch: {source} -> {destination}")
            if not source_chunk:
                return


def _preserve_success_evidence(
    verification_dir: Path,
    requested_destination: Path,
    *,
    binary_dir: Path,
) -> Path:
    destination = _validate_evidence_destination(
        requested_destination,
        binary_dir=binary_dir,
        verification_dir=verification_dir,
    )
    source_before = _capture_evidence_tree(verification_dir)
    shutil.copytree(verification_dir, destination, symlinks=True, dirs_exist_ok=False)
    source_after = _capture_evidence_tree(verification_dir)
    destination_tree = _capture_evidence_tree(destination)
    if source_before != source_after:
        raise RuntimeError("packaged_gui_evidence_source_changed_during_copy")
    if source_after != destination_tree:
        raise RuntimeError(
            f"packaged_gui_evidence_manifest_mismatch: source={source_after!r}; destination={destination_tree!r}"
        )
    for relative_path in source_after[1]:
        _verify_file_bytes_equal(verification_dir / relative_path, destination / relative_path)
    return destination


def _verify_resource_layout(binary_dir: Path) -> None:
    verify_common_resource_layout(binary_dir, error_prefix="packaged_gui")

    missing_assets = missing_files(binary_dir / "assets", _REQUIRED_ASSET_FILES)
    if missing_assets:
        raise RuntimeError(f"packaged_gui_assets_missing: {missing_assets}")


def _read_pyinstaller_module_names(binary_path: Path) -> frozenset[str]:
    """Read Python module names from a PyInstaller executable's embedded PYZ."""

    try:
        from PyInstaller.archive.readers import CArchiveReader

        archive = CArchiveReader(str(binary_path))
        pyz_names = [name for name in archive.toc if name.upper().endswith("PYZ.PYZ")]
        if len(pyz_names) != 1:
            raise RuntimeError(f"expected one embedded PYZ, found {pyz_names}")
        pyz = archive.open_embedded_archive(pyz_names[0])
        return frozenset(str(name) for name in pyz.toc)
    except Exception as exc:
        raise RuntimeError(f"packaged_gui_archive_unreadable: {binary_path}: {exc}") from exc


def _verify_settings_page_archive(binary_path: Path) -> None:
    modules = _read_pyinstaller_module_names(binary_path)
    missing = sorted(_REQUIRED_SETTINGS_PAGE_MODULES - modules)
    if missing:
        raise RuntimeError(f"packaged_gui_settings_modules_missing: {missing}")


def _verify_settings_smoke_report(report_path: Path) -> dict[str, Any]:
    if not report_path.is_file():
        raise RuntimeError(f"packaged_gui_settings_report_missing: {report_path}")
    try:
        report = json.loads(report_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise RuntimeError(f"packaged_gui_settings_report_invalid: {report_path}") from exc

    expected_tabs = list(_REQUIRED_SETTINGS_TAB_KEYS)
    actual_expected = report.get("expectedTabs")
    loaded_tabs = report.get("loadedTabs")
    if (
        report.get("success") is not True
        or not isinstance(actual_expected, list)
        or set(actual_expected) != set(expected_tabs)
        or not isinstance(loaded_tabs, list)
        or set(loaded_tabs) != set(expected_tabs)
        or report.get("failedTabs") != []
        or report.get("missingTabs") != []
        or report.get("unexpectedTabs") != []
        or report.get("error") is not None
    ):
        raise RuntimeError(f"packaged_gui_settings_smoke_failed: {report}")
    return cast(dict[str, Any], report)


def _write_ocr_png(path: Path) -> None:
    from PIL import Image, ImageDraw, ImageFont

    image = Image.new("RGB", (900, 240), "white")
    draw = ImageDraw.Draw(image)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 56)
    except OSError:
        font = ImageFont.load_default()
    draw.text((36, 72), "HELLO DOCWEN OCR", fill="black", font=font)
    image.save(path)


def _run(binary_path: Path, *args: str, cwd: Path) -> subprocess.CompletedProcess[str]:
    env = _base_env(cwd)
    # Set autoclose so the GUI app exits automatically.
    env["DOCWEN_GUI_TEST_AUTOCLOSE_MS"] = os.environ.get("DOCWEN_GUI_TEST_AUTOCLOSE_MS", "2000") or "2000"
    return _run_with_env(binary_path, *args, cwd=cwd, env=env, timeout=30)


def _base_env(cwd: Path) -> dict[str, str]:
    env = os.environ.copy()
    env["PYTHONUTF8"] = "1"
    env["PYTHONIOENCODING"] = "utf-8"
    env["DOCWEN_CONFIG_DIR"] = str(cwd / "config_home")
    env["DOCWEN_LOG_DIR"] = str(cwd / "log_home")
    env["DOCWEN_LOG_TO_TEMP"] = ""

    # In headless CI environments, ensure it can run offscreen if needed
    if sys.platform == "linux" and not env.get("DISPLAY"):
        env["QT_QPA_PLATFORM"] = "offscreen"

    return env


def _run_with_env(
    binary_path: Path,
    *args: str,
    cwd: Path,
    env: dict[str, str],
    timeout: int,
) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [str(binary_path), *args],
        cwd=cwd,
        env=env,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        check=False,
        timeout=timeout,
    )


def _write_office_smoke_inputs(work_dir: Path) -> list[tuple[str, Path, str, tuple[str, ...]]]:
    from docx import Document
    from openpyxl import Workbook

    inputs_dir = work_dir / "inputs"
    inputs_dir.mkdir(parents=True, exist_ok=True)

    docx_path = inputs_dir / "packaged-gui-word.docx"
    document = Document()
    document.add_heading("DOCWEN PACKAGED GUI WORD 2026", level=1)
    document.add_paragraph("Frozen MainWindow Office route semantic readback.")
    document.save(str(docx_path))

    xlsx_path = inputs_dir / "packaged-gui-sheet.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    if sheet is None:
        raise RuntimeError("packaged_gui_office_workbook_sheet_missing")
    sheet.title = "Smoke"
    sheet.append(["DOCWEN PACKAGED GUI SHEET 2026", "Score"])
    sheet.append(["Alpha", 95])
    sheet.column_dimensions["A"].width = 42
    sheet.column_dimensions["B"].width = 12
    workbook.save(str(xlsx_path))
    workbook.close()

    markdown_path = inputs_dir / "packaged-gui-markdown.md"
    markdown_path.write_text(
        "# DOCWEN PACKAGED GUI MARKDOWN 2026\n\nFrozen ActionArea Office route semantic readback.\n",
        encoding="utf-8",
    )

    return [
        ("docx", docx_path, "panel", ("DOCWEN PACKAGED GUI WORD 2026", "semantic readback")),
        ("xlsx", xlsx_path, "panel", ("DOCWEN PACKAGED GUI SHEET 2026", "Alpha", "95")),
        ("markdown", markdown_path, "action", ("DOCWEN PACKAGED GUI MARKDOWN 2026", "semantic readback")),
    ]


def _run_office_smoke(
    binary_path: Path,
    *,
    cwd: Path,
) -> tuple[list[Path], subprocess.CompletedProcess[str]]:
    work_dir = cwd / "gui_office_smoke"
    if work_dir.exists():
        shutil.rmtree(work_dir)
    work_dir.mkdir(parents=True)

    outputs: list[Path] = []
    stdout_parts: list[str] = []
    stderr_parts: list[str] = []
    for case_name, source_path, surface, expected_tokens in _write_office_smoke_inputs(work_dir):
        case_dir = work_dir / case_name
        output_dir = case_dir / "outputs"
        output_dir.mkdir(parents=True)
        report_path = case_dir / "report.json"

        env = _base_env(cwd)
        env["DOCWEN_GUI_TEST_AUTOCLOSE_MS"] = "120000"
        env["DOCWEN_GUI_TEST_CONVERSION_REPORT"] = str(report_path)
        env["DOCWEN_GUI_TEST_CONVERSION_INPUT"] = str(source_path)
        env["DOCWEN_GUI_TEST_CONVERSION_OUTPUT_DIR"] = str(output_dir)
        env["DOCWEN_GUI_TEST_CONVERSION_TARGET"] = "pdf"
        env["DOCWEN_GUI_TEST_CONVERSION_SURFACE"] = surface
        env["DOCWEN_GUI_TEST_CONVERSION_TIMEOUT_MS"] = "90000"
        proc = _run_with_env(binary_path, cwd=cwd, env=env, timeout=120)
        stdout_parts.append(proc.stdout)
        stderr_parts.append(proc.stderr)
        if proc.returncode != 0:
            raise RuntimeError(
                f"packaged_gui_office_process_failed: {case_name}: {proc.returncode}\n"
                f"STDOUT:\n{proc.stdout}\nSTDERR:\n{proc.stderr}"
            )
        output_path = _verify_conversion_smoke_report(report_path, expected_tokens=expected_tokens)
        _verify_office_conversion_metrics(
            report_path,
            case_name=case_name,
            input_path=source_path,
            output_path=output_path,
        )
        outputs.append(output_path)

    return outputs, subprocess.CompletedProcess(
        [str(binary_path), "--office-smoke"],
        0,
        stdout="".join(stdout_parts),
        stderr="".join(stderr_parts),
    )


def _write_presentation_smoke_input(work_dir: Path) -> tuple[Path, str]:
    from pptx import Presentation

    from docwen_core.office_bridge import BridgeCandidate, convert_with_backend_priority

    inputs_dir = work_dir / "inputs"
    inputs_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = inputs_dir / "packaged-gui-presentation-source.pptx"
    ppt_path = inputs_dir / "packaged-gui-presentation.ppt"

    presentation = Presentation()
    presentation.core_properties.title = ""
    first = presentation.slides.add_slide(presentation.slide_layouts[1])
    cast(Any, first.shapes.title).text = "DOCWEN PACKAGED GUI PRESENTATION 2026"
    cast(Any, first.placeholders[1]).text = "Legacy PowerPoint bridge semantic readback."
    second = presentation.slides.add_slide(presentation.slide_layouts[1])
    cast(Any, second.shapes.title).text = "Acceptance Matrix"
    cast(Any, second.placeholders[1]).text = "Alpha\nBeta"
    presentation.save(str(pptx_path))

    presentation_candidates = {
        "wps_presentation": BridgeCandidate("WPS Presentation", "Kwpp.Application", 1, "powerpoint"),
        "msoffice_powerpoint": BridgeCandidate(
            "Microsoft PowerPoint",
            "PowerPoint.Application",
            1,
            "powerpoint",
        ),
    }
    result = convert_with_backend_priority(
        str(pptx_path),
        str(ppt_path),
        source_format="pptx",
        backend_priority=[*presentation_candidates, "libreoffice"],
        com_candidates=presentation_candidates,
        libreoffice_format="ppt",
        com_timeout_s=90.0,
        failure_subject="Packaged GUI presentation fixture backends",
    )
    if not result.success or not ppt_path.is_file():
        raise RuntimeError(f"packaged_gui_presentation_fixture_failed: {result.message}")
    return ppt_path, result.backend


def _run_presentation_smoke(
    binary_path: Path,
    *,
    cwd: Path,
) -> tuple[Path, str, subprocess.CompletedProcess[str]]:
    work_dir = cwd / "gui_presentation_smoke"
    if work_dir.exists():
        shutil.rmtree(work_dir)
    work_dir.mkdir(parents=True)
    source_path, fixture_backend = _write_presentation_smoke_input(work_dir)
    output_dir = work_dir / "outputs"
    output_dir.mkdir()
    report_path = work_dir / "report.json"

    env = _base_env(cwd)
    env["DOCWEN_GUI_TEST_AUTOCLOSE_MS"] = "120000"
    env["DOCWEN_GUI_TEST_CONVERSION_REPORT"] = str(report_path)
    env["DOCWEN_GUI_TEST_CONVERSION_INPUT"] = str(source_path)
    env["DOCWEN_GUI_TEST_CONVERSION_OUTPUT_DIR"] = str(output_dir)
    env["DOCWEN_GUI_TEST_CONVERSION_TARGET"] = "md"
    env["DOCWEN_GUI_TEST_CONVERSION_SURFACE"] = "action"
    env["DOCWEN_GUI_TEST_CONVERSION_TIMEOUT_MS"] = "90000"
    proc = _run_with_env(binary_path, cwd=cwd, env=env, timeout=120)
    if proc.returncode != 0:
        raise RuntimeError(
            f"packaged_gui_presentation_process_failed: {proc.returncode}\n"
            f"STDOUT:\n{proc.stdout}\nSTDERR:\n{proc.stderr}"
        )
    output_path = _verify_markdown_conversion_smoke_report(
        report_path,
        expected_source_stem="packaged-gui-presentation",
        expected_source_format="Ppt",
        expected_tokens=(
            "DOCWEN PACKAGED GUI PRESENTATION 2026",
            "Legacy PowerPoint bridge semantic readback.",
            "Acceptance Matrix",
            "Alpha",
            "Beta",
        ),
    )
    return output_path, fixture_backend, proc


def _write_smartdoc_smoke_inputs(work_dir: Path) -> list[tuple[str, Path, str, tuple[str, ...]]]:
    from docx import Document

    from docwen_core.office_bridge import BridgeCandidate, convert_with_backend_priority

    inputs_dir = work_dir / "inputs"
    inputs_dir.mkdir(parents=True, exist_ok=True)
    source_path = inputs_dir / "packaged-gui-smartdoc-source.docx"
    document = Document()
    document.add_heading("DOCWEN PACKAGED GUI SMARTDOC 2026", level=0)
    document.add_paragraph("External document bridge semantic readback.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Item"
    table.cell(0, 1).text = "Qty"
    table.cell(1, 0).text = "Alpha"
    table.cell(1, 1).text = "7"
    document.save(str(source_path))

    fixture_specs = (
        (
            "doc",
            {
                "wps_writer": BridgeCandidate("WPS Writer", "Kwps.Application", 0, "word"),
                "msoffice_word": BridgeCandidate("Microsoft Word", "Word.Application", 0, "word"),
            },
        ),
        (
            "rtf",
            {
                "wps_writer": BridgeCandidate("WPS Writer", "Kwps.Application", 6, "word"),
                "msoffice_word": BridgeCandidate("Microsoft Word", "Word.Application", 6, "word"),
            },
        ),
        (
            "odt",
            {"msoffice_word": BridgeCandidate("Microsoft Word", "Word.Application", 23, "word")},
        ),
    )
    expected_tokens = (
        "DOCWEN PACKAGED GUI SMARTDOC 2026",
        "External document bridge semantic readback.",
        "Item",
        "Qty",
        "Alpha",
        "7",
    )
    cases: list[tuple[str, Path, str, tuple[str, ...]]] = []
    for source_format, candidates in fixture_specs:
        fixture_path = inputs_dir / f"packaged-gui-smartdoc-{source_format}.{source_format}"
        result = convert_with_backend_priority(
            str(source_path),
            str(fixture_path),
            source_format="docx",
            backend_priority=[*candidates, "libreoffice"],
            com_candidates=candidates,
            libreoffice_format=source_format,
            com_timeout_s=90.0,
            failure_subject=f"Packaged GUI DOCX→{source_format.upper()} fixture backends",
        )
        if not result.success or not fixture_path.is_file():
            raise RuntimeError(f"packaged_gui_smartdoc_fixture_failed: {source_format}: {result.message}")
        cases.append((source_format, fixture_path, result.backend, expected_tokens))
    return cases


def _run_smartdoc_smoke(
    binary_path: Path,
    *,
    cwd: Path,
) -> tuple[list[Path], list[str], subprocess.CompletedProcess[str]]:
    work_dir = cwd / "gui_smartdoc_smoke"
    if work_dir.exists():
        shutil.rmtree(work_dir)
    work_dir.mkdir(parents=True)
    processes_before = _snapshot_relevant_processes()

    outputs: list[Path] = []
    fixture_backends: list[str] = []
    stdout_parts: list[str] = []
    stderr_parts: list[str] = []
    for case_name, source_path, fixture_backend, expected_tokens in _write_smartdoc_smoke_inputs(work_dir):
        case_dir = work_dir / case_name
        output_dir = case_dir / "outputs"
        output_dir.mkdir(parents=True)
        report_path = case_dir / "report.json"

        env = _base_env(cwd)
        env["DOCWEN_GUI_TEST_AUTOCLOSE_MS"] = "120000"
        env["DOCWEN_GUI_TEST_CONVERSION_REPORT"] = str(report_path)
        env["DOCWEN_GUI_TEST_CONVERSION_INPUT"] = str(source_path)
        env["DOCWEN_GUI_TEST_CONVERSION_OUTPUT_DIR"] = str(output_dir)
        env["DOCWEN_GUI_TEST_CONVERSION_TARGET"] = "docx"
        env["DOCWEN_GUI_TEST_CONVERSION_SURFACE"] = "panel"
        env["DOCWEN_GUI_TEST_CONVERSION_TIMEOUT_MS"] = "90000"
        proc = _run_with_env(binary_path, cwd=cwd, env=env, timeout=120)
        stdout_parts.append(proc.stdout)
        stderr_parts.append(proc.stderr)
        if proc.returncode != 0:
            raise RuntimeError(
                f"packaged_gui_smartdoc_process_failed: {case_name}: {proc.returncode}\n"
                f"STDOUT:\n{proc.stdout}\nSTDERR:\n{proc.stderr}"
            )
        outputs.append(
            _verify_docx_conversion_smoke_report(
                report_path,
                expected_name=f"packaged-gui-smartdoc-{case_name}.docx",
                expected_tokens=expected_tokens,
            )
        )
        fixture_backends.append(f"{case_name}={fixture_backend}")

    _wait_for_no_new_relevant_processes(processes_before)

    return (
        outputs,
        fixture_backends,
        subprocess.CompletedProcess(
            [str(binary_path), "--smartdoc-smoke"],
            0,
            stdout="".join(stdout_parts),
            stderr="".join(stderr_parts),
        ),
    )


def _run_successful_warning_smoke(
    binary_path: Path,
    *,
    cwd: Path,
    input_path: Path,
    expected_message: str,
) -> tuple[Path, Path, Path, subprocess.CompletedProcess[str]]:
    source_path = input_path.resolve()
    if not source_path.is_file():
        raise FileNotFoundError(f"packaged_gui_successful_warning_input_missing: {source_path}")
    if not expected_message.strip():
        raise ValueError("packaged_gui_successful_warning_message_empty")

    work_dir = cwd / "gui_successful_warning_smoke"
    output_dir = work_dir / "outputs"
    report_path = work_dir / "successful_warning_smoke.json"
    screenshot_path = work_dir / "successful_warning_info_area.png"
    work_dir.mkdir(parents=True, exist_ok=True)
    output_dir.mkdir(parents=True, exist_ok=True)
    report_path.unlink(missing_ok=True)
    screenshot_path.unlink(missing_ok=True)
    for stale_output in output_dir.iterdir():
        if stale_output.is_file():
            stale_output.unlink()

    env = _base_env(cwd)
    env["DOCWEN_GUI_TEST_CONVERSION_REPORT"] = str(report_path)
    env["DOCWEN_GUI_TEST_CONVERSION_INPUT"] = str(source_path)
    env["DOCWEN_GUI_TEST_CONVERSION_OUTPUT_DIR"] = str(output_dir)
    env["DOCWEN_GUI_TEST_CONVERSION_TARGET"] = "md"
    env["DOCWEN_GUI_TEST_CONVERSION_SURFACE"] = "action"
    env["DOCWEN_GUI_TEST_CONVERSION_ACTION"] = "gongwen"
    env["DOCWEN_GUI_TEST_CONVERSION_EXPECT_WARNING"] = expected_message
    env["DOCWEN_GUI_TEST_CONVERSION_SCREENSHOT"] = str(screenshot_path)
    env["DOCWEN_GUI_TEST_CONVERSION_TIMEOUT_MS"] = "120000"
    env["DOCWEN_GUI_TEST_AUTOCLOSE_MS"] = "135000"

    processes_before = _snapshot_relevant_processes()
    proc = _run_with_env(binary_path, cwd=cwd, env=env, timeout=150)
    output_path = _verify_successful_warning_smoke_report(
        report_path,
        screenshot_path=screenshot_path,
        input_path=source_path,
        expected_message=expected_message,
    )
    _wait_for_no_new_relevant_processes(processes_before)
    return output_path, screenshot_path, report_path, proc


def _probe_successful_warning_contract(
    cli_path: Path,
    *,
    cwd: Path,
    input_path: Path,
    action: str,
    expected_code: str,
    expected_message: str = "",
) -> tuple[str, Path]:
    """Resolve the canonical packaged warning before exercising its GUI projection.

    The GUI smoke owns presentation parity, not Gongwen business wording.  Resolve
    the warning from the packaged CLI using the same candidate, input, action,
    target, locale, and isolated configuration.  An optional explicit message is
    an additional golden oracle and may never override what the candidate emitted.
    """

    source_path = input_path.resolve()
    if not cli_path.is_file():
        raise FileNotFoundError(f"packaged_gui_successful_warning_cli_missing: {cli_path}")
    if not source_path.is_file():
        raise FileNotFoundError(f"packaged_gui_successful_warning_input_missing: {source_path}")
    if not action.strip() or not expected_code.strip():
        raise ValueError("packaged_gui_successful_warning_contract_empty")

    work_dir = cwd / "gui_successful_warning_preflight"
    work_dir.mkdir(parents=True, exist_ok=True)
    output_dir = work_dir / "outputs"
    output_dir.mkdir()
    proc = _run_with_env(
        cli_path,
        "convert",
        str(source_path),
        "--to",
        "md",
        "--optimization",
        action,
        "--output",
        str(output_dir),
        "--lang",
        _SUCCESSFUL_WARNING_LOCALE,
        "--json",
        "--quiet",
        cwd=cwd,
        env=_base_env(cwd),
        timeout=120,
    )
    if proc.returncode != 0:
        raise RuntimeError(
            "packaged_gui_successful_warning_cli_failed: "
            f"exit={proc.returncode}; stdout={proc.stdout!r}; stderr={proc.stderr!r}"
        )
    try:
        payload = json.loads(proc.stdout)
    except json.JSONDecodeError as exc:
        raise RuntimeError(f"packaged_gui_successful_warning_cli_json_invalid: {proc.stdout!r}") from exc
    if not isinstance(payload, dict) or payload.get("protocol_version") != 3 or payload.get("success") is not True:
        raise RuntimeError(f"packaged_gui_successful_warning_cli_unexpected: {payload}")
    warnings = payload.get("warnings")
    if not isinstance(warnings, list):
        raise RuntimeError(f"packaged_gui_successful_warning_cli_warnings_missing: {payload}")
    matching = [
        item
        for item in warnings
        if isinstance(item, dict) and item.get("level") == "warning" and item.get("code") == expected_code
    ]
    if len(matching) != 1:
        raise RuntimeError(
            "packaged_gui_successful_warning_cli_warning_count_unexpected: "
            f"expected_code={expected_code!r}; warnings={warnings!r}"
        )
    canonical_message = matching[0].get("message")
    if not isinstance(canonical_message, str) or not canonical_message.strip():
        raise RuntimeError(f"packaged_gui_successful_warning_cli_message_missing: {matching[0]!r}")
    if expected_message and canonical_message != expected_message:
        raise RuntimeError(
            "packaged_gui_successful_warning_golden_mismatch: "
            f"expected={expected_message!r}; actual={canonical_message!r}"
        )
    data = payload.get("data")
    output_raw = data.get("output") if isinstance(data, dict) else None
    if not isinstance(output_raw, str) or not output_raw:
        raise RuntimeError(f"packaged_gui_successful_warning_cli_output_missing: {payload}")
    output_path = Path(output_raw)
    output_path = (cwd / output_path).resolve() if not output_path.is_absolute() else output_path.resolve()
    if not output_path.is_relative_to(output_dir.resolve()):
        raise RuntimeError(f"packaged_gui_successful_warning_cli_output_escaped: {output_path}")
    _verify_markdown_document_node_path(
        output_path,
        expected_source_stem=source_path.stem,
        expected_source_format=source_path.suffix.removeprefix(".").capitalize(),
    )
    if not output_path.is_file() or output_path.stat().st_size <= 0:
        raise RuntimeError(f"packaged_gui_successful_warning_cli_output_empty: {output_path}")
    return canonical_message, output_path


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(
        description=(
            "Verify packaged DocWen GUI with internal env-var overrides. "
            "DOCWEN_CONFIG_DIR / DOCWEN_LOG_DIR / DOCWEN_LOG_TO_TEMP are internal hooks for CI, "
            "packaging verification, and dev debugging only - not a stable user-facing API."
        )
    )
    parser.add_argument("--binary-dir", required=True, help="Directory containing the packaged DocWen binary.")
    parser.add_argument(
        "--binary-name", default=_default_binary_name(), help="Binary filename inside the package directory."
    )
    parser.add_argument(
        "--notification-smoke",
        action="store_true",
        help=(
            "Also trigger the internal packaged GUI tray-notification path and "
            "verify its JSON report. This narrows, but does not fully prove, "
            "user-visible Windows notification-center delivery."
        ),
    )
    parser.add_argument(
        "--settings-smoke",
        action="store_true",
        help=(
            "Inspect the embedded PyInstaller archive and construct every Settings page "
            "inside the packaged GUI, rejecting missing modules and load placeholders."
        ),
    )
    parser.add_argument(
        "--ocr-smoke",
        action="store_true",
        help=(
            "Also trigger the internal packaged GUI image-to-Markdown OCR workflow "
            "and verify its JSON report/output sidecar. This proves the frozen GUI "
            "can drive the MainWindow action route through bundled OCR models."
        ),
    )
    parser.add_argument(
        "--ipc-smoke",
        action="store_true",
        help=(
            "Also run a packaged runtime/control smoke. This verifies the exact "
            "`gui.settings` info contract, requires no pre-existing GUI, cold-starts "
            "the GUI through `DocWenCLI gui open-settings`, verifies the runtime "
            "handshake and singleton reuse, then opens a file through `gui open` "
            "and verifies that the same primary consumes it."
        ),
    )
    parser.add_argument(
        "--office-smoke",
        action="store_true",
        help=(
            "Run three packaged GUI Office-backed routes (DOCX/XLSX/Markdown to PDF) "
            "through ConversionPanel/ActionArea and verify final PDF text. Requires a "
            "registered Word/Excel/WPS or installed LibreOffice backend."
        ),
    )
    parser.add_argument(
        "--presentation-smoke",
        action="store_true",
        help=(
            "Run a packaged GUI legacy PPT-to-Markdown route through ActionArea and "
            "verify the final Markdown name and semantics. The verifier first creates "
            "a deterministic legacy PPT fixture through a registered PowerPoint/WPS "
            "Presentation or installed LibreOffice backend."
        ),
    )
    parser.add_argument(
        "--smartdoc-smoke",
        action="store_true",
        help=(
            "Run packaged GUI DOC/RTF/ODT-to-DOCX routes through ConversionPanel "
            "and verify final DOCX names and semantic readback. The verifier first "
            "creates deterministic external-format fixtures through registered "
            "Word/WPS Writer or installed LibreOffice backends."
        ),
    )
    parser.add_argument(
        "--successful-warning-smoke",
        nargs="?",
        const=_GENERATED_SUCCESSFUL_WARNING_FIXTURE,
        type=Path,
        metavar="[INPUT]",
        help=(
            "Run a real Gongwen conversion through the packaged MainWindow action route, "
            "verify the successful warning row and persistent summary tone, save an "
            "InfoArea PNG, and retain the final Markdown. Omit INPUT to use the pinned "
            "deterministic fixture. This does not prove Windows notification-center delivery."
        ),
    )
    parser.add_argument(
        "--successful-warning-message",
        default="",
        help=(
            "Optional exact Gongwen warning golden for --successful-warning-smoke. "
            "The packaged CLI remains the canonical source passed to the GUI projection."
        ),
    )
    parser.add_argument(
        "--successful-warning-code",
        default=SUCCESSFUL_WARNING_CODE,
        help="Exact warning code required from the packaged CLI preflight.",
    )
    parser.add_argument(
        "--successful-warning-input-sha256",
        default="",
        help="Optional exact SHA-256 required for an explicit successful-warning input.",
    )
    parser.add_argument(
        "--evidence-dir",
        type=Path,
        help=(
            "After every requested verification succeeds, copy the complete isolated run directory "
            "to this exact new absolute directory and verify every file byte-for-byte and by SHA-256. "
            "The destination must not already exist or be inside the packaged binary directory."
        ),
    )
    args = parser.parse_args(argv)
    if args.ipc_smoke and (
        args.notification_smoke
        or args.settings_smoke
        or args.ocr_smoke
        or args.office_smoke
        or args.presentation_smoke
        or args.smartdoc_smoke
        or args.successful_warning_smoke
    ):
        parser.error("--ipc-smoke must run alone because it starts a long-lived primary instance without autoclose")
    if args.office_smoke and (
        args.notification_smoke
        or args.settings_smoke
        or args.ocr_smoke
        or args.presentation_smoke
        or args.smartdoc_smoke
        or args.successful_warning_smoke
    ):
        parser.error("--office-smoke must run alone because it launches three conversion instances")
    if args.presentation_smoke and (
        args.notification_smoke
        or args.settings_smoke
        or args.ocr_smoke
        or args.smartdoc_smoke
        or args.successful_warning_smoke
    ):
        parser.error("--presentation-smoke must run alone because it launches an external-bridge conversion")
    if args.smartdoc_smoke and (
        args.notification_smoke or args.settings_smoke or args.ocr_smoke or args.successful_warning_smoke
    ):
        parser.error("--smartdoc-smoke must run alone because it launches three external-bridge conversions")
    if args.successful_warning_smoke and (args.notification_smoke or args.settings_smoke or args.ocr_smoke):
        parser.error("--successful-warning-smoke must run alone because it launches a conversion instance")
    if args.settings_smoke and (args.notification_smoke or args.ocr_smoke):
        parser.error("--settings-smoke must run alone so its page-construction report is deterministic")

    binary_dir = Path(args.binary_dir).resolve()
    binary_path = binary_dir / args.binary_name
    if not binary_path.is_file():
        raise FileNotFoundError(f"packaged_gui_not_found: {binary_path}")
    evidence_dir = (
        _validate_evidence_destination(args.evidence_dir, binary_dir=binary_dir)
        if args.evidence_dir is not None
        else None
    )
    _verify_resource_layout(binary_dir)
    if args.settings_smoke:
        _verify_settings_page_archive(binary_path)

    # Ensure it's executable on Unix
    if os.name != "nt":
        binary_path.chmod(binary_path.stat().st_mode | 0o111)

    print(f"Running packaged GUI smoke test: {binary_path} ...")

    verification_dir = Path(tempfile.mkdtemp(prefix="docwen-packaged-gui-verify-"))
    verification_succeeded = False
    try:
        (verification_dir / "config_home").mkdir(parents=True, exist_ok=True)
        (verification_dir / "log_home").mkdir(parents=True, exist_ok=True)
        notification_report = verification_dir / "notification_smoke.json"
        settings_report = verification_dir / "settings_smoke.json"
        ocr_work_dir = verification_dir / "gui_ocr_smoke"
        ocr_report = verification_dir / "gui_ocr_smoke.json"
        ocr_output_dir = ocr_work_dir / "outputs"
        ocr_input = ocr_work_dir / "sample_gui_ocr.png"
        office_outputs: list[Path] = []
        presentation_output: Path | None = None
        presentation_fixture_backend = ""
        smartdoc_outputs: list[Path] = []
        smartdoc_fixture_backends: list[str] = []
        successful_warning_output: Path | None = None
        successful_warning_screenshot: Path | None = None
        successful_warning_report: Path | None = None
        successful_warning_cli_output: Path | None = None
        successful_warning_input_sha256 = ""
        if args.office_smoke:
            office_outputs, proc = _run_office_smoke(binary_path, cwd=verification_dir)
        elif args.presentation_smoke:
            presentation_output, presentation_fixture_backend, proc = _run_presentation_smoke(
                binary_path,
                cwd=verification_dir,
            )
        elif args.smartdoc_smoke:
            smartdoc_outputs, smartdoc_fixture_backends, proc = _run_smartdoc_smoke(
                binary_path,
                cwd=verification_dir,
            )
        elif args.successful_warning_smoke:
            if args.successful_warning_smoke == _GENERATED_SUCCESSFUL_WARNING_FIXTURE:
                successful_warning_input = verification_dir / "fixtures" / "successful-warning.docx"
                successful_warning_input_sha256 = write_successful_warning_fixture(successful_warning_input)
                warning_golden_message = args.successful_warning_message or SUCCESSFUL_WARNING_MESSAGE
            else:
                successful_warning_input = args.successful_warning_smoke.resolve()
                if not successful_warning_input.is_file():
                    raise FileNotFoundError(
                        f"packaged_gui_successful_warning_input_missing: {successful_warning_input}"
                    )
                successful_warning_input_sha256 = hashlib.sha256(successful_warning_input.read_bytes()).hexdigest()
                warning_golden_message = args.successful_warning_message
            expected_input_sha256 = args.successful_warning_input_sha256.strip().lower()
            if expected_input_sha256:
                if not re.fullmatch(r"[0-9a-f]{64}", expected_input_sha256):
                    raise ValueError("packaged_gui_successful_warning_input_sha256_invalid")
                if successful_warning_input_sha256 != expected_input_sha256:
                    raise RuntimeError(
                        "packaged_gui_successful_warning_input_sha256_mismatch: "
                        f"expected={expected_input_sha256}; actual={successful_warning_input_sha256}"
                    )
            canonical_warning_message, successful_warning_cli_output = _probe_successful_warning_contract(
                binary_dir / _default_cli_binary_name(),
                cwd=verification_dir,
                input_path=successful_warning_input,
                action=SUCCESSFUL_WARNING_ACTION,
                expected_code=args.successful_warning_code,
                expected_message=warning_golden_message,
            )
            (
                successful_warning_output,
                successful_warning_screenshot,
                successful_warning_report,
                proc,
            ) = _run_successful_warning_smoke(
                binary_path,
                cwd=verification_dir,
                input_path=successful_warning_input,
                expected_message=canonical_warning_message,
            )
            if successful_warning_output.read_bytes() != successful_warning_cli_output.read_bytes():
                raise RuntimeError("packaged_gui_successful_warning_cli_gui_output_mismatch")
        else:
            old_env = {
                key: os.environ.get(key)
                for key in (
                    "DOCWEN_GUI_TEST_NOTIFICATION_REPORT",
                    "DOCWEN_GUI_TEST_SETTINGS_REPORT",
                    "DOCWEN_GUI_TEST_OCR_REPORT",
                    "DOCWEN_GUI_TEST_OCR_INPUT",
                    "DOCWEN_GUI_TEST_OCR_OUTPUT_DIR",
                    "DOCWEN_GUI_TEST_OCR_EXPECTED_TEXT",
                    "DOCWEN_GUI_TEST_OCR_TIMEOUT_MS",
                    "DOCWEN_GUI_TEST_IPC_REPORT",
                    "DOCWEN_GUI_TEST_IPC_EXPECT_FILE",
                    "DOCWEN_GUI_TEST_IPC_TIMEOUT_MS",
                    "DOCWEN_GUI_TEST_AUTOCLOSE_MS",
                )
            }
            if args.notification_smoke:
                notification_report.unlink(missing_ok=True)
                os.environ["DOCWEN_GUI_TEST_NOTIFICATION_REPORT"] = str(notification_report)
            if args.settings_smoke:
                settings_report.unlink(missing_ok=True)
                os.environ["DOCWEN_GUI_TEST_SETTINGS_REPORT"] = str(settings_report)
            if args.ocr_smoke:
                ocr_report.unlink(missing_ok=True)
                ocr_work_dir.mkdir(parents=True, exist_ok=True)
                ocr_output_dir.mkdir(parents=True, exist_ok=True)
                _write_ocr_png(ocr_input)
                os.environ["DOCWEN_GUI_TEST_OCR_REPORT"] = str(ocr_report)
                os.environ["DOCWEN_GUI_TEST_OCR_INPUT"] = str(ocr_input)
                os.environ["DOCWEN_GUI_TEST_OCR_OUTPUT_DIR"] = str(ocr_output_dir)
                os.environ["DOCWEN_GUI_TEST_OCR_EXPECTED_TEXT"] = "HELLO DOCWEN OCR"
                os.environ["DOCWEN_GUI_TEST_OCR_TIMEOUT_MS"] = "60000"
                os.environ["DOCWEN_GUI_TEST_AUTOCLOSE_MS"] = "90000"
            try:
                proc = (
                    _run_ipc_smoke(binary_path, cwd=verification_dir, binary_dir=binary_dir)
                    if args.ipc_smoke
                    else _run(binary_path, cwd=verification_dir)
                )
            finally:
                for key, old_value in old_env.items():
                    if old_value is None:
                        os.environ.pop(key, None)
                    else:
                        os.environ[key] = old_value
        if proc.returncode != 0:
            raise RuntimeError(
                f"GUI smoke failed with exit code {proc.returncode}:\nSTDOUT:\n{proc.stdout}\nSTDERR:\n{proc.stderr}"
            )

        if args.notification_smoke:
            _verify_notification_smoke_report(notification_report)
        settings_result = None
        if args.settings_smoke:
            settings_result = _verify_settings_smoke_report(settings_report)
        ocr_output = None
        if args.ocr_smoke:
            ocr_output = _verify_ocr_smoke_report(ocr_report)
        ipc_report = None
        if args.ipc_smoke:
            ipc_report = verification_dir / "ipc_smoke.json"
            _verify_ipc_smoke_report(ipc_report)

        log_dir = verification_dir / "log_home" / "logs"
        log_files = list(log_dir.glob("*.log"))
        if not log_files:
            raise RuntimeError(f"packaged_gui_log_missing: {log_dir}")
        _verify_runtime_diagnostics(proc, log_files=log_files)

        retained_evidence = None
        if evidence_dir is not None:
            retained_evidence = _preserve_success_evidence(
                verification_dir,
                evidence_dir,
                binary_dir=binary_dir,
            )

        if settings_result is not None:
            print(f"packaged_gui_smoke_ok: {binary_path.name}; settings_tabs={len(settings_result['loadedTabs'])}")
        elif office_outputs:
            names = ", ".join(path.name for path in office_outputs)
            print(f"packaged_gui_smoke_ok: {binary_path.name}; office -> {names}")
        elif presentation_output is not None:
            print(
                f"packaged_gui_smoke_ok: {binary_path.name}; presentation -> "
                f"{presentation_output.name}; fixture_backend={presentation_fixture_backend}"
            )
        elif smartdoc_outputs:
            names = ", ".join(path.name for path in smartdoc_outputs)
            backends = ", ".join(smartdoc_fixture_backends)
            print(f"packaged_gui_smoke_ok: {binary_path.name}; smartdoc -> {names}; fixture_backends={backends}")
        elif successful_warning_output is not None:
            print(
                f"packaged_gui_smoke_ok: {binary_path.name}; warning -> "
                f"{successful_warning_output.name}; screenshot -> "
                f"{successful_warning_screenshot.name if successful_warning_screenshot else 'missing'}; "
                f"report -> {successful_warning_report.name if successful_warning_report else 'missing'}; "
                f"fixture_sha256={successful_warning_input_sha256}"
            )
        elif ocr_output is not None and ipc_report is not None:
            print(f"packaged_gui_smoke_ok: {binary_path.name}; ocr -> {ocr_output.name}; ipc -> {ipc_report.name}")
        elif ocr_output is not None:
            print(f"packaged_gui_smoke_ok: {binary_path.name}; ocr -> {ocr_output.name}")
        elif ipc_report is not None:
            print(f"packaged_gui_smoke_ok: {binary_path.name}; ipc -> {ipc_report.name}")
        else:
            print(f"packaged_gui_smoke_ok: {binary_path.name}")
        if retained_evidence is not None:
            print(f"packaged_gui_evidence_retained: {retained_evidence}")
        verification_succeeded = True
    except subprocess.TimeoutExpired as exc:
        raise RuntimeError(
            f"GUI smoke test timed out! The application did not auto-close within the expected time.\nSTDOUT:\n{exc.stdout}\nSTDERR:\n{exc.stderr}"
        ) from exc
    finally:
        if verification_succeeded:
            shutil.rmtree(verification_dir, ignore_errors=True)
        else:
            print(f"packaged_gui_failure_artifacts_retained: {verification_dir}", file=sys.stderr)

    return 0


def _run_ipc_smoke(
    binary_path: Path,
    *,
    cwd: Path,
    binary_dir: Path | None = None,
) -> subprocess.CompletedProcess[str]:
    cli_path = (binary_dir or cwd) / _default_cli_binary_name()
    if not cli_path.is_file():
        raise FileNotFoundError(f"packaged_gui_control_cli_missing: {cli_path}")
    report_path = cwd / "ipc_smoke.json"
    input_path = cwd / "ipc_smoke_input.md"
    input_path.write_text("# Control smoke\n\nOpen this file in the primary instance.\n", encoding="utf-8")
    report_path.unlink(missing_ok=True)

    env = _base_env(cwd)
    env.pop("DOCWEN_GUI_TEST_AUTOCLOSE_MS", None)
    env.pop("DOCWEN_GUI_DISABLE_CONTROL", None)
    env["DOCWEN_GUI_TEST_IPC_REPORT"] = str(report_path)
    env["DOCWEN_GUI_TEST_IPC_EXPECT_FILE"] = str(input_path)
    env["DOCWEN_GUI_TEST_IPC_TIMEOUT_MS"] = "30000"

    control_requests: list[subprocess.CompletedProcess[str]] = []
    info_request = _run_with_env(cli_path, "info", "--json", cwd=cwd, env=env, timeout=10)
    _verify_gui_settings_info_response(info_request)
    control_requests.append(info_request)

    # A pre-existing GUI belongs to the user, not this verifier. Fail before
    # entering the owned-process cleanup scope so it can never be terminated.
    stopped_request = _run_with_env(cli_path, "gui", "status", "--json", cwd=cwd, env=env, timeout=10)
    _verify_gui_stopped_response(stopped_request)
    control_requests.append(stopped_request)

    owned_pid: int | None = None
    try:
        # This is intentionally the first process that starts the GUI. It proves
        # the packaged CLI cold-start path instead of pre-starting DocWen.exe and
        # only exercising the already-running control path.
        cold_settings_request = _run_with_env(
            cli_path,
            "gui",
            "open-settings",
            "--section",
            "proofread",
            "--timeout",
            "30",
            "--json",
            cwd=cwd,
            env=env,
            timeout=40,
        )
        status_data = _wait_for_control_ready(
            cli_path=cli_path,
            cwd=cwd,
            env=env,
            timeout=20,
        )
        owned_pid = cast(int, status_data["pid"])
        _verify_open_settings_response(cold_settings_request, expected_reused=False)
        control_requests.append(cold_settings_request)
        _verify_gui_settings_handshake(status_data)
        _wait_for_ipc_report_ready(report_path=report_path, timeout=20)

        settings_request = _run_with_env(
            cli_path,
            "gui",
            "open-settings",
            "--section",
            "proofread",
            "--json",
            cwd=cwd,
            env=env,
            timeout=20,
        )
        _verify_open_settings_response(settings_request, expected_reused=True)
        control_requests.append(settings_request)

        control_request = _run_with_env(
            cli_path,
            "gui",
            "open",
            str(input_path),
            "--json",
            cwd=cwd,
            env=env,
            timeout=20,
        )
        if control_request.returncode != 0:
            raise RuntimeError(
                f"packaged_gui_control_open_failed: {control_request.returncode}\n"
                f"STDOUT:\n{control_request.stdout}\nSTDERR:\n{control_request.stderr}"
            )
        control_requests.append(control_request)
        _wait_for_control_stopped(cli_path=cli_path, cwd=cwd, env=env, timeout=40)
        _verify_ipc_smoke_report(report_path)
    except Exception:
        if owned_pid is not None:
            _terminate_test_gui(cli_path=cli_path, cwd=cwd, env=env, expected_pid=owned_pid)
        raise
    return subprocess.CompletedProcess(
        [str(binary_path), "--control-smoke-cold-start"],
        0,
        stdout="".join(request.stdout for request in control_requests),
        stderr="".join(request.stderr for request in control_requests),
    )


def _wait_for_control_ready(
    *,
    cli_path: Path,
    cwd: Path,
    env: dict[str, str],
    timeout: int,
) -> dict[str, Any]:
    deadline = time.monotonic() + timeout
    while time.monotonic() < deadline:
        status = _run_with_env(cli_path, "gui", "status", "--json", cwd=cwd, env=env, timeout=5)
        payload: dict[str, Any] = {}
        try:
            candidate = json.loads(status.stdout)
            payload = candidate if isinstance(candidate, dict) else {}
            raw_data = payload.get("data")
            data = raw_data if isinstance(raw_data, dict) else {}
        except (json.JSONDecodeError, TypeError):
            data = {}
        pid = data.get("pid")
        if (
            status.returncode == 0
            and payload.get("protocol_version") == 3
            and payload.get("success") is True
            and payload.get("command") == "gui status"
            and data.get("state") == "running"
            and data.get("running") is True
            and data.get("control_ready") is True
            and isinstance(pid, int)
            and not isinstance(pid, bool)
            and pid > 0
        ):
            return cast(dict[str, Any], data)
        time.sleep(0.1)
    raise RuntimeError("packaged_gui_control_ready_timeout")


def _verify_gui_settings_info_response(completed: subprocess.CompletedProcess[str]) -> dict[str, Any]:
    payload = _read_cli_json(completed, error_prefix="packaged_gui_info")
    data = payload.get("data")
    capabilities = data.get("capabilities") if isinstance(data, dict) else None
    settings = (
        [item for item in capabilities if isinstance(item, dict) and item.get("id") == "gui.settings"]
        if isinstance(capabilities, list)
        else []
    )
    expected = {
        "id": "gui.settings",
        "contract_version": 1,
        "state": "runtime_check_required",
        "available": True,
        "platforms": ["windows"],
        "current_platform_supported": True,
        "runtime_check_required": True,
        "details": {"cold_start": True, "sections": ["proofread"]},
    }
    if (
        payload.get("protocol_version") != 3
        or payload.get("success") is not True
        or payload.get("command") != "info"
        or len(settings) != 1
        or settings[0] != expected
    ):
        raise RuntimeError(f"packaged_gui_settings_info_contract_mismatch: {payload}")
    return cast(dict[str, Any], settings[0])


def _verify_gui_stopped_response(completed: subprocess.CompletedProcess[str]) -> None:
    payload = _read_cli_json(completed, error_prefix="packaged_gui_status")
    data = payload.get("data")
    if (
        payload.get("protocol_version") != 3
        or payload.get("success") is not True
        or payload.get("command") != "gui status"
        or not isinstance(data, dict)
        or data.get("state") != "stopped"
        or data.get("running") is not False
        or data.get("control_ready") is not False
    ):
        raise RuntimeError(f"packaged_gui_control_not_stopped_before_cold_start: {payload}")


def _read_cli_json(completed: subprocess.CompletedProcess[str], *, error_prefix: str) -> dict[str, Any]:
    if completed.returncode != 0:
        raise RuntimeError(
            f"{error_prefix}_failed: {completed.returncode}\nSTDOUT:\n{completed.stdout}\nSTDERR:\n{completed.stderr}"
        )
    raw_stdout = completed.stdout
    try:
        candidate = json.loads(raw_stdout)
    except json.JSONDecodeError as exc:
        raise RuntimeError(f"{error_prefix}_json_invalid: {raw_stdout!r}") from exc
    if not isinstance(candidate, dict):
        raise RuntimeError(f"{error_prefix}_json_not_object: {candidate!r}")
    return candidate


def _wait_for_control_stopped(
    *,
    cli_path: Path,
    cwd: Path,
    env: dict[str, str],
    timeout: int,
) -> None:
    deadline = time.monotonic() + timeout
    while time.monotonic() < deadline:
        status = _run_with_env(cli_path, "gui", "status", "--json", cwd=cwd, env=env, timeout=5)
        try:
            _verify_gui_stopped_response(status)
        except RuntimeError:
            time.sleep(0.1)
            continue
        return
    raise RuntimeError("packaged_gui_control_stop_timeout")


def _terminate_test_gui(
    *,
    cli_path: Path,
    cwd: Path,
    env: dict[str, str],
    expected_pid: int,
) -> None:
    with contextlib.suppress(Exception):
        status = _run_with_env(cli_path, "gui", "status", "--json", cwd=cwd, env=env, timeout=5)
        payload = _read_cli_json(status, error_prefix="packaged_gui_cleanup_status")
        data = payload.get("data")
        pid = data.get("pid") if isinstance(data, dict) else None
        if isinstance(pid, int) and not isinstance(pid, bool) and pid == expected_pid:
            os.kill(pid, signal.SIGTERM)


def _verify_gui_settings_handshake(status_data: dict[str, Any]) -> None:
    actions = status_data.get("supported_actions")
    sections = status_data.get("settings_sections")
    if not isinstance(actions, list) or "open_settings" not in actions:
        raise RuntimeError(f"packaged_gui_settings_capability_missing: {status_data}")
    if not isinstance(sections, list) or "proofread" not in sections:
        raise RuntimeError(f"packaged_gui_proofread_settings_section_missing: {status_data}")


def _verify_open_settings_response(
    completed: subprocess.CompletedProcess[str],
    *,
    expected_reused: bool,
) -> dict[str, Any]:
    if completed.returncode != 0:
        raise RuntimeError(
            f"packaged_gui_open_settings_failed: {completed.returncode}\n"
            f"STDOUT:\n{completed.stdout}\nSTDERR:\n{completed.stderr}"
        )
    raw_stdout = completed.stdout
    try:
        candidate = json.loads(raw_stdout)
    except json.JSONDecodeError as exc:
        raise RuntimeError(f"packaged_gui_open_settings_json_invalid: {raw_stdout!r}") from exc
    payload = candidate if isinstance(candidate, dict) else {}
    data = payload.get("data")
    if (
        payload.get("success") is not True
        or payload.get("command") != "gui open-settings"
        or not isinstance(data, dict)
        or data.get("accepted") is not True
        or data.get("running") is not True
        or data.get("action") != "open_settings"
        or data.get("section") != "proofread"
        or data.get("reused") is not expected_reused
    ):
        raise RuntimeError(
            f"packaged_gui_open_settings_contract_mismatch: expected_reused={expected_reused}; payload={payload}"
        )
    return cast(dict[str, Any], data)


def _wait_for_ipc_report_ready(
    *,
    report_path: Path,
    timeout: int,
) -> None:
    deadline = time.monotonic() + timeout
    while time.monotonic() < deadline:
        try:
            report = json.loads(report_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            time.sleep(0.1)
            continue
        status = report.get("status")
        if status == "waiting":
            return
        if status in {"failed", "timed_out", "completed"}:
            raise RuntimeError(f"packaged_gui_ipc_hook_unexpected_early_status: {report}")
        time.sleep(0.1)
    raise RuntimeError(f"packaged_gui_ipc_hook_ready_timeout: {report_path}")


def _verify_notification_smoke_report(report_path: Path) -> None:
    if not report_path.is_file():
        raise RuntimeError(f"packaged_gui_notification_report_missing: {report_path}")
    report = json.loads(report_path.read_text(encoding="utf-8"))
    if report.get("error"):
        raise RuntimeError(f"packaged_gui_notification_probe_error: {report['error']}")
    if not report.get("isSystemTrayAvailable"):
        raise RuntimeError("packaged_gui_notification_tray_unavailable")
    if not report.get("supportsMessages"):
        raise RuntimeError("packaged_gui_notification_messages_unsupported")
    if report.get("defaultTrayIconPresent") is not False:
        raise RuntimeError("packaged_gui_notification_default_tray_icon_present")
    if report.get("probeCreatedTrayIcon") is not True:
        raise RuntimeError("packaged_gui_notification_probe_tray_not_created")
    if not report.get("hasTrayIcon"):
        raise RuntimeError("packaged_gui_notification_tray_icon_missing")
    if not report.get("showMessageCalled"):
        raise RuntimeError("packaged_gui_notification_show_message_not_called")


def _verify_ocr_smoke_report(report_path: Path) -> Path:
    if not report_path.is_file():
        raise RuntimeError(f"packaged_gui_ocr_report_missing: {report_path}")
    report = json.loads(report_path.read_text(encoding="utf-8"))
    if report.get("error"):
        raise RuntimeError(f"packaged_gui_ocr_probe_error: {report['error']}")
    if report.get("status") != "completed":
        raise RuntimeError(f"packaged_gui_ocr_status_unexpected: {report}")
    if not report.get("primaryOutputExists"):
        raise RuntimeError(f"packaged_gui_ocr_primary_output_missing: {report}")

    output_raw = report.get("outputPath")
    if not isinstance(output_raw, str):
        raise RuntimeError(f"packaged_gui_ocr_output_path_missing: {report}")
    output_path = Path(output_raw)
    if not output_path.is_file():
        raise RuntimeError(f"packaged_gui_ocr_output_file_missing: {output_path}")
    expected_text = str(report.get("expectedText") or "HELLO DOCWEN OCR")
    primary_text = output_path.read_text(encoding="utf-8", errors="replace")
    sidecar_raw = report.get("sidecarPath")
    sidecar_path = Path(sidecar_raw) if isinstance(sidecar_raw, str) and sidecar_raw else None
    sidecar_text = (
        sidecar_path.read_text(encoding="utf-8", errors="replace")
        if sidecar_path is not None and sidecar_path.is_file()
        else ""
    )
    expected_words = " ".join(expected_text.split())

    def contains_expected_words(text: str) -> bool:
        without_quote_prefixes = re.sub(r"(?m)^\s*>+\s?", "", text)
        return expected_words in " ".join(without_quote_prefixes.split())

    primary_contains_expected = bool(report.get("primaryContainsExpectedText")) or contains_expected_words(primary_text)
    sidecar_contains_expected = bool(report.get("sidecarContainsExpectedText")) or contains_expected_words(sidecar_text)
    if not (primary_contains_expected or sidecar_contains_expected):
        raise RuntimeError(f"packaged_gui_ocr_output_missing_expected_text: {report}")
    if sidecar_contains_expected:
        if not report.get("sidecarOutputExists") and not (sidecar_path and sidecar_path.is_file()):
            raise RuntimeError(f"packaged_gui_ocr_sidecar_output_missing: {report}")
        if not report.get("primaryReferencesSidecar") and sidecar_path and sidecar_path.name not in primary_text:
            raise RuntimeError(f"packaged_gui_ocr_primary_missing_sidecar_reference: {report}")
    return output_path


def _verify_successful_warning_smoke_report(
    report_path: Path,
    *,
    screenshot_path: Path,
    input_path: Path,
    expected_message: str,
) -> Path:
    output_path = _read_conversion_smoke_output(report_path)
    report = json.loads(report_path.read_text(encoding="utf-8"))
    if output_path.stat().st_size <= 0 or int(report.get("outputBytes") or 0) != output_path.stat().st_size:
        raise RuntimeError(f"packaged_gui_successful_warning_output_empty_or_untracked: {report}")
    _verify_markdown_document_node_path(
        output_path,
        expected_source_stem=input_path.stem,
        expected_source_format=input_path.suffix.removeprefix(".").capitalize(),
    )
    if report.get("actionName") != "gongwen" or report.get("surface") != "action":
        raise RuntimeError(f"packaged_gui_successful_warning_route_unexpected: {report}")
    if report.get("expectedWarningMessage") != expected_message:
        raise RuntimeError(f"packaged_gui_successful_warning_message_unexpected: {report}")
    history_rows = report.get("historyRows")
    if not isinstance(history_rows, list) or not any(
        isinstance(row, dict) and row.get("message") == expected_message and row.get("messageType") == "warning"
        for row in history_rows
    ):
        raise RuntimeError(f"packaged_gui_successful_warning_history_missing: {report}")
    if (
        report.get("warningRowTone") != "warning"
        or report.get("warningRowTooltip") != expected_message
        or report.get("warningRowVisible") is not True
    ):
        raise RuntimeError(f"packaged_gui_successful_warning_row_unexpected: {report}")
    warning_row_screenshot_path = screenshot_path.with_name(
        f"{screenshot_path.stem}_warning_row{screenshot_path.suffix}"
    )
    if (
        report.get("warningRowScreenshotPath") != str(warning_row_screenshot_path)
        or report.get("warningRowScreenshotSaved") is not True
        or not warning_row_screenshot_path.is_file()
        or not warning_row_screenshot_path.read_bytes().startswith(b"\x89PNG\r\n\x1a\n")
        or int(report.get("warningRowScreenshotBytes") or 0) != warning_row_screenshot_path.stat().st_size
        or int(report.get("warningRowScreenshotWidth") or 0) <= 0
        or int(report.get("warningRowScreenshotHeight") or 0) <= 0
    ):
        raise RuntimeError(f"packaged_gui_successful_warning_row_screenshot_invalid: {report}")
    summary = report.get("taskSummary")
    if not isinstance(summary, dict) or (
        summary.get("state") != "success"
        or summary.get("tone") != "warning"
        or summary.get("completedCount") != 1
        or summary.get("failedCount") != 0
        or summary.get("navigatePath") != str(output_path)
    ):
        raise RuntimeError(f"packaged_gui_successful_warning_summary_unexpected: {report}")
    if report.get("statusSource") != "task" or report.get("statusTone") != "warning":
        raise RuntimeError(f"packaged_gui_successful_warning_status_unexpected: {report}")
    if report.get("infoAreaVisible") is not True:
        raise RuntimeError(f"packaged_gui_successful_warning_info_area_hidden: {report}")
    if report.get("screenshotPath") != str(screenshot_path) or report.get("screenshotSaved") is not True:
        raise RuntimeError(f"packaged_gui_successful_warning_screenshot_unreported: {report}")
    if not screenshot_path.is_file() or not screenshot_path.read_bytes().startswith(b"\x89PNG\r\n\x1a\n"):
        raise RuntimeError(f"packaged_gui_successful_warning_screenshot_invalid: {screenshot_path}")
    if (
        int(report.get("screenshotBytes") or 0) != screenshot_path.stat().st_size
        or int(report.get("screenshotWidth") or 0) <= 0
        or int(report.get("screenshotHeight") or 0) <= 0
    ):
        raise RuntimeError(f"packaged_gui_successful_warning_screenshot_metadata_invalid: {report}")
    return output_path


def _verify_ipc_smoke_report(report_path: Path) -> None:
    if not report_path.is_file():
        raise RuntimeError(f"packaged_gui_ipc_report_missing: {report_path}")
    report = json.loads(report_path.read_text(encoding="utf-8"))
    if report.get("error"):
        raise RuntimeError(f"packaged_gui_ipc_probe_error: {report['error']}")
    if report.get("status") != "completed":
        raise RuntimeError(f"packaged_gui_ipc_status_unexpected: {report}")
    if report.get("success") is not True:
        raise RuntimeError(f"packaged_gui_ipc_smoke_failed: {report}")
    if report.get("expectedReceived") is not True:
        raise RuntimeError(f"packaged_gui_ipc_expected_file_not_received: {report}")
    if report.get("expectedInFiles") is not True:
        raise RuntimeError(f"packaged_gui_ipc_expected_file_not_added: {report}")
    activation_count = report.get("activationCount")
    if not isinstance(activation_count, int) or isinstance(activation_count, bool) or activation_count < 1:
        raise RuntimeError(f"packaged_gui_ipc_activation_missing: {report}")

    expected_raw = report.get("expectedFile")
    if not isinstance(expected_raw, str) or not expected_raw:
        raise RuntimeError(f"packaged_gui_ipc_expected_file_missing: {report}")
    expected_path = Path(expected_raw)
    if not expected_path.is_file():
        raise RuntimeError(f"packaged_gui_ipc_expected_file_absent: {expected_path}")

    received_files = report.get("receivedFiles")
    files = report.get("files")
    if not isinstance(received_files, list) or expected_raw not in received_files:
        raise RuntimeError(f"packaged_gui_ipc_received_files_unexpected: {report}")
    if not isinstance(files, list) or expected_raw not in files:
        raise RuntimeError(f"packaged_gui_ipc_files_unexpected: {report}")


def _read_conversion_smoke_output(report_path: Path) -> Path:
    if not report_path.is_file():
        raise RuntimeError(f"packaged_gui_conversion_report_missing: {report_path}")
    report = json.loads(report_path.read_text(encoding="utf-8"))
    if report.get("error"):
        raise RuntimeError(f"packaged_gui_conversion_probe_error: {report['error']}")
    if report.get("status") != "completed" or not report.get("success"):
        raise RuntimeError(f"packaged_gui_conversion_status_unexpected: {report}")
    if not report.get("outputExists"):
        raise RuntimeError(f"packaged_gui_conversion_output_missing: {report}")

    output_raw = report.get("outputPath")
    if not isinstance(output_raw, str) or not output_raw:
        raise RuntimeError(f"packaged_gui_conversion_output_path_missing: {report}")
    output_path = Path(output_raw)
    if not output_path.is_file():
        raise RuntimeError(f"packaged_gui_conversion_output_missing: {output_path}")
    return output_path


def _verify_office_conversion_metrics(
    report_path: Path,
    *,
    case_name: str,
    input_path: Path,
    output_path: Path,
) -> None:
    report = json.loads(report_path.read_text(encoding="utf-8"))
    metrics = report.get("conversionMetrics") if isinstance(report, dict) else None
    if not isinstance(metrics, dict):
        raise RuntimeError(f"packaged_gui_office_metrics_invalid: {case_name}: {metrics!r}")

    input_size = input_path.stat().st_size
    output_size = output_path.stat().st_size
    duration_ms = metrics.get("durationMs")
    input_bytes = metrics.get("inputBytes")
    output_bytes = metrics.get("outputBytes")
    reported_output_bytes = report.get("outputBytes")
    engine = metrics.get("engine")
    backend = metrics.get("backend")
    duration_valid = (
        isinstance(duration_ms, (int, float))
        and not isinstance(duration_ms, bool)
        and math.isfinite(float(duration_ms))
        and duration_ms >= 0
    )
    input_valid = (
        isinstance(input_bytes, int)
        and not isinstance(input_bytes, bool)
        and input_bytes > 0
        and input_bytes == input_size
        and report.get("inputPath") == str(input_path)
    )
    output_valid = (
        isinstance(output_bytes, int)
        and not isinstance(output_bytes, bool)
        and output_bytes > 0
        and output_bytes == output_size
        and isinstance(reported_output_bytes, int)
        and not isinstance(reported_output_bytes, bool)
        and reported_output_bytes == output_size
        and report.get("outputPath") == str(output_path)
    )
    engine_valid = engine == "office_bridge"
    backend_valid = isinstance(backend, str) and bool(backend.strip())
    if not all((duration_valid, input_valid, output_valid, engine_valid, backend_valid)):
        raise RuntimeError(f"packaged_gui_office_metrics_invalid: {case_name}: {metrics!r}")


def _snapshot_relevant_processes() -> dict[int, str]:
    if os.name == "nt":
        return _snapshot_windows_relevant_processes()

    processes: dict[int, str] = {}
    result = subprocess.run(
        ["ps", "-eo", "pid=,comm="],
        capture_output=True,
        text=True,
        errors="replace",
        check=False,
    )
    if result.returncode != 0:
        raise RuntimeError(f"packaged_gui_process_snapshot_failed: {result.stderr}")
    for line in result.stdout.splitlines():
        fields = line.split(maxsplit=1)
        if len(fields) != 2:
            continue
        name = Path(fields[1]).name.lower()
        if name in _RELEVANT_PROCESS_NAMES:
            try:
                processes[int(fields[0])] = name
            except ValueError:
                continue
    return processes


def _snapshot_windows_relevant_processes() -> dict[int, str]:
    processes: dict[int, str] = {}
    tasklist_result = subprocess.run(
        ["tasklist", "/FO", "CSV", "/NH"],
        capture_output=True,
        text=True,
        errors="replace",
        check=False,
    )
    if tasklist_result.returncode == 0:
        for row in csv.reader(tasklist_result.stdout.splitlines()):
            if len(row) < 2:
                continue
            name = _normalize_relevant_process_name(row[0])
            if name is None:
                continue
            try:
                processes[int(row[1].replace(",", ""))] = name
            except ValueError:
                continue
        return processes

    # Managed desktop sessions can deny ``tasklist`` while the current user
    # can still enumerate process identities through Get-Process. Keep the
    # residue gate fail-closed without requiring an elevated system query.
    powershell_result = subprocess.run(
        [
            "powershell.exe",
            "-NoLogo",
            "-NoProfile",
            "-NonInteractive",
            "-Command",
            "$ErrorActionPreference='Stop'; Get-Process | ForEach-Object { '{0}`t{1}' -f $_.Id, $_.ProcessName }",
        ],
        capture_output=True,
        text=True,
        errors="replace",
        check=False,
    )
    if powershell_result.returncode != 0:
        raise RuntimeError(
            "packaged_gui_process_snapshot_failed: "
            f"tasklist={tasklist_result.stderr.strip()!r}; "
            f"get_process={powershell_result.stderr.strip()!r}"
        )
    for line in powershell_result.stdout.splitlines():
        fields = line.split("\t", maxsplit=1)
        if len(fields) != 2:
            continue
        name = _normalize_relevant_process_name(fields[1])
        if name is None:
            continue
        try:
            processes[int(fields[0])] = name
        except ValueError:
            continue
    return processes


def _normalize_relevant_process_name(raw_name: str) -> str | None:
    name = Path(raw_name.strip()).name.lower()
    if name in _RELEVANT_PROCESS_NAMES:
        return name
    executable_name = f"{name}.exe"
    return executable_name if executable_name in _RELEVANT_PROCESS_NAMES else None


def _verify_no_new_relevant_processes(before: dict[int, str], after: dict[int, str]) -> None:
    new_processes = {pid: name for pid, name in after.items() if pid not in before}
    if new_processes:
        raise RuntimeError(f"packaged_gui_process_residue: {new_processes}")


def _wait_for_no_new_relevant_processes(before: dict[int, str], *, settle_seconds: float = 6.0) -> None:
    time.sleep(settle_seconds)
    _verify_no_new_relevant_processes(before, _snapshot_relevant_processes())


def _verify_runtime_diagnostics(
    proc: subprocess.CompletedProcess[str],
    *,
    log_files: list[Path],
) -> None:
    sources = [proc.stdout, proc.stderr]
    for log_file in log_files:
        sources.append(log_file.read_text(encoding="utf-8", errors="replace"))
    combined = "\n".join(sources)
    failures = [marker for marker in _RUNTIME_FAILURE_MARKERS if marker in combined]
    if failures:
        raise RuntimeError(f"packaged_gui_runtime_diagnostics_failed: markers={failures}")


def _verify_markdown_conversion_smoke_report(
    report_path: Path,
    *,
    expected_source_stem: str,
    expected_source_format: str,
    expected_tokens: tuple[str, ...],
) -> Path:
    output_path = _read_conversion_smoke_output(report_path)
    _verify_markdown_document_node_path(
        output_path,
        expected_source_stem=expected_source_stem,
        expected_source_format=expected_source_format,
    )
    try:
        output_text = output_path.read_text(encoding="utf-8")
    except (OSError, UnicodeError) as exc:
        raise RuntimeError(f"packaged_gui_conversion_markdown_invalid: {output_path}: {exc}") from exc
    missing_tokens = [token for token in expected_tokens if token not in output_text]
    if missing_tokens:
        raise RuntimeError(
            f"packaged_gui_conversion_markdown_semantics_missing: missing={missing_tokens}, output={output_path}"
        )
    return output_path


def _verify_markdown_document_node_path(
    output_path: Path,
    *,
    expected_source_stem: str,
    expected_source_format: str,
) -> None:
    root_name = output_path.stem
    expected_root = re.fullmatch(
        rf"{re.escape(expected_source_stem)}_(\d{{8}}_\d{{6}})_from{re.escape(expected_source_format)}",
        root_name,
    )
    if expected_root is None or output_path.parent.name != root_name:
        raise RuntimeError(
            "packaged_gui_conversion_markdown_name_unexpected: "
            f"expected={expected_source_stem}_YYYYMMDD_HHMMSS_from{expected_source_format}/"
            f"{expected_source_stem}_YYYYMMDD_HHMMSS_from{expected_source_format}.md, "
            f"actual={output_path.parent.name}/{output_path.name}"
        )
    try:
        time.strptime(expected_root.group(1), "%Y%m%d_%H%M%S")
    except ValueError as exc:
        raise RuntimeError(f"packaged_gui_conversion_markdown_timestamp_invalid: {output_path}") from exc


def _verify_docx_conversion_smoke_report(
    report_path: Path,
    *,
    expected_name: str,
    expected_tokens: tuple[str, ...],
) -> Path:
    output_path = _read_conversion_smoke_output(report_path)
    if output_path.name != expected_name:
        raise RuntimeError(
            f"packaged_gui_conversion_docx_name_unexpected: expected={expected_name}, actual={output_path.name}"
        )

    from docx import Document

    try:
        document = Document(str(output_path))
    except (OSError, ValueError) as exc:
        raise RuntimeError(f"packaged_gui_conversion_docx_invalid: {output_path}: {exc}") from exc
    text_parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            text_parts.extend(cell.text for cell in row.cells)
    output_text = " ".join(text_parts)
    missing_tokens = [token for token in expected_tokens if token not in output_text]
    if missing_tokens:
        raise RuntimeError(
            f"packaged_gui_conversion_docx_semantics_missing: missing={missing_tokens}, output={output_path}"
        )
    return output_path


def _verify_conversion_smoke_report(report_path: Path, *, expected_tokens: tuple[str, ...]) -> Path:
    output_path = _read_conversion_smoke_output(report_path)
    if not output_path.read_bytes().startswith(b"%PDF"):
        raise RuntimeError(f"packaged_gui_conversion_pdf_invalid: {output_path}")

    import fitz

    with fitz.open(output_path) as document:
        page_count = document.page_count
        text_parts: list[str] = []
        for page in document:
            page_text = page.get_text("text")
            if not isinstance(page_text, str):
                raise RuntimeError(f"packaged_gui_conversion_pdf_text_invalid: {output_path}")
            text_parts.append(" ".join(page_text.split()))
    output_text = " ".join(text_parts)
    missing_tokens = [token for token in expected_tokens if token not in output_text]
    if page_count < 1 or missing_tokens:
        raise RuntimeError(
            f"packaged_gui_conversion_pdf_semantics_missing: pages={page_count}, "
            f"missing={missing_tokens}, output={output_path}"
        )
    return output_path


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
